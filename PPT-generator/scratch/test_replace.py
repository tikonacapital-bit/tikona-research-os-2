from pptx import Presentation

def replace_text_in_pptx(prs_path, out_path, replacements):
    prs = Presentation(prs_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    full_text = paragraph.text
                    changed = False
                    for k, v in replacements.items():
                        var = f"{{{{{k}}}}}"
                        if var in full_text:
                            full_text = full_text.replace(var, str(v) if v is not None else "")
                            changed = True
                    
                    if changed:
                        p_font = None
                        if paragraph.runs:
                            p_font = paragraph.runs[0].font
                            p_font_name = p_font.name
                            p_font_size = p_font.size
                            p_font_bold = p_font.bold
                            p_font_color = None
                            if getattr(p_font.color, 'type', None) == 1:
                                try: p_font_color = p_font.color.rgb
                                except: pass
                        
                        paragraph.clear()
                        new_run = paragraph.add_run()
                        new_run.text = full_text
                        if p_font is not None:
                            if p_font_name: new_run.font.name = p_font_name
                            if p_font_size: new_run.font.size = p_font_size
                            if p_font_bold is not None: new_run.font.bold = p_font_bold
                            if p_font_color: new_run.font.color.rgb = p_font_color

            # Check inside tables
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            full_text = paragraph.text
                            changed = False
                            for k, v in replacements.items():
                                var = f"{{{{{k}}}}}"
                                if var in full_text:
                                    full_text = full_text.replace(var, str(v) if v is not None else "")
                                    changed = True
                            
                            if changed:
                                paragraph.text = full_text
    prs.save(out_path)

if __name__ == "__main__":
    replacements = {
        "company_name": "RELIANCE INDUSTRIES",
        "cmp": "3000",
        "target": "3500",
        "investment_thesis": "Reliance is a great company and it's growing rapidly."
    }
    replace_text_in_pptx(r"C:\tikona-research-os-2\master_template.pptx", r"C:\tikona-research-os-2\PPT-generator\scratch\test_output.pptx", replacements)
    print("Done")
