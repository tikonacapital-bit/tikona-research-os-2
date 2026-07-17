"""
Excel-to-image injector (cross-platform).

Reads Excel sheets with openpyxl, renders them as styled table images
using matplotlib, and returns PNG bytes keyed by sheet name.
Works on Linux/Docker — no COM/win32com dependency.
"""
import io
import logging
from pathlib import Path
from typing import Any

import openpyxl
from openpyxl.utils import get_column_letter

logger = logging.getLogger(__name__)

# ── Brand colours (matches financial_model_v5 + chart_generators) ─────────
_NAVY = "#1F4690"
_ORANGE = "#FFA500"
_BLUE = "#3A5BA0"
_CREAM = "#FFE5B4"
_GREY = "#F2F2F2"
_WHITE = "#FFFFFF"
_DGREEN = "#006400"
_DRED = "#C00000"
_LIGHT_BG = "#f8f9fc"

# Map placeholder token → sheet name(s) to try
PLACEHOLDER_SHEET_MAP: dict[str, list[str]] = {
    "{{earnings_forecast_table}}": ["Earnings_Forecast"],
    "{{financials_table}}": ["Financials_Table"],
    "{{valuations_table}}": ["Valuations_Table"],
    "{{key_risks_table}}": ["Key_Risks"],
    "{{peer_comparision}}": ["Peer_Compare", "Peer_Comparison"],
    "{{financial_model_from_excel_operational_sheet}}": ["Operational_Data"],
    "{{governance_table}}": ["Governance"],
    "{{timeline}}": ["Timeline"],
    "{{financial_model_from_excel}}": ["Op_Charts", "Fin_Summary"],
}


def _mpl():
    """Lazy-load matplotlib."""
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import numpy as np
        return plt, np
    except ImportError:
        logger.warning("matplotlib not installed — Excel table rendering disabled")
        return None, None


def _openpyxl_color_to_hex(fill) -> str | None:
    """Extract hex colour from an openpyxl PatternFill, or None."""
    if fill is None or fill.patternType is None:
        return None
    fg = fill.fgColor
    if fg is None:
        return None
    if fg.type == "rgb" and fg.rgb and fg.rgb != "00000000":
        rgb = fg.rgb
        # openpyxl sometimes returns AARRGGBB
        if len(rgb) == 8:
            rgb = rgb[2:]
        return f"#{rgb}"
    if fg.type == "indexed":
        # Common indexed colours
        _IDX = {
            0: "#000000", 1: "#FFFFFF", 2: "#FF0000", 3: "#00FF00",
            4: "#0000FF", 5: "#FFFF00", 9: "#800000",
        }
        return _IDX.get(fg.indexed)
    if fg.type == "theme":
        # Can't easily resolve theme — return None
        return None
    return None


def _font_color_hex(font) -> str:
    """Extract font colour, default black."""
    if font is None or font.color is None:
        return "#000000"
    c = font.color
    if c.type == "rgb" and c.rgb and c.rgb != "00000000":
        rgb = c.rgb
        if len(rgb) == 8:
            rgb = rgb[2:]
        return f"#{rgb}"
    return "#000000"


def _format_cell_value(cell) -> str:
    """Format a cell's value for display."""
    v = cell.value
    if v is None:
        return "—"
    nf = cell.number_format or "General"

    if isinstance(v, (int, float)):
        try:
            if "%" in nf:
                return f"{v * 100:.1f}%" if abs(v) < 10 else f"{v:.1f}%"
            if "x" in nf.lower():
                return f"{v:.1f}x"
            if "#,##0.00" in nf:
                return f"{v:,.2f}"
            if "#,##0" in nf or "," in nf:
                if abs(v) >= 1:
                    return f"{v:,.0f}"
                return f"{v:.2f}"
            if isinstance(v, float):
                if v == int(v) and abs(v) < 1e12:
                    return f"{int(v):,}"
                return f"{v:,.1f}"
            return f"{v:,}"
        except (ValueError, TypeError):
            return str(v)
    return str(v)


def _read_sheet_data(ws, max_rows: int = 60, max_cols: int = 15) -> tuple[list[list[str]], list[list[dict]], int, int]:
    """Read sheet into a grid of display strings + style info.
    
    Returns: (data_grid, style_grid, n_rows, n_cols)
    """
    # Find actual used range
    used_rows = min(ws.max_row or 1, max_rows)
    used_cols = min(ws.max_column or 1, max_cols)

    # Skip entirely empty leading rows
    first_data_row = 1
    for r in range(1, used_rows + 1):
        if any(ws.cell(r, c).value is not None for c in range(1, used_cols + 1)):
            first_data_row = r
            break

    data: list[list[str]] = []
    styles: list[list[dict]] = []

    for r in range(first_data_row, used_rows + 1):
        row_data: list[str] = []
        row_styles: list[dict] = []
        for c in range(1, used_cols + 1):
            cell = ws.cell(r, c)
            row_data.append(_format_cell_value(cell))

            bg = _openpyxl_color_to_hex(cell.fill) or _WHITE
            fg = _font_color_hex(cell.font)
            bold = bool(cell.font and cell.font.bold)
            align = "center"
            if cell.alignment and cell.alignment.horizontal:
                align = cell.alignment.horizontal

            row_styles.append({
                "bg": bg, "fg": fg, "bold": bold, "align": align,
            })
        data.append(row_data)
        styles.append(row_styles)

    return data, styles, len(data), used_cols


def render_sheet_as_image(ws, title: str = "", max_rows: int = 55) -> bytes | None:
    """Render an openpyxl worksheet as a styled PNG table image.
    
    Returns PNG bytes, or None if matplotlib unavailable or sheet empty.
    """
    plt, np = _mpl()
    if plt is None:
        return None

    data, styles, n_rows, n_cols = _read_sheet_data(ws, max_rows=max_rows)
    if n_rows == 0 or n_cols == 0:
        return None

    # Determine figure size based on content
    col_width = 1.6
    row_height = 0.32
    fig_w = max(8, n_cols * col_width + 1)
    fig_h = max(3, n_rows * row_height + 1.2)
    fig_h = min(fig_h, 24)  # cap height

    fig, ax = plt.subplots(figsize=(fig_w, fig_h))
    ax.set_xlim(0, n_cols)
    ax.set_ylim(0, n_rows)
    ax.invert_yaxis()
    ax.axis("off")
    fig.patch.set_facecolor(_LIGHT_BG)

    # Draw cells
    for ri in range(n_rows):
        for ci in range(n_cols):
            style = styles[ri][ci]
            text = data[ri][ci]
            bg = style["bg"]
            fg = style["fg"]
            bold = style["bold"]
            ha = style["align"] if style["align"] in ("left", "center", "right") else "center"

            # Draw cell background
            rect = plt.Rectangle((ci, ri), 1, 1,
                                  facecolor=bg, edgecolor="#D5DCE8",
                                  linewidth=0.5)
            ax.add_patch(rect)

            # Draw text
            fontsize = 6.5 if n_cols > 8 else 7.5
            weight = "bold" if bold else "normal"
            ax.text(ci + 0.5, ri + 0.5, text,
                    ha="center", va="center",
                    fontsize=fontsize, color=fg,
                    fontweight=weight,
                    fontfamily="sans-serif",
                    clip_on=True)

    # Title
    if title:
        fig.suptitle(title, fontsize=9, color=_NAVY, fontweight="bold", y=0.98)

    fig.tight_layout(pad=0.3)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=180, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    plt.close(fig)
    return buf.getvalue()


def _find_sheet(wb, names: list[str]):
    """Find the first matching sheet from a list of candidate names."""
    wb_sheets_lower = {s.lower(): s for s in wb.sheetnames}
    for name in names:
        # Exact match
        if name in wb.sheetnames:
            return wb[name]
        # Case-insensitive
        real = wb_sheets_lower.get(name.lower())
        if real:
            return wb[real]
    return None


def render_all_excel_sheets(excel_path: str) -> dict[str, bytes]:
    """Open an Excel file and render each target sheet as a PNG image.
    
    Returns a dict mapping placeholder token → PNG bytes.
    """
    path = Path(excel_path)
    if not path.exists():
        logger.warning("Excel file not found: %s", excel_path)
        return {}

    try:
        wb = openpyxl.load_workbook(str(path), data_only=True, read_only=False)
    except Exception as e:
        logger.error("Failed to open Excel workbook: %s", e)
        return {}

    logger.info("Excel opened. Sheets: %s", wb.sheetnames)
    results: dict[str, bytes] = {}

    for token, sheet_names in PLACEHOLDER_SHEET_MAP.items():
        ws = _find_sheet(wb, sheet_names)
        if ws is None:
            logger.warning("No sheet found for %s (tried %s)", token, sheet_names)
            continue

        logger.info("Rendering sheet '%s' for placeholder %s", ws.title, token)
        try:
            img = render_sheet_as_image(ws, title=ws.title)
            if img:
                results[token] = img
                logger.info("Rendered %s → %d bytes", token, len(img))
            else:
                logger.warning("Sheet '%s' rendered empty for %s", ws.title, token)
        except Exception as e:
            logger.error("Failed to render sheet '%s': %s", ws.title, e)

    wb.close()
    return results


def inject_excel_visuals_into_pptx(pptx_path: str, images: dict[str, bytes]) -> int:
    """Replace Excel placeholder shapes in a PPTX with rendered images.
    
    Uses python-pptx (cross-platform, no COM needed).
    
    Args:
        pptx_path: Path to the PPTX file to modify in-place.
        images: Dict mapping placeholder token → PNG bytes.
    
    Returns: Number of placeholders replaced.
    """
    from pptx import Presentation

    if not images:
        logger.info("No images to inject — skipping")
        return 0

    prs = Presentation(pptx_path)
    replaced = 0
    # Collect (slide, shape, token) tuples first to avoid mutation during iteration
    to_replace: list[tuple] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or not shape.text_frame:
                continue
            try:
                text = shape.text_frame.text.strip()
            except Exception:
                continue
            if text in images:
                to_replace.append((slide, shape, text))

    for slide, shape, token in to_replace:
        try:
            img_bytes = images[token]
            left, top, width, height = shape.left, shape.top, shape.width, shape.height

            # Remove original placeholder shape
            sp_elem = shape._element
            sp_elem.getparent().remove(sp_elem)

            # Insert the rendered image at the same position/size
            slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width, height)
            replaced += 1
            logger.info("Replaced placeholder %s with image (%dx%d)",
                        token, width, height)
        except Exception as e:
            logger.error("Failed to replace placeholder %s: %s", token, e)

    if replaced > 0:
        prs.save(pptx_path)
        logger.info("Saved PPTX with %d Excel visual injections", replaced)

    return replaced


# ── Legacy COM-based injector (Windows only, kept as fallback) ────────────

def inject_excel_visuals_into_ppt(excel_path: str, ppt_path: str) -> int:
    """Try COM-based injection (Windows), fall back to openpyxl rendering.
    
    Returns the number of visuals injected.
    """
    # Try COM first (only works on Windows with Excel + PowerPoint installed)
    try:
        import win32com.client  # noqa: F401
        count = _com_inject(excel_path, ppt_path)
        if count > 0:
            return count
        logger.info("COM injection returned 0 — falling back to openpyxl renderer")
    except ImportError:
        logger.info("win32com not available — using openpyxl renderer (cross-platform)")

    # Cross-platform fallback: render sheets as images and inject
    images = render_all_excel_sheets(excel_path)
    if images:
        return inject_excel_visuals_into_pptx(ppt_path, images)
    return 0


def extract_json_from_excel(excel_path: str) -> dict | None:
    """Extract JSON data from Excel file for financial model processing.
    
    Looks for a hidden '_json' sheet or parses data from visible sheets
    to construct the financial model JSON structure.
    """
    path = Path(excel_path)
    if not path.exists():
        logger.warning("Excel file not found: %s", excel_path)
        return None

    try:
        wb = openpyxl.load_workbook(str(path), data_only=True, read_only=False)
    except Exception as e:
        logger.error("Failed to open Excel workbook: %s", e)
        return None

    # Try to find hidden JSON sheet first
    json_sheet = None
    for sheet_name in wb.sheetnames:
        if sheet_name.lower() in ('_json', 'json', 'model_json', 'financial_model_json'):
            json_sheet = wb[sheet_name]
            break
    
    if json_sheet:
        logger.info("Found JSON sheet: %s", json_sheet.title)
        try:
            # Try to parse JSON from the first cell
            json_text = None
            for row in json_sheet.iter_rows(values_only=True):
                if row and row[0]:
                    json_text = str(row[0])
                    break
            
            if json_text:
                import json
                return json.loads(json_text)
        except Exception as e:
            logger.warning("Failed to parse JSON from sheet: %s", e)
    
    # Fallback: construct basic JSON structure from visible sheets
    logger.info("No JSON sheet found, constructing basic structure from visible sheets")
    result = {
        "assumptions": {},
        "projections": {},
        "historical_ratios": {},
        "metrics": {},
        "valuation": {}
    }
    
    # Try to extract basic data from common sheets
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        if sheet_name.lower() in ('assumptions', 'inputs'):
            result["assumptions"] = _extract_sheet_data(ws)
        elif sheet_name.lower() in ('projections', 'forecast'):
            result["projections"] = _extract_sheet_data(ws)
        elif sheet_name.lower() in ('historical', 'ratios'):
            result["historical_ratios"] = _extract_sheet_data(ws)
    
    wb.close()
    return result if any(result.values()) else None


def _extract_sheet_data(ws) -> dict:
    """Extract basic key-value data from a worksheet."""
    data = {}
    for row in ws.iter_rows(max_row=20, values_only=True):
        if row and row[0] and row[1]:
            key = str(row[0]).strip()
            value = row[1]
            if value is not None:
                try:
                    data[key] = float(value) if isinstance(value, (int, float)) else str(value)
                except (ValueError, TypeError):
                    data[key] = str(value)
    return data


def _com_inject(excel_path: str, ppt_path: str) -> int:
    """COM-based injection (Windows only). Returns count of visuals injected."""
    import os
    import time
    import win32com.client

    excel_path = os.path.abspath(excel_path)
    ppt_path = os.path.abspath(ppt_path)

    if not os.path.exists(excel_path) or not os.path.exists(ppt_path):
        return 0

    excel = None
    ppt = None
    wb = None
    prs = None

    try:
        excel = win32com.client.DispatchEx("Excel.Application")
        ppt = win32com.client.DispatchEx("PowerPoint.Application")
        excel.Visible = False
        excel.DisplayAlerts = False

        wb = excel.Workbooks.Open(excel_path)
        prs = ppt.Presentations.Open(ppt_path, WithWindow=False)

        shapes_replaced = 0
        placeholder_map = {
            "{{earnings_forecast_table}}": "Earnings_Forecast",
            "{{financials_table}}": "Financials_Table",
            "{{valuations_table}}": "Valuations_Table",
            "{{key_risks_table}}": "Key_Risks",
            "{{peer_comparision}}": "Peer_Compare",
            "{{financial_model_from_excel_operational_sheet}}": "Operational_Data",
            "{{governance_table}}": "Governance",
            "{{timeline}}": "Timeline",
            "{{financial_model_from_excel}}": "Op_Charts",
        }

        for slide_index in range(1, prs.Slides.Count + 1):
            slide = prs.Slides(slide_index)
            targets = []

            for si in range(1, slide.Shapes.Count + 1):
                try:
                    shape = slide.Shapes(si)
                    if getattr(shape, "HasTextFrame", False) and shape.TextFrame.HasText:
                        text = shape.TextFrame.Text.strip()
                        if text in placeholder_map:
                            sheet_name = placeholder_map[text]
                            shape.Name = f"excel_table:{sheet_name}"
                            targets.append(shape.Name)
                except Exception:
                    pass

            for shape_name in targets:
                try:
                    shape = slide.Shapes(shape_name)
                    top, left, width, height = shape.Top, shape.Left, shape.Width, shape.Height
                    sheet_target = shape_name.replace("excel_table:", "")

                    try:
                        wb.Sheets(sheet_target).UsedRange.Copy()
                    except Exception:
                        continue

                    time.sleep(0.3)
                    pasted = slide.Shapes.PasteSpecial(DataType=2)
                    if getattr(pasted, "Count", 0) > 0:
                        pasted = pasted(1)
                    pasted.Top = top
                    pasted.Left = left
                    pasted.Width = width
                    pasted.Height = height
                    shape.Delete()
                    shapes_replaced += 1
                except Exception as e:
                    logger.error("COM inject failed for %s: %s", shape_name, e)

        if shapes_replaced > 0:
            prs.Save()
        return shapes_replaced

    except Exception as e:
        logger.error("COM injection error: %s", e)
        return 0
    finally:
        if wb:
            try: wb.Close(False)
            except: pass
        if prs:
            try: prs.Close()
            except: pass
        if excel:
            try: excel.Quit()
            except: pass
        if ppt:
            try: ppt.Quit()
            except: pass
