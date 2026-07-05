"""Orchestrates PPTX generation: pulls report+session data from Supabase, builds
the four reportgen input files, runs the pipeline, and uploads results."""
from __future__ import annotations

import json
import logging
import os
import re
import shutil
import tempfile
import time
import urllib.error
import urllib.request
import io
from contextlib import contextmanager
from datetime import date, datetime, timezone
from pathlib import Path
from typing import Any

from supabase_client import get_service_client
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from openpyxl.utils import column_index_from_string, get_column_letter
from openpyxl.utils.cell import coordinate_from_string

from chart_generators import generate_chart
import excel_injector

import logging
logger = logging.getLogger(__name__)

PPTX_BUCKET = "research-reports-pptx"
MODEL_BUCKET = os.environ.get("SUPABASE_FINANCIAL_MODEL_BUCKET", "research-reports-html")

PPTX_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
PDF_CONTENT_TYPE = "application/pdf"

_NUM_RE = re.compile(r"[\d,]+\.?\d*")

# Fields whose values should render with inline **bold** segments preserved.
# Sanitiser on the TS side keeps the `**...**` markers for these keys; the
# Python writer below splits them into bold/non-bold runs.
_INLINE_BOLD_KEYS: set[str] = {"investment_thesis_s1"}

# SAARTHI per-dimension score placeholders on slide 16 are SINGLE-brace tokens
# in the template ({s_s}, {a1_s}, ...) — the renderer's standard {{key}} pass
# misses them. We do a second pass that also tries the single-brace form
# for these specific keys.
_SAARTHI_SCORE_PLACEHOLDERS: set[str] = {"s_s", "a1_s", "a2_s", "r_s", "t_s", "h_s", "i_s"}

# Bullet-shaped placeholders on slides 5 + 19. If the TS sanitiser missed (or
# the override JSON was authored elsewhere), these may arrive as a single
# paragraph with bullet markers embedded mid-string. We split defensively so
# the renderer's per-line paragraph emitter receives one bullet per line.
_BULLET_LIST_KEYS: set[str] = {
    "industry_structure",
    "key_industry_tailwinds",
    "key_industry_risks",
    "entry_strategy",
    "review_strategy",
    "exit_strategy",
    "entry_strategy_1",
    "review_strategy_2",
    "exit_strategy_3",
    "key_industry",
    "key_industry_risk",
}


def _ensure_bullet_lines(value: str) -> str:
    """Normalise a bullets value to newline-separated `• <line>` lines.

    Accepts: (a) already-newline-separated; (b) one-paragraph with embedded
    bullet markers; (c) plain prose with sentence-level periods.
    """
    if not isinstance(value, str) or not value.strip():
        return value
    text = value.strip()
    if "\n" in text:
        lines = [ln.strip() for ln in text.split("\n") if ln.strip()]
    else:
        # Split on any embedded bullet marker.
        parts = re.split(r"(?=[•●▪◦])\s*", text)
        parts = [p.strip() for p in parts if p.strip()]
        if len(parts) <= 1:
            # No bullet markers — split on sentence boundary as a last resort.
            parts = [s.strip() for s in re.split(r"(?<=[.!?])\s+", text) if s.strip()]
        lines = parts
    out_lines: list[str] = []
    for ln in lines:
        cleaned = re.sub(r"^[\-\*•●▪◦·#>\s]+", "", ln).strip()
        if cleaned:
            out_lines.append(f"• {cleaned}")
    return "\n".join(out_lines) if out_lines else value

# Pure-numeric value? Used by fmt_number to decide whether to format or pass through.
_PURE_NUM_RE = re.compile(r"^-?\d[\d,]*\.?\d*$")


def fmt_number(value, *, force_decimal: bool = False) -> str:
    """Format a number with commas. Drop decimals when |value| > 150,
    keep one decimal otherwise. Pass through non-numeric strings unchanged.

    Examples:
        1809.3   -> '1,809'
        129.7    -> '129.7'
        12460    -> '12,460'
        0.45     -> '0.5'
        'BUY'    -> 'BUY'
        None,''  -> ''
    """
    if value is None:
        return ""
    if isinstance(value, str):
        s = value.strip()
        if not s:
            return ""
        # Allow leading currency / sign characters when checking for pure number
        candidate = s.replace(",", "").replace("₹", "").replace("$", "").strip()
        if not _PURE_NUM_RE.match(candidate.lstrip("-")) and not _PURE_NUM_RE.match(candidate):
            # Not a pure number string — pass through (e.g. "BUY", "₹1,200 Cr").
            return s
        try:
            num = float(candidate)
        except ValueError:
            return s
    else:
        try:
            num = float(value)
        except (TypeError, ValueError):
            return str(value)
    if abs(num) > 150 and not force_decimal:
        return f"{num:,.0f}"
    return f"{num:,.1f}"
_PLACEHOLDER_RE = re.compile(r"not included in the generated report|content pending", re.I)
_HOUSE_PLANNER_PATCHED = False
_ORPHAN_NUMBER_RE = re.compile(
    r"(?<!FY)(?<!fy)\b\d[\d,]*(?:\.\d+)?\s*(?:%|x\b|cr\b|crore\b|bn\b|billion\b|lakh\b|bps\b)"
    r"|(?:₹|\$|€|£|INR\s|USD\s|EUR\s|GBP\s)\s*\d[\d,]*(?:\.\d+)?",
    re.IGNORECASE,
)

# Placeholder tokens that the Excel injector needs to find intact.
# fill_master_template() must NOT replace text in shapes containing these.
_EXCEL_INJECTION_TOKENS: set[str] = {
    "{{financial_model_from_excel}}",
    "{{financial_model_from_excel_operational_sheet}}",
    "{{financial_charts}}",
    "{{operational_charts}}",
    "{{financial_summary_image}}",
    "{{earnings_forecast_table}}",
    "{{financials_table}}",
    "{{valuations_table}}",
    "{{key_risks_table}}",
    "{{peer_comparision}}",
    "{{governance_table}}",
    "{{timeline}}",
    "{{company_timeline}}",
    "{{competitive_chart_1}}",
    "{{competitive_chart_2}}",
    "{{peer_comparison_chart_1}}",
    "{{peer_comparison_chart_2}}",
    "{{pie_chart_1}}",
    "{{pie_chart_2}}",
    "{{percentage_revenue_pie_chart}}",
    "{{percentage_EBIT_pie_chart}}",
    "{{catalyst_timeline_chart}}",
    "{{probability_weight_table}}",
}


# ─────────── helpers ──────────────────────────────────────────────────────────


def _parse_number(val: Any) -> float | None:
    """Mirror the frontend parseNumber: extract first numeric token, strip commas."""
    if val is None:
        return None
    s = str(val)
    m = _NUM_RE.search(s)
    if not m:
        return None
    try:
        return float(m.group(0).replace(",", ""))
    except ValueError:
        return None


def _normalize_rating(raw: str) -> str:
    up = (raw or "").upper()
    if "SELL" in up:
        return "SELL"
    if "HOLD" in up:
        return "HOLD"
    if "ACCUMULATE" in up:
        return "ACCUMULATE"
    return "BUY"


def _section_value(sections: list[dict], key: str) -> str:
    for s in sections:
        if s.get("section_key") == key:
            value = (s.get("content") or "").strip()
            return "" if _PLACEHOLDER_RE.search(value) else value
    return ""


def _prefer(*vals: Any) -> str:
    for v in vals:
        if v is not None:
            value = str(v).strip()
            if value and not _PLACEHOLDER_RE.search(value):
                return value
    return ""


def _clean_prose(text: str, *, max_len: int = 500) -> str:
    text = (text or "").replace("â¹", "₹").replace("â‚¹", "₹")
    # Strip markdown bold/italic markers
    cleaned = re.sub(r"\*\*|__", "", text)
    cleaned = re.sub(r"\s+", " ", cleaned).strip()
    if len(cleaned) <= max_len:
        return cleaned
    # Truncate at word boundary so we don't cut mid-word
    truncated = cleaned[:max_len]
    last_space = truncated.rfind(" ")
    if last_space > int(max_len * 0.75):
        return truncated[:last_space].strip()
    return truncated.strip()


def _section_by_any(sections: list[dict], keywords: list[str]) -> str:
    for s in sections:
        haystack = f"{s.get('section_key') or ''} {s.get('section_title') or ''}".casefold()
        if any(k in haystack for k in keywords):
            value = (s.get("content") or "").strip()
            if value and not _PLACEHOLDER_RE.search(value):
                return value
    return ""


def _split_entry_review_exit(text: str) -> tuple[str, str, str]:
    """Split a combined Entry/Review/Exit strategy section into three parts.

    Looks for heading markers like:
    - **Entry Strategy:**  /  **Review Strategy:**  /  **Exit Strategy:**
    - ## Entry Strategy  /  ## Review Strategy  /  ## Exit Strategy
    - Entry Strategy:  /  Review Strategy:  /  Exit Strategy:

    Returns (entry_text, review_text, exit_text).  Any part not found returns "".
    """
    if not text or not text.strip():
        return ("", "", "")

    # Build a regex that captures the three sub-sections.  The pattern matches
    # any of the common heading formats the Stage-2 prompt produces.
    _heading = r"(?:#{1,3}\s*|[-*]\s*)?(?:\*\*\s*)?"   # optional #/*/- prefix
    _trail   = r"(?:\s*\*\*)?(?:\s*:?\s*)"              # optional **/ : suffix

    entry_pat  = re.compile(_heading + r"entry\s+strategy" + _trail, re.I)
    review_pat = re.compile(_heading + r"review\s+strategy" + _trail, re.I)
    exit_pat   = re.compile(_heading + r"exit\s+strategy" + _trail, re.I)

    # Find the start positions of each heading
    entry_m  = entry_pat.search(text)
    review_m = review_pat.search(text)
    exit_m   = exit_pat.search(text)

    # Collect (start_of_content, pattern_name) tuples, sorted by position
    markers: list[tuple[int, str]] = []
    if entry_m:
        markers.append((entry_m.end(), "entry"))
    if review_m:
        markers.append((review_m.end(), "review"))
    if exit_m:
        markers.append((exit_m.end(), "exit"))

    if not markers:
        # No sub-headings found — cannot split
        return ("", "", "")

    markers.sort(key=lambda x: x[0])

    parts: dict[str, str] = {}
    for i, (start, name) in enumerate(markers):
        end = markers[i + 1][0] if i + 1 < len(markers) else len(text)
        # Walk backwards from end to skip the next heading's prefix
        if i + 1 < len(markers):
            next_match = [m for m in (entry_m, review_m, exit_m) if m and m.end() == markers[i + 1][0]]
            if next_match:
                end = next_match[0].start()
        parts[name] = text[start:end].strip()

    return (parts.get("entry", ""), parts.get("review", ""), parts.get("exit", ""))

def _truncate_words(text: str, max_words: int, *, max_len: int = 240) -> str:
    cleaned = _clean_prose(text, max_len=max(max_len, len(text or "")))
    words = cleaned.split()
    if len(words) <= max_words:
        return cleaned
    return " ".join(words[:max_words]).rstrip(",;:.") + "."


def _sentences(text: str, *, limit: int = 4, max_len: int = 260) -> list[str]:
    cleaned = _clean_prose(text, max_len=max(800, len(text or "")))
    parts = re.split(r'(?<=[.!?])\s+', cleaned)
    out: list[str] = []
    for part in parts:
        item = part.strip()
        if len(item) < 20:
            continue
        out.append(_clean_prose(item, max_len=max_len))
        if len(out) >= limit:
            break
    return out


def _format_two_paragraphs(text_content: str) -> str:
    if not text_content:
        return "Financial forecast data shows stable projections.\n\nAnalysts remain constructive on the medium-term outlook."
    sents = _sentences(text_content, limit=4, max_len=220)
    if len(sents) == 0:
        return "Financial forecast data shows stable projections.\n\nAnalysts remain constructive on the medium-term outlook."
    elif len(sents) == 1:
        return sents[0]
    elif len(sents) == 2:
        return sents[0] + "\n\n" + sents[1]
    elif len(sents) == 3:
        return sents[0] + " " + sents[1] + "\n\n" + sents[2]
    else:
        return sents[0] + " " + sents[1] + "\n\n" + sents[2] + " " + sents[3]


def _metric_chips(metadata: dict, fin_model: dict, company: dict) -> list[str]:
    operational = fin_model.get("operational") or {}
    chips: list[str] = []

    def add(value: str) -> None:
        value = value.strip()
        if value and value not in chips:
            chips.append(value)

    mcap = _fmt_mcap(metadata.get("market_cap", ""))
    if mcap:
        add(mcap)

    cmp_val = metadata.get("cmp", "")
    if cmp_val:
        add(f"₹{cmp_val} CMP")

    target = metadata.get("target_price", "")
    if target:
        add(f"₹{target} TP")

    upside = str(metadata.get("upside_pct", "") or "")
    if upside:
        add(upside if upside.endswith("%") else f"{upside}% Upside")

    years = operational.get("years") or []
    utils = operational.get("capacity_utilisation_pct") or []
    if years and utils and utils[-1] is not None:
        try:
            add(f"{float(utils[-1]):.0f}% Util.")
        except (ValueError, TypeError):
            pass

    employees = operational.get("employees")
    if employees:
        try:
            add(f"{int(float(employees)):,} Staff")
        except (ValueError, TypeError):
            pass

    plants_india = operational.get("plants_india") or []
    plants_over = operational.get("plants_overseas") or []
    if plants_india or plants_over:
        try:
            india = int(plants_india[-1]) if (plants_india and plants_india[-1] is not None) else 0
            over = int(plants_over[-1]) if (plants_over and plants_over[-1] is not None) else 0
            total = india + over
            if total > 0:
                add(f"{total} Plants")
        except (ValueError, TypeError):
            pass

    countries = operational.get("countries_of_operation") or []
    if countries and countries[-1] is not None:
        try:
            add(f"{int(countries[-1])} Countries")
        except (ValueError, TypeError):
            pass

    thesis = fin_model.get("thesis") or {}
    score = thesis.get("saarthi_total")
    if score is not None:
        add(f"{float(score):.0f} SAARTHI")

    shares = metadata.get("shares_cr") or fin_model.get("metrics", {}).get("shares_cr")
    if shares:
        try:
            add(f"{float(str(shares)):.1f} Cr Shrs")
        except (TypeError, ValueError):
            pass

    sector = company.get("sector") or ""
    if sector:
        words = str(sector).split()
        add(" ".join(words[:2]))

    return chips[:6]


def _bullets_from_text(text: str, *, limit: int = 5, max_len: int = 180) -> list[str]:
    if not text:
        return ["Analyst narrative supports continued monitoring of the core investment case."]

    # Try splitting on **Bold header:** sections first — common in research reports
    bold_headers = re.findall(r'\*\*([^*]+?):\*\*', text)
    bold_parts   = re.split(r'\*\*[^*]+?:\*\*', text)
    if len(bold_headers) >= 2 and len(bold_parts) >= 2:
        bullets: list[str] = []
        for header, content in zip(bold_headers, bold_parts[1:]):
            # First sentence of content after the header
            first_sent = re.split(r'(?<=[.!?])\s+', content.strip())[0] if content.strip() else ""
            combined = f"{header}: {first_sent.strip()}" if first_sent.strip() else header
            item = _clean_prose(combined.lstrip("-*•: "), max_len=max_len)
            if len(item) >= 12:
                bullets.append(item)
            if len(bullets) >= limit:
                break
        if len(bullets) >= 2:
            return bullets

    # Fall back: split on newlines and sentence endings
    parts = re.split(r"(?:\n+|(?<=[.!?])\s+)", text)
    bullets = []
    for part in parts:
        item = _clean_prose(part.lstrip("-*•** ").strip(), max_len=max_len)
        if len(item) >= 12:
            bullets.append(item)
        if len(bullets) >= limit:
            break
    return bullets or ["Analyst narrative supports continued monitoring of the core investment case."]


def _all_sections_text(sections: list[dict]) -> str:
    parts: list[str] = []
    for s in sections:
        title = s.get("section_title") or s.get("section_key") or ""
        content = (s.get("content") or "").strip()
        if content and not _PLACEHOLDER_RE.search(content):
            parts.append(f"{title}\n{content}")
    return "\n\n".join(parts)


def _extract_labeled_number(text: str, patterns: list[str]) -> str:
    for pattern in patterns:
        match = re.search(pattern, text, re.I | re.S)
        if match:
            value = _parse_number(match.group(1))
            if value is not None and value > 0:
                return str(value)
    return ""


def _extract_rating(text: str) -> str:
    patterns = [
        r"\b(BUY|SELL|HOLD|ACCUMULATE)\b\s+(?:rating|recommendation)\b",
        r"\b(?:rating|recommendation)\s*(?:of|:|is|=)?\s*\b(BUY|SELL|HOLD|ACCUMULATE)\b",
        r"\bsupporting\s+a\s+\b(BUY|SELL|HOLD|ACCUMULATE)\b\s+rating\b",
    ]
    for pattern in patterns:
        match = re.search(pattern, text, re.I)
        if match:
            return match.group(1)
    return ""


def _strip_qmark(url: str) -> str:
    return url.rstrip("?") if url else url


# ─────────── input builders ───────────────────────────────────────────────────


def _build_company(report: dict, session: dict, sections: list[dict]) -> dict:
    description = _prefer(
        report.get("company_description"),
        session.get("company_description"),
        _section_by_any(sections, ["company", "business", "overview"]),
        _all_sections_text(sections),
    )
    return {
        "name": report.get("company_name") or "Unknown",
        "ticker": (report.get("nse_symbol") or "UNKNOWN").upper(),
        "exchange": "NSE",
        "sector": session.get("sector") or "General",
        "industry": session.get("industry") or session.get("sector") or "General",
        "country": "India",
        "description": _clean_prose(description, max_len=900),
        "peer_list": [],
    }


def _build_metadata(report: dict, sections: list[dict]) -> dict:
    narrative = _all_sections_text(sections)
    rating_raw = _prefer(report.get("cs_rating"), _section_value(sections, "rating"), _extract_rating(narrative))
    if not rating_raw:
        rating_raw = "BUY"
    rating = _normalize_rating(rating_raw)

    cmp_raw = _prefer(
        report.get("cs_current_market_price"),
        report.get("current_market_price"),
        _section_value(sections, "current_market_price"),
        _extract_labeled_number(
            narrative,
            [
                r"(?:current\s+(?:market\s+)?price|CMP)\s*(?:of|at|:|=)?[^\d]{0,20}([\d,]+(?:\.\d+)?)",
                r"(?:from|vs\.?)\s+current\s+(?:market\s+)?price\s*(?:of|at|:|=)?[^\d]{0,20}([\d,]+(?:\.\d+)?)",
            ],
        ),
    )
    tp_raw = _prefer(
        report.get("cs_target_price"),
        report.get("target_price"),
        _section_value(sections, "target_price"),
        _extract_labeled_number(
            narrative,
            [
                r"probability-weighted\s+target\s+price.*?=\s*[^\d]{0,20}([\d,]+(?:\.\d+)?)",
                r"probability-weighted\s+target\s+price\s*(?:of|at|:|=)?[^\d]{0,20}([\d,]+(?:\.\d+)?)",
                r"(?:BUY|ACCUMULATE|HOLD|SELL)\s+recommendation\s+with\s+[^\d]{0,20}([\d,]+(?:\.\d+)?)\s+target\s+price",
                r"base\s+case.*?implied\s+target\s+price\s*(?:of|at|:|=)?[^\d]{0,20}([\d,]+(?:\.\d+)?)",
                r"implied\s+target\s+price\s*(?:of|at|:|=)?[^\d]{0,20}([\d,]+(?:\.\d+)?)",
                r"target\s+price\s*(?:of|at|:|=)?[^\d]{0,20}([\d,]+(?:\.\d+)?)",
                r"[^\d]{0,20}([\d,]+(?:\.\d+)?)\s+target\s+price",
            ],
        ),
    )
    mcap_raw = _prefer(report.get("cs_market_cap"), _section_value(sections, "market_cap"))

    cmp = _parse_number(cmp_raw)
    target = _parse_number(tp_raw)
    if cmp is None or cmp <= 0:
        raise ValueError("CMP could not be extracted from report — fill cs_current_market_price column")
    if target is None or target <= 0:
        raise ValueError("target_price could not be extracted from report — fill cs_target_price column")

    mcap = _parse_number(mcap_raw)

    # Recompute upside_pct from cmp+target so it's always consistent with the validator's check.
    # (Stored upside_pct can be stale; validator allows ≤0.3% tolerance so recomputing avoids the error.)
    upside = round(((target - cmp) / cmp) * 100, 1)

    meta: dict[str, Any] = {
        "rating": rating,
        "currency": "INR",
        "cmp": str(cmp),
        "target_price": str(target),
        "upside_pct": str(upside),
        "analyst": _prefer(report.get("user_email"), "Tikona Research"),
        "report_date": date.today().isoformat(),
        "report_type": "Initiation",
    }
    meta["market_cap"] = str(mcap if mcap is not None and mcap > 0 else 0)
    return meta


def _series(name: str, unit: str, periods: list[str], values: list) -> dict | None:
    """Build a FinancialSeries dict, dropping None-only values."""
    if not periods or not values or len(periods) != len(values):
        return None
    return {
        "name": name,
        "unit": unit,
        "periods": periods,
        "values": [None if v is None else str(v) for v in values],
    }


def _list_of_dicts(value: Any) -> list[dict]:
    if isinstance(value, list):
        return [item for item in value if isinstance(item, dict)]
    return []


def _copy_if_present(target: dict, source: dict, key: str) -> None:
    value = source.get(key)
    if value:
        target[key] = value


def _placeholder_series() -> list[dict]:
    return [
        {
            "name": "Revenue",
            "unit": "INR Cr",
            "periods": ["FY24A", "FY25A", "FY26E"],
            "values": ["0", "0", "0"],
        }
    ]


def _build_financial_model(ticker: str, model_json: dict | None, warnings: list[str]) -> dict:
    """Map the v5 financial model JSON onto the FinancialModelSnapshot schema."""
    base: dict[str, Any] = {
        "model_name": f"{ticker} Base Model",
        "model_version": "v5.1",
        "currency": "INR",
        "fiscal_year_end": "March",
    }

    if not model_json:
        warnings.append("financial-model JSON sidecar missing; using minimal placeholder model")
        base["metrics"] = {"placeholder": "0"}
        base["series"] = _placeholder_series()
        return base

    asmp = model_json.get("assumptions") or {}
    proj = model_json.get("projections") or {}
    hist = model_json.get("historical_ratios") or {}

    proj_years: list[str] = [str(y) for y in (asmp.get("projection_years") or proj.get("years") or [])]
    hist_years: list[str] = [str(y) for y in (hist.get("years") or [])]

    # All series MUST share the same period list — combine hist + proj into one unified timeline.
    seen: set[str] = set()
    all_years: list[str] = []
    for y in hist_years + proj_years:
        if y not in seen:
            seen.add(y)
            all_years.append(y)

    series: list[dict] = []
    metrics: dict[str, str] = {}

    # ── Absolute P&L series (Revenue, EBITDA, PAT in ₹ Cr) ──────────────────
    # financial_model_v5._embed_pl_series() stores historical + projected
    # absolute values so charts show real numbers rather than placeholders.
    hist_pl = model_json.get("historical_pl") or {}
    proj_pl = model_json.get("projected_pl") or {}
    h_yrs   = [str(y) for y in (hist_pl.get("years") or [])]
    p_yrs   = [str(y) for y in (proj_pl.get("years") or [])]
    combined_pl_years = h_yrs + p_yrs

    if combined_pl_years:
        for label, h_key, p_key in [
            ("Revenue", "revenue", "revenue"),
            ("EBITDA",  "ebitda",  "ebitda"),
            ("PAT",     "pat",     "pat"),
        ]:
            vals = list(hist_pl.get(h_key) or []) + list(proj_pl.get(p_key) or [])
            if any(v is not None for v in vals):
                s = _series(label, "INR Cr", combined_pl_years, vals)
                if s:
                    series.append(s)

    if all_years:
        rev_growth = asmp.get("revenue_growth_pct") or {}
        # Revenue growth % series (used as fallback when absolute values absent)
        rev_vals = [rev_growth.get(y) for y in all_years]
        if any(v is not None for v in rev_vals):
            s = _series("Revenue Growth", "%", all_years, rev_vals)
            if s:
                series.append(s)

        # Historical ratio series — align to unified timeline (None for proj years)
        hist_idx = {y: i for i, y in enumerate(hist_years)}
        for label, key, unit in [
            ("EBITDA Margin", "ebitda_margin_pct", "%"),
            ("PAT Margin", "pat_margin_pct", "%"),
            ("ROE", "roe_pct", "%"),
            ("ROCE", "roce_pct", "%"),
        ]:
            raw_vals = hist.get(key)
            if not raw_vals:
                continue
            aligned = [raw_vals[hist_idx[y]] if y in hist_idx and hist_idx[y] < len(raw_vals) else None
                       for y in all_years]
            if any(v is not None for v in aligned):
                s = _series(label, unit, all_years, aligned)
                if s:
                    series.append(s)
            last = next((raw_vals[hist_idx[y]] for y in reversed(hist_years)
                         if y in hist_idx and hist_idx[y] < len(raw_vals) and raw_vals[hist_idx[y]] is not None), None)
            if last is not None:
                metrics[f"{key}_latest"] = str(last)

    # Headline numbers (skip upside_pct — recomputed from cmp+target to avoid rounding mismatch)
    for k in ("cmp", "target_price", "shares_cr"):
        v = model_json.get(k)
        if v is not None:
            metrics[k] = str(v)

    val = model_json.get("valuation") or {}
    for k in ("dcf_fair_value", "pe_fair_value", "ev_ebitda_fair_value", "blended_fair_value"):
        v = val.get(k)
        if v is not None:
            metrics[k] = str(v)

    if not metrics:
        metrics = {"placeholder": "0"}

    base["metrics"] = metrics
    base["series"] = series or _placeholder_series()
    if not series:
        warnings.append("financial-model JSON contained no usable annual series; using placeholder series")

    # Preserve richer reportgen-compatible model fields when the sidecar already
    # has them. The earlier mapper only kept a narrow v5 forecast slice, which
    # made the deck look empty even when the source JSON had SAARTHI, scenarios,
    # segments, risks, or strategy fields.
    for key in (
        "quarterly_series",
        "segments",
        "peers",
        "valuation_bands",
        "scenarios",
        "ratios",
        "saarthi",
        "operational",
        "governance",
        "timeline_events",
        "risk_items",
        "peers_detailed",
        "management_team",
        "forensic",
        "key_highlights",
        "competitive_advantages",
        "industry_tailwinds",
        "industry_risks",
        "trading_strategy",
        "historical_ratios",
        "historical_pl",
        "projected_pl",
    ):
        _copy_if_present(base, model_json, key)

    if "saarthi" not in base:
        if model_json.get("saarthi_scorecard"):
            base["saarthi"] = model_json["saarthi_scorecard"]
        elif "thesis" in model_json:
            thesis_data = model_json["thesis"] or {}
            if "saarthi_dimensions" in thesis_data or "saarthi_total" in thesis_data:
                base["saarthi"] = {
                    "total_score": thesis_data.get("saarthi_total", 70),
                    "max_score": 100,
                    "rating": thesis_data.get("saarthi_rating") or "BUY",
                    "dimensions": [
                        {
                            "code": d.get("key"),
                            "name": d.get("name"),
                            "score": d.get("score", 10),
                            "max_score": d.get("max_score", 15),
                            "assessment": d.get("rationale") or "",
                            "key_evidence": d.get("rationale") or "",
                        }
                        for d in thesis_data.get("saarthi_dimensions") or []
                    ]
                }

    if "segments" not in base:
        business_summary = model_json.get("business_summary") or {}
        if isinstance(business_summary, dict) and business_summary.get("segments"):
            base["segments"] = business_summary["segments"]

    if "valuation_bands" not in base:
        target_range = model_json.get("target_price_range") or {}
        if isinstance(target_range, dict) and target_range.get("base"):
            base["valuation_bands"] = [
                {
                    "method": "Target Price Range",
                    "low": str(target_range.get("low") or target_range.get("base")),
                    "base": str(target_range.get("base")),
                    "high": str(target_range.get("high") or target_range.get("base")),
                    "weight_pct": "100",
                    "notes": "Range supplied by the financial model sidecar.",
                }
            ]

    if "key_highlights" not in base and _list_of_dicts(model_json.get("key_highlights")):
        base["key_highlights"] = model_json["key_highlights"]

    if "competitive_advantages" not in base:
        comp = model_json.get("competitive_advantages")
        if isinstance(comp, list) and comp:
            base["competitive_advantages"] = comp
    return base


def _enrich_financial_model_for_house_deck(
    fin_model: dict,
    report: dict,
    sections: list[dict],
    metadata: dict,
) -> dict:
    """Guarantee renderer data refs for the fixed 15-slide Tikona deck.

    Claude and the mock planner both skip slides when refs such as scenarios,
    SAARTHI, or forensic data are absent. The approved report often contains the
    narrative even when the financial-model sidecar is thin, so we seed compact
    structured placeholders from the narrative to keep the house format intact.
    """
    narrative = _all_sections_text(sections)
    company_name = report.get("company_name") or "Company"
    cmp_val = _parse_number(metadata.get("cmp")) or 100
    target_val = _parse_number(metadata.get("target_price")) or cmp_val
    low_val = round(min(cmp_val, target_val) * 0.85, 1)
    high_val = round(max(cmp_val, target_val) * 1.15, 1)

    periods = (fin_model.get("series") or _placeholder_series())[0].get("periods") or ["FY24A", "FY25A", "FY26E"]

    if not fin_model.get("quarterly_series"):
        fin_model["quarterly_series"] = [
            {
                "name": "Quarterly Performance",
                "unit": "INR Cr",
                "periods": ["Q1", "Q2", "Q3", "Q4"],
                "values": ["0", "0", "0", "0"],
            }
        ]

    if not fin_model.get("ratios"):
        fin_model["ratios"] = [
            {
                "name": "Return Profile",
                "unit": "%",
                "periods": periods,
                "values": ["0" for _ in periods],
            }
        ]

    if not fin_model.get("valuation_bands"):
        fin_model["valuation_bands"] = [
            {
                "method": "Bear Case",
                "low": str(low_val),
                "base": str(low_val),
                "high": str(cmp_val),
                "weight_pct": "25",
                "notes": "Conservative execution and valuation assumptions.",
            },
            {
                "method": "Base Case",
                "low": str(cmp_val),
                "base": str(target_val),
                "high": str(high_val),
                "weight_pct": "50",
                "notes": "Analyst-approved target price and core thesis assumptions.",
            },
            {
                "method": "Bull Case",
                "low": str(target_val),
                "base": str(high_val),
                "high": str(high_val),
                "weight_pct": "25",
                "notes": "Upside scenario if demand drivers and execution improve.",
            },
        ]

    if not fin_model.get("scenarios"):
        fin_model["scenarios"] = [
            {"name": "Bear", "target_price": str(low_val), "probability_pct": "25", "notes": "Execution slows and valuation support weakens."},
            {"name": "Base", "target_price": str(target_val), "probability_pct": "50", "notes": "Core investment thesis plays out as expected."},
            {"name": "Bull", "target_price": str(high_val), "probability_pct": "25", "notes": "Catalysts accelerate and market confidence improves."},
        ]

    if not fin_model.get("segments"):
        business_text = _section_by_any(sections, ["business", "model", "company"]) or narrative
        fin_model["segments"] = [
            {
                "name": "Core Business",
                "description": _clean_prose(business_text, max_len=240),
            },
            {
                "name": "Growth Drivers",
                "description": _clean_prose(_section_by_any(sections, ["demand", "driver", "catalyst"]) or business_text, max_len=240),
            },
        ]

    if not fin_model.get("saarthi"):
        dimensions = [
            ("S", "Scalability of Core Engine"),
            ("A", "Addressable Market"),
            ("A", "Asymmetric Pricing Power"),
            ("R", "Reinvestment Quality"),
            ("T", "Track Record Through Adversity"),
            ("H", "Human Capital and Governance"),
            ("I", "Inflection Point Identification"),
        ]
        fin_model["saarthi"] = {
            "total_score": 70,
            "max_score": 100,
            "rating": metadata.get("rating") or "BUY",
            "dimensions": [
                {
                    "code": code,
                    "name": name,
                    "score": 10,
                    "max_score": 15,
                    "assessment": "Derived from approved analyst narrative.",
                    "key_evidence": _clean_prose(narrative, max_len=160),
                }
                for code, name in dimensions
            ],
        }

    if not fin_model.get("management_team"):
        mgmt_text = _section_by_any(sections, ["management", "governance", "forensic"]) or narrative
        fin_model["management_team"] = [
            {
                "name": "Management Team",
                "role": "Company leadership",
                "bio": _clean_prose(mgmt_text, max_len=260),
            }
        ]

    if not fin_model.get("forensic"):
        forensic_text = _section_by_any(sections, ["forensic", "governance", "risk"]) or narrative
        fin_model["forensic"] = {
            "category": "Monitor",
            "overall_assessment": _clean_prose(forensic_text, max_len=360),
            "violations": [
                {
                    "title": "Governance and forensic review",
                    "description": _clean_prose(forensic_text, max_len=260),
                    "severity": "MEDIUM",
                }
            ],
        }

    if not fin_model.get("key_highlights"):
        thesis_text = _section_by_any(sections, ["investment", "thesis", "idea"]) or narrative
        fin_model["key_highlights"] = [
            {"title": f"{company_name} investment idea", "body": item}
            for item in _bullets_from_text(thesis_text, limit=5, max_len=220)
        ]

    if not fin_model.get("competitive_advantages"):
        comp_text = _section_by_any(sections, ["competitive", "moat", "advantage"]) or narrative
        fin_model["competitive_advantages"] = _bullets_from_text(comp_text, limit=5, max_len=180)

    if not fin_model.get("industry_tailwinds"):
        industry_text = _section_by_any(sections, ["industry", "sector", "market"]) or narrative
        fin_model["industry_tailwinds"] = _bullets_from_text(industry_text, limit=4, max_len=180)

    if not fin_model.get("industry_risks"):
        risk_text = _section_by_any(sections, ["risk"]) or narrative
        fin_model["industry_risks"] = _bullets_from_text(risk_text, limit=4, max_len=180)

    if not fin_model.get("trading_strategy"):
        strategy_text = _section_by_any(sections, ["trading", "strategy", "exit", "entry"]) or narrative
        # Parse the section into Entry / Review / Exit sub-parts using heading markers
        entry_part, review_part, exit_part = _split_entry_review_exit(strategy_text)
        fin_model["trading_strategy"] = {
            "entry_range": "Accumulate selectively",
            "entry_rationale": _clean_prose(entry_part or strategy_text, max_len=1200),
            "position_size": "Risk-managed position",
            "review_frequency": _clean_prose(review_part, max_len=1200) if review_part else "Review quarterly against thesis milestones.",
            "review_metrics": _bullets_from_text(review_part or strategy_text, limit=6, max_len=200),
            "upside_exit": ["Review after target achievement"],
            "downside_exit": _clean_prose(exit_part, max_len=1200) if exit_part else "Exit if thesis invalidation triggers materialise.",
            "thesis_breaking_exits": _bullets_from_text(exit_part or _section_by_any(sections, ["risk"]) or strategy_text, limit=5, max_len=200),
        }

    return fin_model


def _build_approved_report_md(report: dict, sections: list[dict]) -> str:
    lines = [f"# {report.get('company_name') or 'Company'} — Investment Research Report", ""]
    for s in sections:
        title = (s.get("section_title") or s.get("section_key") or "Section").strip()
        body = (s.get("content") or "").strip()
        if not body:
            continue
        lines.append(f"## {title}")
        lines.append("")
        lines.append(body)
        lines.append("")
    if len(lines) <= 2:
        # Schema requires non-empty markdown with at least some prose
        lines.append("## Summary")
        lines.append("")
        lines.append("Report content pending.")
    return "\n".join(lines)


# ─────────── supabase i/o ─────────────────────────────────────────────────────


def _fetch_inputs(client, report_id: str, session_id: str) -> tuple[dict, dict, list[dict]]:
    rep = client.table("research_reports").select("*").eq("report_id", report_id).single().execute()
    if not rep.data:
        raise ValueError(f"research_reports row not found for report_id={report_id}")

    sess = client.table("research_sessions").select("*").eq("session_id", session_id).single().execute()
    if not sess.data:
        raise ValueError(f"research_sessions row not found for session_id={session_id}")

    secs = (
        client.table("research_sections")
        .select("section_key, section_title, content, sort_order")
        .eq("session_id", session_id)
        .eq("stage", "stage2")
        .order("sort_order")
        .execute()
    )
    return rep.data, sess.data, secs.data or []


def _download_model_json(client, ticker: str, warnings: list[str]) -> dict | None:
    """Fetch financial-models/{TICKER}/{TICKER}_model.json from research-reports-html."""
    path = f"financial-models/{ticker}/{ticker}_model.json"
    try:
        # supabase-py returns bytes
        data = client.storage.from_(MODEL_BUCKET).download(path)
        return json.loads(data.decode("utf-8"))
    except Exception as exc:
        logger.warning("model JSON download failed (%s): %s", path, exc)
        warnings.append(f"Could not load financial model JSON sidecar: {exc}")
        return None


def _download_model_excel(client, ticker: str, warnings: list[str], output_dir: Path) -> Path | None:
    """Fetch financial-models/{TICKER}/{TICKER}_model.xlsx from research-reports-html."""
    path = f"financial-models/{ticker}/{ticker}_model.xlsx"
    try:
        data = client.storage.from_(MODEL_BUCKET).download(path)
        out_path = output_dir / f"{ticker}_model.xlsx"
        with open(out_path, "wb") as f:
            f.write(data)
        return out_path
    except Exception as exc:
        logger.warning("model EXCEL download failed (%s): %s", path, exc)
        warnings.append(f"Could not load financial model EXCEL sidecar (required for native charts/tables): {exc}")
        return None


def _upload(client, local: Path, key: str, content_type: str) -> tuple[str, str]:
    with open(local, "rb") as fh:
        body = fh.read()
    client.storage.from_(PPTX_BUCKET).upload(
        path=key,
        file=body,
        file_options={"content-type": content_type, "upsert": "true"},
    )
    public = client.storage.from_(PPTX_BUCKET).get_public_url(key)
    return key, _strip_qmark(public)


def _upload_to_gdrive(local_path: Path, filename: str, folder_id: str, subfolder_name: str = None) -> dict | None:
    """Uploads file to Google Drive via n8n webhook upload-document."""
    import base64
    import requests
    import os
    
    # n8n webhook URL
    n8n_base = os.environ.get("N8N_BASE_URL") or "https://n8n.tikonacapital.com/webhook"
    url = f"{n8n_base.rstrip('/')}/upload-document"
    
    try:
        with open(local_path, "rb") as fh:
            file_bytes = fh.read()
        file_base64 = base64.b64encode(file_bytes).decode("utf-8")
        
        payload = {
            "folder_id": folder_id,
            "file_name": filename,
            "file_base64": file_base64
        }
        if subfolder_name:
            payload["subfolder_name"] = subfolder_name
        
        logger.info("Uploading %s to Google Drive folder %s", filename, folder_id)
        res = requests.post(url, json=payload, timeout=60)
        if res.status_code == 200:
            data = res.json()
            # Normalize response format from n8n
            # Webhook returns { status: "success", file: { id, webViewLink, ... } }
            if data.get("status") == "success" or "file" in data or "id" in data:
                file_info = data.get("file", data)
                if isinstance(file_info, list) and len(file_info) > 0:
                    file_info = file_info[0]
                file_id = file_info.get("id")
                slides_url = f"https://docs.google.com/presentation/d/{file_id}/edit"
                return {
                    "id": file_id,
                    "url": slides_url
                }
        logger.warning("Upload to Google Drive failed: %s %s", res.status_code, res.text)
    except Exception as e:
        logger.error("Error uploading to Google Drive: %s", e)
    return None


def sync_slides_to_pdf(report_id: str, ppt_file_id: str) -> dict:
    """
    Triggers n8n convert-to-pdf webhook to download the updated PPTX from Drive,
    converts it to PDF locally using LibreOffice, uploads it to Supabase storage,
    and updates research_reports.pptx_pdf_file_url.
    """
    import requests
    import os
    import tempfile
    from pathlib import Path
    from datetime import datetime, timezone
    
    client = get_service_client()
    
    # 1. Fetch report to get ticker symbol
    report_res = client.table("research_reports").select("nse_symbol").eq("report_id", report_id).execute()
    if not report_res.data:
        raise ValueError(f"Report not found for ID: {report_id}")
    ticker = (report_res.data[0].get("nse_symbol") or "UNKNOWN").upper()
    
    n8n_base = os.environ.get("N8N_BASE_URL") or "https://n8n.tikonacapital.com/webhook"
    
    # 2. Trigger convert-to-pdf n8n webhook
    convert_url = f"{n8n_base.rstrip('/')}/convert-to-pdf"
    logger.info("Triggering n8n convert-to-pdf (download) for file: %s", ppt_file_id)
    convert_res = requests.post(convert_url, json={"pptFileId": ppt_file_id}, timeout=180)
    if convert_res.status_code != 200:
        raise RuntimeError(f"n8n convert-to-pdf download failed ({convert_res.status_code}): {convert_res.text}")
    
    pptx_bytes = convert_res.content
    logger.info("Successfully downloaded PPTX binary from Drive via n8n (%d bytes)", len(pptx_bytes))
    
    # 3. Save PPTX bytes and convert to PDF locally using LibreOffice
    with tempfile.TemporaryDirectory(prefix="pptsync_") as tmp:
        tmp_dir = Path(tmp)
        temp_pptx = tmp_dir / f"report_{report_id}.pptx"
        temp_pptx.write_bytes(pptx_bytes)
        
        logger.info("Converting downloaded PPTX to PDF using LibreOffice...")
        pdf_path = _convert_pptx_to_pdf(temp_pptx, tmp_dir)
        if not pdf_path or not pdf_path.exists():
            raise RuntimeError("LibreOffice PDF conversion failed during sync")
            
        pdf_bytes = pdf_path.read_bytes()
        logger.info("Successfully converted PPTX to PDF locally (%d bytes)", len(pdf_bytes))
        
        # 4. Upload PDF binary to Supabase Storage
        ts = int(time.time())
        pdf_key = f"{ticker}/{report_id}/report_{ts}_sync.pdf"
        
        client.storage.from_(PPTX_BUCKET).upload(
            path=pdf_key,
            file=pdf_bytes,
            file_options={"content-type": "application/pdf", "upsert": "true"},
        )
        public_url_raw = client.storage.from_(PPTX_BUCKET).get_public_url(pdf_key)
        pdf_url = _strip_qmark(public_url_raw)
        logger.info("Uploaded synced PDF to Supabase storage: %s", pdf_url)
        
        # 5. Update research_reports table
        now_iso = datetime.now(timezone.utc).isoformat()
        client.table("research_reports").update({
            "pptx_pdf_file_path": pdf_key,
            "pptx_pdf_file_url": pdf_url,
            "updated_at": now_iso
        }).eq("report_id", report_id).execute()
        
        return {
            "status": "success",
            "message": "Slides synced and PDF updated successfully",
            "pptx_pdf_file_url": pdf_url,
            "pptx_pdf_file_path": pdf_key
        }



# ─────────── orchestrator ─────────────────────────────────────────────────────

def _classify_market_cap(market_cap_raw: str) -> str:
    mcap = _parse_number(market_cap_raw)
    if mcap is None or mcap <= 0:
        return "Large Cap"
    if mcap >= 20000:
        return "Large Cap"
    if mcap >= 5000:
        return "Mid Cap"
    return "Small Cap"


def _fmt_mcap(market_cap_raw: str) -> str:
    mcap = _parse_number(market_cap_raw)
    if mcap is None or mcap <= 0:
        return market_cap_raw or ""
    return f"₹{mcap:,.0f} Cr"


_SAARTHI_LETTER_ORDER = ("S", "A", "A", "R", "T", "H", "I")
_SAARTHI_CARD_KEYS = ("s", "a1", "a2", "r", "t", "h", "i")
_SAARTHI_DIMENSION_NAMES = {
    "s":  "Scalability of Core Engine",
    "a1": "Addressable Market & Adjacency",
    "a2": "Asymmetric Pricing Power",
    "r":  "Reinvestment Quality",
    "t":  "Track Record Through Adversity",
    "h":  "Human Capital & Institutional DNA",
    "i":  "Inflection Point Identification",
}


def _split_saarthi_framework(text: str, saarthi: dict) -> dict[str, str]:
    """Split a SAARTHI framework narrative into the 7 letter cards.

    Looks for letter prefixes like ``S —`` / ``A —`` / ``R —`` etc. at the start
    of lines or after sentence breaks. Falls back to the structured
    ``saarthi.dimensions`` data, then to a generic prompt, so every card always
    gets distinct copy even when the source narrative is missing or single-letter.
    """
    result: dict[str, str] = {k: "" for k in _SAARTHI_CARD_KEYS}

    cleaned_text = (text or "").strip()
    if cleaned_text:
        pattern = re.compile(
            r"(?:^|\n|(?<=[.!?\)\]]\s))\s*(?:\*\*|__)?\s*([SARTHIsarthi])\s*[—–\-:]\s+(?=[A-Za-z])",
        )
        matches = list(pattern.finditer(cleaned_text))
        segments: list[tuple[str, str]] = []
        for idx, match in enumerate(matches):
            start = match.start()
            end = matches[idx + 1].start() if idx + 1 < len(matches) else len(cleaned_text)
            body = cleaned_text[start:end].strip()
            segments.append((match.group(1).upper(), body))

        # Greedy mapping: walk expected letter sequence S,A,A,R,T,H,I and consume
        # the first matching segment we have not yet used.
        used = [False] * len(segments)
        for card_key, expected in zip(_SAARTHI_CARD_KEYS, _SAARTHI_LETTER_ORDER):
            for idx, (letter, body) in enumerate(segments):
                if used[idx] or letter != expected:
                    continue
                used[idx] = True
                result[card_key] = body
                break

    dims = saarthi.get("dimensions") or []
    for idx, card_key in enumerate(_SAARTHI_CARD_KEYS):
        if result[card_key]:
            continue
        if idx < len(dims):
            d = dims[idx]
            name = str(d.get("name") or _SAARTHI_DIMENSION_NAMES[card_key]).strip()
            score = d.get("score")
            max_score = d.get("max_score") or 15
            fmt_score = f"{int(score)}" if score is not None and score.is_integer() else f"{score}" if score is not None else ""
            fmt_max = f"{int(max_score)}" if max_score is not None and max_score.is_integer() else f"{max_score}" if max_score is not None else ""
            score_part = f"{fmt_score}/{fmt_max}" if fmt_score else ""
            head = f"{name}: {score_part}".strip(": ").rstrip()
            result[card_key] = (head + (". " + evidence if evidence else "")).strip()
        else:
            result[card_key] = _SAARTHI_DIMENSION_NAMES[card_key]

    return result


def _synthesise_saarthi_assessment(saarthi: dict) -> str:
    """Derive a short overall_assessment from dimension evidence when the field is absent."""
    dims = saarthi.get("dimensions") or []
    # Try to pick the first dimension that has real evidence text
    for d in dims:
        evidence = (d.get("key_evidence") or d.get("assessment") or "").strip()
        if evidence and len(evidence) > 20 and "Derived from" not in evidence:
            return evidence
    # Fall back to assembling from all dimension assessments
    snippets = []
    for d in dims:
        a = (d.get("assessment") or "").strip()
        if a and len(a) > 10:
            snippets.append(a)
    if snippets:
        return _clean_prose(" ".join(snippets), max_len=600)
    return "Quality score driven by strong scalability and resilient track record."


def _extract_score_from_text(text: str) -> tuple[float, float] | None:
    if not text:
        return None
    match = re.search(r"\b(\d+(?:\.\d+)?)\s*/\s*(\d+(?:\.\d+)?)\b", text)
    if match:
        try:
            return float(match.group(1)), float(match.group(2))
        except (ValueError, TypeError):
            pass
    return None


def map_replacements(company, metadata, fin_model, sections):
    thesis   = _section_by_any(sections, ["investment", "thesis", "summary"])
    rationale = _section_value(sections, "investment_rationale")
    industry = _section_by_any(sections, ["industry", "sector", "market"])
    idea     = _section_by_any(sections, ["idea", "model", "demand", "driver"])
    catalyst = _section_by_any(sections, ["catalyst", "driver"])
    competitive = _section_by_any(sections, ["competitive", "moat", "advantage"])
    mgmt     = _section_by_any(sections, ["management", "forensic", "governance"])
    # Broaden peer search: catch all common DB naming conventions
    peer     = _section_by_any(sections, [
        "peer", "competition", "comparable",
        "competitive_analysis", "competitive_landscape",
        "sector_analysis", "industry_analysis",
    ])
    saarthi  = fin_model.get("saarthi") or {}

    # Extract SAARTHI card texts and parse the scores from text to keep box values in sync
    saarthi_text = _section_by_any(sections, ["saarthi"]) or _section_value(sections, "saarthi_framework")
    saarthi_cards_text = _split_saarthi_framework(saarthi_text, saarthi)
    dims = saarthi.get("dimensions") or []
    computed_total = 0
    has_custom_scores = False

    for idx, card_key in enumerate(_SAARTHI_CARD_KEYS):
        card_content = saarthi_cards_text.get(card_key, "")
        score_info = _extract_score_from_text(card_content)
        if score_info:
            score, max_score = score_info
            has_custom_scores = True
            if idx < len(dims):
                dims[idx]["score"] = score
                dims[idx]["max_score"] = max_score
            computed_total += score
        else:
            if idx < len(dims):
                existing_score = dims[idx].get("score")
                if existing_score is not None:
                    try:
                        computed_total += float(existing_score)
                    except (ValueError, TypeError):
                        pass

    if has_custom_scores or dims:
        saarthi["total_score"] = int(round(computed_total))

    thesis_bullets = _bullets_from_text(thesis, limit=4)

    # 6 idea bullets: combine idea + catalyst + competitive + key_highlights
    idea_combined = idea or catalyst or competitive or _all_sections_text(sections)
    idea_bullets  = _bullets_from_text(idea_combined, limit=6)
    # pad from key_highlights if we have fewer than 6
    highlights = fin_model.get("key_highlights") or []
    if isinstance(highlights, list) and len(idea_bullets) < 6:
        for h in highlights:
            body = h.get("body", h) if isinstance(h, dict) else str(h)
            bullet = _clean_prose(str(body), max_len=180)
            if bullet and bullet not in idea_bullets:
                idea_bullets.append(bullet)
            if len(idea_bullets) >= 6:
                break

    # Competitive advantage bullets — strip leading JSON/markdown artefacts
    comp_raw = fin_model.get("competitive_advantages") or []
    if isinstance(comp_raw, list) and comp_raw:
        comp_bullets = [
            _clean_prose(re.sub(r'^[\s{"\' *\-•]+', '', str(c)), max_len=280)
            for c in comp_raw[:4]
        ]
    else:
        comp_text   = competitive or idea or _all_sections_text(sections)
        comp_bullets = _bullets_from_text(comp_text, limit=4)

    # Synthesise saarthi overall_assessment once (used by both summary + content fields)
    saarthi_assessment = (
        saarthi.get("overall_assessment")
        or _synthesise_saarthi_assessment(saarthi)
    )

    # Peer text — broadened fallback so valuation/sector sections are used when
    # no section is explicitly labelled "peer" or "competition".
    peer_text = peer or _section_by_any(sections, [
        "valuation", "peers", "sector", "industry", "market",
    ])
    peer_sents = _bullets_from_text(peer_text, limit=4, max_len=220) if peer_text else []

    # Market cap helpers
    mcap_raw   = metadata.get("market_cap", "")
    m_category = _classify_market_cap(mcap_raw)
    m_cap_disp = _fmt_mcap(mcap_raw)

    # Upside display — round to 1 decimal so we don't render "-55.529999%".
    upside_raw = str(metadata.get("upside_pct", "") or "")
    _upside_num = _parse_number(upside_raw)
    if _upside_num is not None:
        upside_disp = f"{_upside_num:.1f}%"
    else:
        upside_disp = (upside_raw + "%") if upside_raw and not upside_raw.endswith("%") else upside_raw

    today_str = date.today().strftime("%d %b %Y")
    thesis_sentences = _sentences(thesis, limit=4, max_len=260)
    thesis_box_headings = [
        "Market Position",
        "Expansion Trigger",
        "Margin Outlook",
        "Why It Matters",
    ]
    thesis_box_texts = [
        _truncate_words(thesis_sentences[0] if len(thesis_sentences) > 0 else thesis_bullets[0] if len(thesis_bullets) > 0 else thesis, 28, max_len=240),
        _truncate_words(thesis_sentences[1] if len(thesis_sentences) > 1 else thesis_bullets[1] if len(thesis_bullets) > 1 else thesis, 28, max_len=240),
        _truncate_words(thesis_sentences[2] if len(thesis_sentences) > 2 else thesis_bullets[2] if len(thesis_bullets) > 2 else thesis, 28, max_len=240),
    ]
    thesis_bottom_summary = _truncate_words(
        thesis_sentences[3] if len(thesis_sentences) > 3 else saarthi_assessment or thesis,
        45,
        max_len=360,
    )
    industry_tailwinds_text = _section_by_any(sections, ["tailwind", "tailwinds"]) or "\n\n".join(fin_model.get("industry_tailwinds") or [])
    industry_risks_text = _section_by_any(sections, ["industry_risks", "risks", "risk"]) or "\n\n".join(fin_model.get("industry_risks") or [])
    metric_chips = _metric_chips(metadata, fin_model, company)
    business_model_text = _section_value(sections, "business_model") or idea or competitive or thesis
    business_cards = _sentences(business_model_text, limit=6, max_len=240)
    while len(business_cards) < 6 and len(metric_chips) > len(business_cards):
        business_cards.append(metric_chips[len(business_cards)])
    business_cards = [
        _truncate_words(item, 30, max_len=240)
        for item in business_cards[:6]
    ]
    # Combine thesis and rationale paragraphs, preventing duplicate insertion 
    # if both variables resolved to the same underlying section.
    thesis_parts = []
    if rationale:
        thesis_parts.append(rationale)
    if thesis and thesis != rationale:
        thesis_parts.append(thesis)
    if metadata.get("target_price") and metadata.get("cmp"):
        thesis_parts.append(f"Target price ₹{metadata.get('target_price', '')} versus CMP ₹{metadata.get('cmp', '')} implies {upside_disp} upside.")

    thesis_panel_text = _clean_prose(
        "\n\n".join(thesis_parts),
        max_len=1650,
    )
    company_overview_text = _clean_prose(company.get("description", ""), max_len=1500)
    business_model_text = _clean_prose(_section_value(sections, "business_model"), max_len=900)
    top_overview_parts = _sentences(company_overview_text, limit=6, max_len=240)
    if business_model_text:
        top_overview_parts.extend(_sentences(business_model_text, limit=3, max_len=220))
    top_overview = " ".join(top_overview_parts[:7]).strip() or company_overview_text

    op = fin_model.get("operational") or {}
    overview_metrics: list[str] = []
    countries = op.get("countries_of_operation") or []
    plants_india = op.get("plants_india") or []
    plants_overseas = op.get("plants_overseas") or []
    utils = op.get("capacity_utilisation_pct") or []
    volumes = op.get("volume_segments") or {}
    latest_total_volume = 0
    if volumes:
        latest_idx = _last_actual_index([str(y).strip() for y in (op.get("years") or [])])
        if latest_idx is not None:
            for vals in volumes.values():
                seq = [float(v) if v is not None else 0.0 for v in (vals or [])]
                if latest_idx < len(seq):
                    latest_total_volume += seq[latest_idx]
    if countries and countries[-1] is not None:
        try:
            overview_metrics.append(f"Operates across {int(countries[-1])} countries")
        except (ValueError, TypeError):
            pass
    if plants_india or plants_overseas:
        try:
            p_ind = int(plants_india[-1]) if (plants_india and plants_india[-1] is not None) else 0
            p_ovr = int(plants_overseas[-1]) if (plants_overseas and plants_overseas[-1] is not None) else 0
            total_plants = p_ind + p_ovr
            if total_plants:
                overview_metrics.append(f"{total_plants} recycling plants")
        except (ValueError, TypeError):
            pass
    if latest_total_volume:
        overview_metrics.append(f"{int(latest_total_volume):,} MT latest throughput")
    if utils and utils[-1] is not None:
        try:
            latest_util = float(utils[-1])
            if latest_util <= 1.0:
                latest_util *= 100
            overview_metrics.append(f"{latest_util:.0f}% utilisation")
        except (ValueError, TypeError):
            pass
    if m_cap_disp:
        overview_metrics.append(f"{m_cap_disp} market cap")
    bottom_overview = " | ".join(overview_metrics[:4])
    if business_model_text:
        business_snips = _sentences(business_model_text, limit=2, max_len=180)
        if business_snips:
            bottom_overview = f"{bottom_overview}. {' '.join(business_snips)}".strip(". ")
    bottom_overview = _clean_prose(bottom_overview, max_len=520)
    management_sentences = _sentences(mgmt, limit=10, max_len=180)
    management_headings = [
        "Capital Allocation",
        "Execution Track Record",
        "Expansion Discipline",
        "Leadership Quality",
        "Funding Approach",
        "Margin Focus",
        "Strategic Vision",
        "Risk Controls",
        "Communication Style",
        "Shareholder Alignment",
    ]
    # Card width on slide 11 fits ~35 words cleanly; lower than that produces
    # ugly mid-clause cut-offs like "focusing on." even when distinct sentences
    # are available.
    management_cards = [_truncate_words(text, 35, max_len=230) for text in management_sentences[:10]]
    while len(management_cards) < 10:
        management_cards.append(_truncate_words(mgmt, 35, max_len=230))
    governance_text = _section_by_any(sections, ["governance", "forensic", "indicator"]) or mgmt
    governance_cards = _sentences(governance_text, limit=6, max_len=170)
    while len(governance_cards) < 6:
        governance_cards.append(_truncate_words(governance_text, 24, max_len=170))

    replacements = {
        # ── Slide 1: Cover ────────────────────────────────────────────────────
        "company_name": company.get("name", "Company"),
        "nse_code":     company.get("ticker", ""),
        "cmp":          fmt_number(metadata.get("cmp", "")),
        "target":       fmt_number(metadata.get("target_price", "")),
        "m_cap":        m_cap_disp,
        "m_category":   m_category,
        "saarthi_s":    str(saarthi.get("total_score", "70")),
        "tagline":      "Initiation Report",
        # ── Slide 1 atomic placeholders the cover template references ────────
        # Both upside and rating are exposed as separate keys here (in addition
        # to the legacy `up` / `m_category` / etc.) because the cover slide
        # uses {{upside}} and {{rating}} tokens directly.
        "upside":       upside_disp,
        "rating":       _normalize_rating(metadata.get("rating", "") or fin_model.get("rating", "") or "BUY"),
        # ── Slides 1, 4: Investment Thesis ────────────────────────────────────
        "investment_thesis_heading": "Investment Thesis",
        "investment_thesis":         thesis_panel_text,
        "saarthi_summary_heading":   "SAARTHI Overview",
        "saarthi_summary":           thesis_bottom_summary,
        "1": thesis_bullets[0] if len(thesis_bullets) > 0 else "",
        "2": thesis_bullets[1] if len(thesis_bullets) > 1 else "",
        "3": thesis_bullets[2] if len(thesis_bullets) > 2 else "",
        "4": thesis_bullets[3] if len(thesis_bullets) > 3 else "",
        # ── Slide 5: Industry Analysis ────────────────────────────────────────
        "date":     today_str,
        " date ":   today_str,
        "cell":     fmt_number(metadata.get("cmp", "")),
        "cell_cap": m_category,
        "mod_cap":  m_cap_disp,
        "mod":      metadata.get("upside_pct", ""),
        "tar_pr":   fmt_number(metadata.get("target_price", "")),
        "tar":      fmt_number(metadata.get("target_price", "")),  # Slide 5 uses {{tar}}
        "buy":      fmt_number(metadata.get("cmp", "")),
        "up":       upside_disp,
        "industry_structure": "\n".join(_bullets_from_text(industry, limit=5, max_len=145)),
        "key_industry":       "\n".join(_bullets_from_text(industry_tailwinds_text, limit=5, max_len=145)),
        "key_industry_risk":  "\n".join(_bullets_from_text(industry_risks_text, limit=5, max_len=145)),
        # ── Slide 6: Company Overview ─────────────────────────────────────────
        "COMPANY_OVERVIEW": company_overview_text,
        "__slide6_top_overview": top_overview,
        "__slide6_bottom_overview": bottom_overview,
        # ── Slide 7: Business Ideas ───────────────────────────────────────────
        "p1": business_cards[0] if len(business_cards) > 0 else "",
        "p2": business_cards[1] if len(business_cards) > 1 else "",
        "p3": business_cards[2] if len(business_cards) > 2 else "",
        "p4": business_cards[3] if len(business_cards) > 3 else "",
        "p5": business_cards[4] if len(business_cards) > 4 else "",
        "p6": business_cards[5] if len(business_cards) > 5 else "",
        # ── Slide 8: Competitive Advantages ──────────────────────────────────
        "competitive_advantage_1": comp_bullets[0] if len(comp_bullets) > 0 else "",
        "competitive_advantage_2": comp_bullets[1] if len(comp_bullets) > 1 else "",
        "competitive_advantage_3": comp_bullets[2] if len(comp_bullets) > 2 else "",
        "competitive_advantage_4": comp_bullets[3] if len(comp_bullets) > 3 else "",
        "industry_tailwinds":     _clean_prose(industry, max_len=400),
        # ── Slide 9: Peer Comparison ──────────────────────────────────────────
        "peer_comparision": _clean_prose(peer_text, max_len=600),
        "peer_para1": _clean_prose(
            peer_sents[0] if peer_sents else (peer_text[:280] if peer_text else ""), max_len=300),
        "peer_para2": _clean_prose(
            peer_sents[1] if len(peer_sents) > 1 else (peer_text[280:560] if peer_text and len(peer_text) > 280 else ""),
            max_len=300,
        ),
        # ── Slide 10: Management ──────────────────────────────────────────────
        "management_commentry_heading": "Management Analysis",
        "management_content":           _clean_prose(mgmt, max_len=1500),
        "__slide11_headings":          management_headings,
        "__slide11_contents":          management_cards,
        # ── Slide 11: Governance ──────────────────────────────────────────────
        "indicators": _clean_prose(
            governance_text,
            max_len=400,
        ),
        "__slide12_indicator_cards": governance_cards[:6],
        # ── Slide 15: SAARTHI ─────────────────────────────────────────────────
        "saarthi_heading":  "SAARTHI Score Analysis",
        "saarthi_content":  _clean_prose(saarthi_assessment, max_len=600),
        # ── Slide 7: Company Timeline ─────────────────────────────────────────
        "COMPANY_TIMELINE": _clean_prose(
            _section_by_any(sections, ["timeline", "history", "milestones", "journey"]) or
            _section_by_any(sections, ["company", "overview", "business"]),
            max_len=800,
        ),
        # ── Slide 7: Business Idea Paragraphs (para_1..para_6) ────────────────
        "para_1": metric_chips[0] if len(metric_chips) > 0 else "",
        "para_2": metric_chips[1] if len(metric_chips) > 1 else "",
        "para_3": metric_chips[2] if len(metric_chips) > 2 else "",
        "para_4": metric_chips[3] if len(metric_chips) > 3 else "",
        "para_5": metric_chips[4] if len(metric_chips) > 4 else "",
        "para_6": metric_chips[5] if len(metric_chips) > 5 else "",
        # ── Slide 14: Financial Commentary ────────────────────────────────────
        "financial_commentry": _format_two_paragraphs(
            _section_by_any(sections, ["financial", "earnings", "revenue", "profit"]) or
            _all_sections_text(sections)
        ),
        # ── Slide 15: Valuations Commentary ──────────────────────────────────
        "commentry": _format_two_paragraphs(
            _section_by_any(sections, ["valuation", "dcf", "pe_ratio", "fair_value"]) or
            _section_by_any(sections, ["investment", "thesis"]) or
            _all_sections_text(sections)
        ),
    }

    replacements.update({
        "investment_thesis_s1": thesis_panel_text,
        "investment_ideas_1": thesis_bullets[0] if len(thesis_bullets) > 0 else "",
        "investment_ideas_2": thesis_bullets[1] if len(thesis_bullets) > 1 else "",
        "investment_ideas_3": thesis_bullets[2] if len(thesis_bullets) > 2 else "",
        "investment_ideas_4": thesis_bullets[3] if len(thesis_bullets) > 3 else "",
        "investment_thesis_heading_s4": "Investment Thesis",
        "investment_thesis_s4": thesis_panel_text,
        "key_catalyst_1": thesis_box_texts[0] if len(thesis_box_texts) > 0 else "",
        "key_catalyst_2": thesis_box_texts[1] if len(thesis_box_texts) > 1 else "",
        "key_catalyst_3": thesis_box_texts[2] if len(thesis_box_texts) > 2 else "",
        "key_catalyst_heading_1": thesis_box_headings[0] if len(thesis_box_headings) > 0 else "Key Catalyst 1",
        "key_catalyst_heading_2": thesis_box_headings[1] if len(thesis_box_headings) > 1 else "Key Catalyst 2",
        "key_catalyst_heading_3": thesis_box_headings[2] if len(thesis_box_headings) > 2 else "Key Catalyst 3",
        "key catalyst_heading_1": thesis_box_headings[0] if len(thesis_box_headings) > 0 else "Key Catalyst 1",
        "key catalyst_heading_2": thesis_box_headings[1] if len(thesis_box_headings) > 1 else "Key Catalyst 2",
        "key catalyst_heading_3": thesis_box_headings[2] if len(thesis_box_headings) > 2 else "Key Catalyst 3",
        "saarthi_summary_s4": thesis_bottom_summary,
        "key_industry_tailwainds": _clean_prose(industry_tailwinds_text, max_len=900),
        "key_industry_tailwinds": _clean_prose(industry_tailwinds_text, max_len=900),
        "key_industry_risks": _clean_prose(industry_risks_text, max_len=900),
        "KPI_heading_1": "Market Cap",
        "KPI_heading_2": "CMP",
        "KPI_heading_3": "Target Price",
        "KPI_heading_4": "Upside",
        "KPI_heading_5": "Category",
        "KPI_heading_6": "SAARTHI",
        "KPI_1": m_cap_disp,
        "KPI_2": fmt_number(metadata.get("cmp", "")),
        "KPI_3": fmt_number(metadata.get("target_price", "")),
        "KPI_4": upside_disp,
        "KPI_5": m_category,
        "KPI_6": str(saarthi.get("total_score", "70")),
        "company_overview": top_overview,
        "competitive_moat_1": comp_bullets[0] if len(comp_bullets) > 0 else "",
        "competitive_moat_2": comp_bullets[1] if len(comp_bullets) > 1 else "",
        "key_insights": bottom_overview,
        "company_timeline": replacements["COMPANY_TIMELINE"],
        "investment_thesis_detailed": thesis_panel_text,
        # Right-top box on slide 10. Must be visibly distinct from the left
        # `investment_thesis_detailed` panel — pull from the catalyst / driver
        # section first, then fall back to industry tailwinds, then to the
        # short thesis snippets as a last resort.
        "key_catalyst": _clean_prose(
            catalyst
            or industry_tailwinds_text
            or " ".join(thesis_box_texts[:3]),
            max_len=900,
        ),
        "business_model_1": business_cards[0] if len(business_cards) > 0 else "",
        "business_model_2": business_cards[1] if len(business_cards) > 1 else "",
        "business_model_3": business_cards[2] if len(business_cards) > 2 else "",
        "business_model_4": business_cards[3] if len(business_cards) > 3 else "",
        "business_model_5": business_cards[4] if len(business_cards) > 4 else "",
        "business_model_6": business_cards[5] if len(business_cards) > 5 else "",
        "competitive_advantage": _clean_prose(" ".join(comp_bullets[:4]), max_len=1200),
        # New narrative heading placeholders introduced on slides 9 + 10.
        # If the LLM copywriter does not populate these via cs_ppt_data, the
        # template would otherwise show the literal {{...}} token — write a
        # sensible default so the slide remains visually clean.
        "competitive_advantage_heading": "Strategic Edge",
        "investment_thesis_detailed_heading": "Long-Term Thesis Drivers",
        "management_commentry_heading_1": management_headings[0],
        "management_commentry_heading_2": management_headings[1],
        "management_commentry_heading_3": management_headings[2],
        "management_commentry_heading_4": management_headings[3],
        "management_commentry_heading_5": management_headings[4],
        "management_commentry_heading_6": management_headings[5],
        "management_commentry_heading_7": management_headings[6],
        "management_commentry_heading_8": management_headings[7],
        "management_content_1": management_cards[0],
        "management_content_2": management_cards[1],
        "management_content_3": management_cards[2],
        "management_content_4": management_cards[3],
        "management_content_5": management_cards[4],
        "management_content_6": management_cards[5],
        "management_content_7": management_cards[6],
        "management_content_8": management_cards[7],
        "indicators_1": governance_cards[0],
        "indicators_2": governance_cards[1],
        "indicators_3": governance_cards[2],
        "indicators_4": governance_cards[3],
        "indicators_5": governance_cards[4],
        "indicators_6": governance_cards[5],
        "forecast_assumptions": _format_two_paragraphs(
            _section_by_any(sections, ["forecast", "assumption", "growth", "capex", "working capital"]) or
            business_model_text or thesis_panel_text
        ),
        "financial_commentary": replacements["financial_commentry"],
        "valuation_commentary": replacements["commentry"],
        **{
            f"saarthi_{key}_content": _clean_prose(value, max_len=420)
            for key, value in _split_saarthi_framework(
                _section_by_any(sections, ["saarthi"]) or _section_value(sections, "saarthi_framework"),
                saarthi,
            ).items()
        },
        "saarthi_summary_s16": _truncate_words(saarthi_assessment, 65, max_len=450),
    })

    # ── SAARTHI per-dimension scores (Slide 16: {{s_s}}, {{a1_s}}, ...) ──────
    _saarthi_dims_list = saarthi.get("dimensions") or []
    for _idx, _card_key in enumerate(_SAARTHI_CARD_KEYS):
        _placeholder = f"{_card_key}_s"
        if _idx < len(_saarthi_dims_list):
            _d = _saarthi_dims_list[_idx]
            _score = _d.get("score")
            _max = _d.get("max_score") or 15
            if _score is not None:
                try:
                    replacements[_placeholder] = f"{float(_score):.0f}/{int(_max)}"
                    continue
                except (TypeError, ValueError):
                    pass
        replacements[_placeholder] = ""

    # ── Scenario data (Slides 16, 18) ─────────────────────────────────────────
    scenarios = fin_model.get("scenarios") or []
    bear_tp_f = 0.0
    for s in scenarios:
        name  = str(s.get("name", "")).lower()
        notes = _clean_prose(str(s.get("notes", "")), max_len=200)
        prob  = str(s.get("probability_pct", "") or "")
        tp    = str(s.get("target_price", "") or "")
        prob_disp = (prob + "%") if prob and not prob.endswith("%") else prob
        tp_fmt = fmt_number(tp)
        if "bull" in name:
            replacements["valuation_bull"] = tp_fmt
            replacements["bull"]           = tp_fmt
            replacements["bull_p"]         = prob_disp
            replacements["bull_content"]   = notes
        elif "bear" in name:
            replacements["valuation_bear"] = tp_fmt
            replacements["bear"]           = tp_fmt
            replacements["bear_p"]         = prob_disp
            replacements["bear_content"]   = notes
            bear_tp_f = _parse_number(tp) or 0.0
        elif "base" in name:
            replacements["base"]           = tp_fmt
            replacements["valuation_base"] = tp_fmt
            replacements["base_p"]         = prob_disp
            replacements["base_content"]   = notes

    # ── Trading Strategy (Slide 18) ───────────────────────────────────────────
    trading = fin_model.get("trading_strategy") or {}
    entry_text = _clean_prose(
        str(trading.get("entry_rationale") or trading.get("entry_range") or
            "Accumulate at current market price with defined risk."),
        max_len=1200,
    )
    review_metrics = trading.get("review_metrics") or []
    review_text = _clean_prose(
        str(trading.get("review_frequency") or
            "; ".join(str(m) for m in review_metrics[:2]) or
            "Review quarterly against thesis milestones."),
        max_len=1200,
    )
    exits = trading.get("thesis_breaking_exits") or []
    exit_text = _clean_prose(
        str(trading.get("downside_exit") or
            "; ".join(str(x) for x in exits[:2]) or
            "Exit on thesis invalidation or sustained breach of support."),
        max_len=1200,
    )
    replacements["entry_strategy_1"]  = entry_text
    replacements["review_strategy_2"] = review_text
    replacements["exit_strategy_3"]   = exit_text
    replacements["entry_strategy"] = entry_text
    replacements["review_strategy"] = review_text
    replacements["exit_strategy"] = exit_text

    # Slide 18 price analytics — downside % and stop-loss derived from bear scenario
    cmp_val_f = _parse_number(metadata.get("cmp") or "") or 0.0
    if bear_tp_f <= 0:
        bear_tp_f = round(cmp_val_f * 0.85, 1)
    stp_loss_val = round(bear_tp_f, 1)
    if cmp_val_f > 0:
        down_pct  = round((bear_tp_f - cmp_val_f) / cmp_val_f * 100, 1)
        # Template already provides the trailing "%", so emit a bare number.
        down_disp = f"{down_pct:.1f}"
    else:
        down_disp = ""
    replacements["stp_loss"] = fmt_number(stp_loss_val)
    replacements["down"]     = down_disp
    upside_num = _parse_number(metadata.get("upside_pct") or "") or 0.0
    downside_abs = abs(_parse_number(down_disp) or 0.0)
    rr_ratio = round((upside_num / downside_abs), 1) if downside_abs > 0 else 0.0
    replacements["pnt"]      = f"{rr_ratio:.1f}" if rr_ratio else ""
    replacements["up"]       = f"{upside_num:.1f}" if upside_num else ""
    replacements["__slide4_right_headings"] = thesis_box_headings
    replacements["__slide4_right_texts"] = thesis_box_texts

    # Normalise bullet-shaped placeholders so the renderer's per-line paragraph
    # emitter (in _replace_text_in_frame) gets one bullet per line — without
    # this, a single paragraph with embedded "• " markers would render as one
    # run of prose on the slide instead of the visible bullet list the
    # template was designed for.
    for _bk in _BULLET_LIST_KEYS:
        if _bk in replacements and isinstance(replacements[_bk], str):
            replacements[_bk] = _ensure_bullet_lines(replacements[_bk])

    return replacements


def _sync_equivalent_keys(d: dict) -> None:
    """Synchronise canonical schema keys and alternative UI/template keys.

    If a canonical key is present, its value is copied to the alternative key.
    If an alternative key is present, its value is copied to the canonical key.
    """
    mappings = {
        "entry_strategy": "entry_strategy_1",
        "review_strategy": "review_strategy_2",
        "exit_strategy": "exit_strategy_3",
        "financial_commentary": "financial_commentry",
        "valuation_commentary": "commentry",
        "key_industry_tailwinds": "key_industry",
        "key_industry_risks": "key_industry_risk",
        "company_overview": "COMPANY_OVERVIEW",
        "saarthi_summary_s16": "saarthi_summary",
    }
    # Sync canonical -> alternative
    for canonical, alt in mappings.items():
        if canonical in d and d[canonical]:
            d[alt] = d[canonical]
    # Sync alternative -> canonical
    for canonical, alt in mappings.items():
        if alt in d and d[alt] and (canonical not in d or not d[canonical]):
            d[canonical] = d[alt]


# ── Placeholder preview (called by /preview-placeholders) ─────────────────────

def preview_ppt_placeholders(report_id: str, session_id: str, *, ignore_overrides: bool = False) -> dict:
    """Compute all text placeholder values without generating the PPTX.

    Also merges any previously saved overrides from `cs_ppt_data` so the UI
    shows the last confirmed values when the user re-opens the panel.
    """
    warnings: list[str] = []
    client = get_service_client()
    report, session, sections = _fetch_inputs(client, report_id, session_id)
    ticker = (report.get("nse_symbol") or "UNKNOWN").upper()

    company  = _build_company(report, session, sections)
    metadata = _build_metadata(report, sections)
    model_json = _download_model_json(client, ticker, warnings)
    fin_model  = _build_financial_model(ticker, model_json, warnings)
    fin_model  = _enrich_financial_model_for_house_deck(fin_model, report, sections, metadata)

    placeholders = map_replacements(company, metadata, fin_model, sections)

    # Apply slide-specific copy from ppt_content_json if available
    ppt_copy = session.get("ppt_content_json")
    if isinstance(ppt_copy, str):
        try:
            ppt_copy = json.loads(ppt_copy)
        except Exception:
            ppt_copy = None
    if isinstance(ppt_copy, dict) and ppt_copy:
        applied = 0
        for k, v in ppt_copy.items():
            if v is None:
                continue
            value = str(v).strip()
            if not value:
                continue
            placeholders[k] = value
            applied += 1
        _sync_equivalent_keys(placeholders)
        logger.info("Preview: Applied %d slide-copy values from ppt_content_json", applied)

    # Merge previously-saved overrides so the UI shows confirmed values
    saved_raw = report.get("cs_ppt_data") or ""
    has_saved = bool(saved_raw)
    if saved_raw and not ignore_overrides:
        try:
            saved = json.loads(saved_raw)
            if isinstance(saved, dict):
                placeholders.update(saved)
                _sync_equivalent_keys(placeholders)
        except Exception:
            pass

    return {
        "status": "success",
        "placeholders": placeholders,
        "has_saved_overrides": has_saved if not ignore_overrides else False,
        "warnings": warnings,
    }

# ── Per-slide chart data mapping ──────────────────────────────────────────────
# Maps slide_type → which financial series names to include in native PPTX charts.
# The first series' periods become the category labels.
_SLIDE_CHART_SERIES: dict[str, list[str]] = {
    "cover":               ["Revenue", "EBITDA", "PAT"],
    "story_charts":        ["Revenue", "EBITDA", "PAT"],
    "earnings_forecast":   ["Revenue", "EBITDA", "PAT"],
    "financial_highlights":["EBITDA Margin", "PAT Margin", "ROE"],
    "valuation":           ["Revenue", "EBITDA"],
    "business_segments":   ["Revenue", "EBITDA"],
    "industry":            ["Revenue"],
}

# Maps table_key → which financial series to write into a PPTX table shape.
# "table_key" is either the shape name suffix after "table:" or the slide_type.
_TABLE_SERIES: dict[str, list[str]] = {
    "earnings_forecast":   ["Revenue", "EBITDA", "PAT", "EBITDA Margin", "PAT Margin"],
    "financial_highlights":["EBITDA Margin", "PAT Margin", "ROE", "ROCE"],
    "valuation":           ["Revenue", "EBITDA", "PAT"],
}


def _detect_slide_type(slide) -> str:
    """Read slide type from a shape named 'slide_type' in the template.

    In PowerPoint, open the Selection Pane and name one shape exactly
    'slide_type', with its text content set to e.g. 'earnings_forecast'.
    Falls back to 'generic' when no such shape is found.
    """
    for shape in slide.shapes:
        if (shape.name or "").lower().strip() == "slide_type":
            if hasattr(shape, "text_frame") and shape.text_frame:
                return shape.text_frame.text.strip().lower()
    return "generic"


def _build_pptx_chart_data(slide_type: str, fin_model: dict) -> "CategoryChartData | None":
    """Build a CategoryChartData for native PPTX chart replacement, per slide type."""
    wants = _SLIDE_CHART_SERIES.get(slide_type, ["Revenue", "EBITDA", "PAT"])
    series_list = fin_model.get("series") or []
    matching = [
        s for s in series_list
        if any(w.lower() in (s.get("name") or "").lower() for w in wants)
    ]
    if not matching:
        matching = series_list[:3]
    if not matching:
        return None

    periods = matching[0].get("periods") or []
    if not periods:
        return None

    cd = CategoryChartData()
    cd.categories = periods
    for s in matching[:3]:
        vals = []
        for v in (s.get("values") or []):
            try:
                vals.append(float(str(v).replace(",", "")) if v else 0.0)
            except (ValueError, TypeError):
                vals.append(0.0)
        cd.add_series(s.get("name", "Metric"), vals)
    return cd


def _fill_table_from_model(table, table_key: str, fin_model: dict) -> None:
    """Write financial series data into an existing PPTX table shape.

    Row 0 is treated as the header (Metric | period1 | period2 ...).
    Subsequent rows receive one series each. Only writes; never adds rows/cols.
    """
    wants = _TABLE_SERIES.get(table_key)
    if not wants:
        return

    series_list = fin_model.get("series") or []
    rows_data = [
        s for s in series_list
        if any(w.lower() in (s.get("name") or "").lower() for w in wants)
    ]
    if not rows_data:
        return

    periods = rows_data[0].get("periods") or []
    n_data_cols = min(len(periods), len(table.columns) - 1)
    n_data_rows = min(len(rows_data), len(table.rows) - 1)

    if n_data_cols <= 0 or n_data_rows <= 0:
        return

    # Header row
    try:
        table.cell(0, 0).text = "Metric"
        for ci, p in enumerate(periods[:n_data_cols]):
            table.cell(0, ci + 1).text = str(p)
    except Exception as e:
        logger.warning("Table header write failed: %s", e)
        return

    # Data rows
    for ri, s in enumerate(rows_data[:n_data_rows]):
        try:
            table.cell(ri + 1, 0).text = s.get("name", "")
            for ci, val in enumerate((s.get("values") or [])[:n_data_cols]):
                table.cell(ri + 1, ci + 1).text = str(val) if val is not None else "—"
        except Exception as e:
            logger.warning("Table row %d write failed: %s", ri, e)


def _get_fallback_font(font_name: str | None) -> str | None:
    return font_name


def _apply_run_font(run, saved_font: dict, *, bold_override: bool | None = None) -> None:
    if not saved_font:
        if bold_override is not None:
            run.font.bold = bold_override
        return
    font_name = _get_fallback_font(saved_font.get("name"))
    if font_name:
        run.font.name = font_name
    if saved_font.get("size"):
        run.font.size = saved_font["size"]
    if bold_override is not None:
        run.font.bold = bold_override
    elif saved_font.get("bold") is not None:
        run.font.bold = saved_font["bold"]
    
    color_type = saved_font.get("color_type")
    if color_type == 1 and saved_font.get("color"):
        try:
            run.font.color.rgb = saved_font["color"]
        except Exception:
            pass
    elif color_type == 2 and saved_font.get("color_theme") is not None:
        try:
            run.font.color.theme_color = saved_font["color_theme"]
            if saved_font.get("color_brightness") is not None:
                run.font.color.brightness = saved_font["color_brightness"]
        except Exception:
            pass
    elif saved_font.get("color"):
        # Fallback for general RGB color
        try:
            run.font.color.rgb = saved_font["color"]
        except Exception:
            pass


def _write_paragraph_text(paragraph, text: str, saved_font: dict, *, allow_bold: bool) -> None:
    """Write `text` into a cleared paragraph.

    If `allow_bold` is True and the text contains `**...**` segments, those
    segments are emitted as bold runs while the rest inherit the saved font.
    Otherwise any stray `**` markers are stripped and the text becomes one run.
    """
    if allow_bold and "**" in text:
        parts = re.split(r"\*\*([^*]+?)\*\*", text)
        # parts alternates: [plain, bold, plain, bold, ..., plain]
        first = True
        for idx, seg in enumerate(parts):
            if not seg:
                continue
            is_bold = (idx % 2 == 1)
            run = paragraph.add_run() if not first else paragraph.add_run()
            first = False
            run.text = seg
            _apply_run_font(run, saved_font, bold_override=True if is_bold else None)
        if first:  # nothing was written (all empty parts) — emit empty run
            run = paragraph.add_run()
            run.text = ""
            _apply_run_font(run, saved_font)
        return
    # Defensive: strip any stray markdown bold for non-bold fields.
    clean = text.replace("**", "").replace("__", "")
    run = paragraph.add_run()
    run.text = clean
    _apply_run_font(run, saved_font)


def _replace_text_in_frame(text_frame, replacements: dict) -> None:
    """Replace {{key}} tokens in a text frame, preserving run font properties.

    Special cases:
      - When the paragraph contains a placeholder whose key is in
        ``_INLINE_BOLD_KEYS`` (currently just ``investment_thesis_s1``), any
        ``**phrase**`` segments inside the substituted value are rendered as
        bold runs while surrounding text uses the original paragraph font.
      - For all other fields, stray ``**`` markers are stripped defensively.
      - Paragraph breaks inside the value (``\\n\\n``) are honoured by emitting
        additional paragraphs in the same text frame.
    """
    paragraphs = list(text_frame.paragraphs)
    for paragraph in paragraphs:
        full_text = paragraph.text
        present_keys = [k for k in replacements if f"{{{{{k}}}}}" in full_text]
        # Also detect single-brace SAARTHI score placeholders ({s_s} etc).
        single_brace_keys = [
            k for k in _SAARTHI_SCORE_PLACEHOLDERS
            if k in replacements and f"{{{k}}}" in full_text and f"{{{{{k}}}}}" not in full_text
        ]
        if not present_keys and not single_brace_keys:
            continue
        present_keys = present_keys + single_brace_keys

        # Capture font from first run before clearing
        saved_font: dict = {}
        if paragraph.runs:
            f = paragraph.runs[0].font
            saved_font["name"]  = f.name
            saved_font["size"]  = f.size
            saved_font["bold"]  = f.bold
            saved_font["color_type"] = getattr(f.color, "type", None)
            saved_font["color"] = None
            if saved_font["color_type"] == 1:
                try:
                    saved_font["color"] = f.color.rgb
                except Exception:
                    pass
            elif saved_font["color_type"] == 2:
                try:
                    saved_font["color_theme"] = f.color.theme_color
                    saved_font["color_brightness"] = f.color.brightness
                except Exception:
                    pass

        allow_bold = any(k in _INLINE_BOLD_KEYS for k in present_keys)

        for k, v in replacements.items():
            full_text = full_text.replace(f"{{{{{k}}}}}", str(v) if v is not None else "")
        # Single-brace pass for the small whitelisted set of SAARTHI score keys
        # (the template author chose `{s_s}` etc. instead of `{{s_s}}`).
        for k in _SAARTHI_SCORE_PLACEHOLDERS:
            v = replacements.get(k)
            if v is None:
                continue
            full_text = full_text.replace(f"{{{k}}}", str(v))

        # Honour explicit paragraph breaks (\n\n) and single newlines (bullets).
        # python-pptx paragraphs cannot contain newlines, so split and emit one
        # paragraph per line. Blank lines collapse into an empty paragraph.
        lines = full_text.split("\n")
        
        # Capture the original paragraph properties (pPr) element before clearing.
        # It contains bullet styles, indents, levels, and spacing.
        original_pPr = paragraph._p.pPr

        paragraph.clear()
        if not lines:
            return
        _write_paragraph_text(paragraph, lines[0], saved_font, allow_bold=allow_bold)
        from pptx.oxml.ns import qn  # local import to avoid top-level dep noise
        from pptx.text.text import _Paragraph
        from copy import deepcopy

        current_p = paragraph
        for extra in lines[1:]:
            # Create a sibling <a:p> right after the current paragraph element.
            # Copy original paragraph properties if present to keep formatting/bullets.
            new_p_el = current_p._p.makeelement(qn("a:p"), {})
            if original_pPr is not None:
                new_p_el.append(deepcopy(original_pPr))
            current_p._p.addnext(new_p_el)
            new_para = _Paragraph(new_p_el, text_frame)
            _write_paragraph_text(new_para, extra, saved_font, allow_bold=allow_bold)
            current_p = new_para  # so next iteration inserts after this one


def _apply_literal_text_subs(slide, subs: list[tuple[str, str]]) -> None:
    """Apply literal find→replace passes against every paragraph on a slide.

    Used for hard-coded strings baked into the master template (e.g. a stale
    company name in the disclosure clause) that the {{key}} replacer cannot
    address. Run-level edits are skipped because the literal may straddle
    multiple runs; we rewrite the paragraph text instead and keep the first
    run's font.
    """
    def _process_text_frame(tf) -> None:
        for paragraph in tf.paragraphs:
            full_text = paragraph.text
            if not any(old in full_text for old, _ in subs):
                continue
            new_text = full_text
            for old, new in subs:
                new_text = new_text.replace(old, new)
            if new_text == full_text:
                continue
            saved_font: dict = {}
            if paragraph.runs:
                f = paragraph.runs[0].font
                saved_font["name"] = f.name
                saved_font["size"] = f.size
                saved_font["bold"] = f.bold
            paragraph.text = new_text
            if paragraph.runs and saved_font:
                f = paragraph.runs[0].font
                if saved_font.get("name"):  f.name = saved_font["name"]
                if saved_font.get("size"):  f.size = saved_font["size"]
                if saved_font.get("bold") is not None: f.bold = saved_font["bold"]

    for shape in slide.shapes:
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if hasattr(cell, "text_frame") and cell.text_frame:
                        _process_text_frame(cell.text_frame)
        elif hasattr(shape, "text_frame") and shape.text_frame:
            _process_text_frame(shape.text_frame)


def _replace_text_in_table(shape, replacements: dict) -> None:
    """Replace {{key}} tokens inside every cell of a table shape."""
    for row in shape.table.rows:
        for cell in row.cells:
            if not hasattr(cell, "text_frame") or not cell.text_frame:
                continue
            for paragraph in cell.text_frame.paragraphs:
                full_text = paragraph.text
                if not any(f"{{{{{k}}}}}" in full_text for k in replacements):
                    continue
                for k, v in replacements.items():
                    full_text = full_text.replace(
                        f"{{{{{k}}}}}", str(v) if v is not None else "")
                paragraph.text = full_text


def _replace_slide4_thesis_shapes(slide, replacements: dict) -> None:
    """Apply slide-4-specific content sizing for the repeated thesis placeholders.

    The template reuses {{investment_thesis}} for one large left box and three
    small right boxes, so generic replacement makes the right boxes overflow.
    """
    headings = list(replacements.get("__slide4_right_headings") or [])
    right_texts = list(replacements.get("__slide4_right_texts") or [])
    thesis_count = 0
    heading_count = 0

    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or not shape.text_frame:
            continue
        text = shape.text_frame.text.strip()
        if text == "{{investment_thesis}}":
            thesis_count += 1
            if thesis_count == 1:
                _replace_text_in_frame(shape.text_frame, replacements)
            elif thesis_count <= 4:
                scoped = dict(replacements)
                scoped["investment_thesis"] = right_texts[thesis_count - 2] if thesis_count - 2 < len(right_texts) else ""
                _replace_text_in_frame(shape.text_frame, scoped)
        elif text == "{{investment_thesis_heading}}":
            heading_count += 1
            scoped = dict(replacements)
            scoped["investment_thesis_heading"] = headings[heading_count - 1] if heading_count - 1 < len(headings) else "Investment Thesis"
            _replace_text_in_frame(shape.text_frame, scoped)


def _replace_slide6_overview_shapes(slide, replacements: dict) -> None:
    overview_shapes = [
        shape
        for shape in slide.shapes
        if hasattr(shape, "text_frame")
        and shape.text_frame
        and "{{COMPANY_OVERVIEW}}" in shape.text_frame.text
    ]
    if not overview_shapes:
        return

    overview_shapes.sort(key=lambda s: (int(s.top), int(s.left)))
    top_text = replacements.get("__slide6_top_overview") or replacements.get("COMPANY_OVERVIEW", "")
    bottom_text = replacements.get("__slide6_bottom_overview") or replacements.get("COMPANY_OVERVIEW", "")

    if len(overview_shapes) >= 1:
        _replace_text_in_frame(overview_shapes[0].text_frame, {"COMPANY_OVERVIEW": top_text})
    if len(overview_shapes) >= 2:
        _replace_text_in_frame(overview_shapes[1].text_frame, {"COMPANY_OVERVIEW": bottom_text})


def _replace_slide11_management_shapes(slide, replacements: dict) -> None:
    heading_shapes = [
        shape for shape in slide.shapes
        if hasattr(shape, "text_frame") and shape.text_frame and "{{management_commentry_heading}}" in shape.text_frame.text
    ]
    content_shapes = [
        shape for shape in slide.shapes
        if hasattr(shape, "text_frame") and shape.text_frame and "{{management_content}}" in shape.text_frame.text
    ]
    heading_shapes.sort(key=lambda s: (int(s.top), int(s.left)))
    content_shapes.sort(key=lambda s: (int(s.top), int(s.left)))

    headings = replacements.get("__slide11_headings") or []
    contents = replacements.get("__slide11_contents") or []
    default_heading = replacements.get("management_commentry_heading", "Management Analysis")
    default_content = replacements.get("management_content", "")

    for idx, shape in enumerate(heading_shapes):
        heading = headings[idx] if idx < len(headings) else default_heading
        _replace_text_in_frame(shape.text_frame, {"management_commentry_heading": heading})
    for idx, shape in enumerate(content_shapes):
        content = contents[idx] if idx < len(contents) else default_content
        _replace_text_in_frame(shape.text_frame, {"management_content": content})


def _replace_slide12_indicator_shapes(slide, replacements: dict) -> None:
    indicator_shapes = [
        shape for shape in slide.shapes
        if hasattr(shape, "text_frame") and shape.text_frame and "{{indicators}}" in shape.text_frame.text
    ]
    indicator_shapes.sort(key=lambda s: (int(s.top), int(s.left)))
    cards = replacements.get("__slide12_indicator_cards") or []
    fallback = replacements.get("indicators", "")
    for idx, shape in enumerate(indicator_shapes):
        content = cards[idx] if idx < len(cards) else fallback
        _replace_text_in_frame(shape.text_frame, {"indicators": content})


def _last_actual_index(years: list[str]) -> int | None:
    actual_idx = [idx for idx, year in enumerate(years) if "E" not in str(year).upper()]
    return actual_idx[-1] if actual_idx else None


def _build_slide6_pie_data(fin_model: dict) -> tuple[dict[str, float], dict[str, float]]:
    operational = fin_model.get("operational") or {}
    years = [str(y).strip() for y in (operational.get("years") or [])]
    latest_actual_idx = _last_actual_index(years)

    revenue_mix_raw = operational.get("revenue_mix_pct") or {}
    revenue_mix = {
        str(k): float(v) * 100.0 if float(v) <= 1.0 else float(v)
        for k, v in revenue_mix_raw.items()
        if v not in (None, "")
    }

    volume_segments = operational.get("volume_segments") or {}
    ebit_mix: dict[str, float] = {}
    if latest_actual_idx is not None and volume_segments:
        total = 0.0
        raw_points: dict[str, float] = {}
        for raw_name, series in volume_segments.items():
            vals = [float(v) if v is not None else 0.0 for v in (series or [])]
            if latest_actual_idx >= len(vals):
                continue
            name = str(raw_name).replace(" Recycling", "").strip()
            value = vals[latest_actual_idx]
            raw_points[name] = value
            total += value
        if total > 0:
            ebit_mix = {name: (value / total) * 100.0 for name, value in raw_points.items()}

    if revenue_mix and "Others" not in ebit_mix:
        known = sum(ebit_mix.values())
        other_pct = max(0.0, 100.0 - known)
        if other_pct > 0.1:
            ebit_mix["Others"] = other_pct

    if not revenue_mix:
        revenue_mix = dict(ebit_mix)
    if not ebit_mix:
        ebit_mix = dict(revenue_mix)

    return revenue_mix, ebit_mix


def _render_pie_chart(title: str, data: dict[str, float]) -> bytes | None:
    if not data:
        return None
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception as exc:
        logger.warning("matplotlib unavailable for pie chart render: %s", exc)
        return None

    labels = list(data.keys())
    values = [max(0.0, float(v)) for v in data.values()]
    total = sum(values)
    if total <= 0:
        return None

    colors = ["#1F4690", "#FFA500", "#3A5BA0", "#7BC8A4", "#D9E2F3", "#F7C873"]
    fig, ax = plt.subplots(figsize=(3.5, 3.5), facecolor="white")
    wedges, texts, autotexts = ax.pie(
        values,
        labels=labels,
        autopct=lambda pct: f"{pct:.0f}%" if pct >= 4 else "",
        startangle=90,
        colors=colors[: len(labels)],
        textprops={"fontsize": 9, "color": "#173B73", "fontweight": "bold"},
        wedgeprops={"linewidth": 1, "edgecolor": "white"},
    )
    for auto in autotexts:
        auto.set_color("white")
        auto.set_fontsize(9)
        auto.set_fontweight("bold")
    ax.set_title(title, fontsize=11, color="#1F4690", fontweight="bold", pad=10)
    ax.axis("equal")
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=220, bbox_inches="tight", pad_inches=0.04, facecolor=fig.get_facecolor())
    plt.close(fig)
    return buf.getvalue()


def inject_company_overview_slide(pptx_path: str, fin_model: dict) -> int:
    revenue_mix, ebit_mix = _build_slide6_pie_data(fin_model)
    pie_images = {
        "{{pie_chart_1}}": _render_pie_chart("Revenue Mix %", revenue_mix),
        "{{pie_chart_2}}": _render_pie_chart("EBIT Mix %", ebit_mix),
    }
    pie_images = {token: img for token, img in pie_images.items() if img}
    if not pie_images:
        return 0

    prs = Presentation(pptx_path)
    if len(prs.slides) < 6:
        return 0
    slide = prs.slides[5]
    injected = 0
    fallback_targets = {
        "{{pie_chart_1}}": "Segment breakdown — see Excel model for details.",
        "{{pie_chart_2}}": "Segment breakdown — see Excel model for details.",
    }
    for token, img_bytes in pie_images.items():
        target_shape = None
        for shape in slide.shapes:
            if (
                hasattr(shape, "text_frame")
                and shape.text_frame
                and shape.text_frame.text.strip() in {token, fallback_targets.get(token, "")}
            ):
                target_shape = shape
                break
        if target_shape is None:
            continue
        _insert_image_into_shape(slide, target_shape, img_bytes)
        injected += 1
    if injected:
        prs.save(pptx_path)
    return injected


def _remove_shape(shape) -> None:
    try:
        sp_elem = shape._element
        sp_elem.getparent().remove(sp_elem)
    except Exception:
        pass


def _read_timeline_rows(
    excel_path: str,
    *,
    sheet_name: str = "Timeline",
) -> list[tuple[str, str, str, str]]:
    from openpyxl import load_workbook

    wb = load_workbook(excel_path, data_only=True)
    if sheet_name not in wb.sheetnames:
        return []
    ws = wb[sheet_name]
    rows: list[tuple[str, str, str, str]] = []
    for r in range(5, ws.max_row + 1):
        year = ws.cell(r, 1).value
        category = ws.cell(r, 2).value
        description = ws.cell(r, 3).value
        impact = ws.cell(r, 4).value
        if not year or not category or not description:
            continue
        desc_text = _truncate_words(str(description), 30)
        impact_text = _truncate_words(str(impact or ""), 25)
        rows.append((str(year), str(category), desc_text, impact_text))
    return rows


def _set_cell_left_align(cell, padding: float = 0.02) -> None:
    """Set cell text to left alignment and reposition the anchor so it does not overflow."""
    cell.get_text().set_ha("left")
    cell.get_text().set_x(padding)


def _wrap_table_cells(
    cell_rows: list[list[Any]],
    col_widths: list[float],
    *,
    fig_width_inches: float,
    chars_per_inch: float = 11.0,
    skip_header: bool = True,
) -> tuple[list[list[str]], list[int]]:
    """Pre-wrap every cell's text so it fits inside its column width.

    matplotlib's ``Table`` does not wrap text natively — overflow is clipped at
    the column edge, which is the "text getting cut" problem visible on the
    timeline / risks / governance tables. This helper inserts explicit ``\\n``
    line breaks using ``textwrap`` sized to each column's share of the figure
    width, and returns the wrapped rows alongside a per-row line count so the
    caller can scale cell heights via :func:`_apply_row_line_heights`.

    ``chars_per_inch`` is an empirical calibration for fontsize ~8-9pt on the
    standard serif/sans render — bump it down for larger fonts or up for
    tighter ones. ``skip_header`` keeps the first row at 1 line which usually
    looks better; pass ``False`` if a header label is itself too long.
    """
    import textwrap
    wrapped: list[list[str]] = []
    line_counts: list[int] = []
    full_chars = max(20.0, fig_width_inches * chars_per_inch)
    for ri, row in enumerate(cell_rows):
        new_row: list[str] = []
        max_lines = 1
        for ci, raw in enumerate(row):
            text = "" if raw is None else str(raw)
            col_w = col_widths[ci] if ci < len(col_widths) else 1.0 / max(1, len(row))
            width = max(8, int(full_chars * col_w))
            if "\n" in text:
                new_row.append(text)
                max_lines = max(max_lines, text.count("\n") + 1)
                continue
            if len(text) <= width:
                new_row.append(text)
                continue
            lines = textwrap.wrap(text, width=width, break_long_words=False, break_on_hyphens=True) or [text]
            new_row.append("\n".join(lines))
            if len(lines) > max_lines:
                max_lines = len(lines)
        wrapped.append(new_row)
        if skip_header and ri == 0:
            line_counts.append(1)
        else:
            line_counts.append(max_lines)
    return wrapped, line_counts


def _apply_row_line_heights(table, line_counts: list[int]) -> None:
    """Scale each cell's height proportional to its row's wrapped-line count.

    matplotlib ``Table`` cell heights are fractions of the axes that sum to
    ~1.0 by default (equal heights). We rebuild those fractions so rows with
    more wrapped lines receive more vertical space — eliminating the
    "one cramped row clipping multi-line text" look. We also rebuild and
    apply the vertical y-coordinates of the cells (via cell.set_y()) so that
    cells do not overlap or leave gaps when their heights are resized.
    We also tighten the cell's internal padding (default ``PAD=0.1`` reserves
    20% of cell height for top+bottom margin) so wrapped text actually fits
    inside its row, and centre-align vertically so the visual baseline of
    multi-line cells lines up with single-line cells in the same row.
    """
    total = sum(line_counts)
    if total <= 0:
        return

    # Calculate the bottom y-position for each row (top of table is y=1.0)
    y_positions = []
    current_y = 1.0
    for lines in line_counts:
        h = lines / total
        current_y -= h
        y_positions.append(current_y)

    cells = table.get_celld()
    for (r, _c), cell in cells.items():
        if 0 <= r < len(line_counts):
            h = line_counts[r] / total
            cell.set_height(h)
            cell.set_y(y_positions[r])
            cell.set_text_props(verticalalignment="center")
            # Differential padding by row type:
            #   - Multi-line wrapped rows (risks / timeline / governance):
            #     keep PAD tight so the wrapped text uses the full vertical
            #     area of the cell — was 0.03, still 0.03.
            #   - Single-line rows (numeric tables like financials, valuations,
            #     PE sensitivity): use a slightly looser PAD so numbers don't
            #     visually kiss the left / right cell borders. The default
            #     matplotlib PAD is 0.1 which leaves too much; 0.06 hits a
            #     comfortable middle.
            cell.PAD = 0.03 if line_counts[r] > 1 else 0.06


def _render_timeline_table(
    excel_path: str,
    *,
    sheet_name: str = "Timeline",
) -> bytes | None:
    rows = _read_timeline_rows(excel_path, sheet_name=sheet_name)
    if not rows:
        return None
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception as exc:
        logger.warning("matplotlib unavailable for timeline render: %s", exc)
        return None

    category_colors = {
        "Founding": "#8CBF88",
        "Expansion": "#5B9BD5",
        "IPO": "#9E9E9E",
        "International Expansion": "#7A7A7A",
        "Strategy": "#FFA500",
        "New Vertical": "#9C7AE3",
        "Milestone": "#31B0B0",
        "Outlook": "#2F5597",
    }

    cell_rows = [["Year", "Event Category", "Description", "Strategic Impact"]]
    cell_rows.extend([[y, c, d, i] for y, c, d, i in rows])
    styles = ["header"] + ["data"] * len(rows)

    fig_width = 11.2
    col_widths = [0.08, 0.16, 0.42, 0.34]
    cell_rows, line_counts = _wrap_table_cells(cell_rows, col_widths, fig_width_inches=fig_width)
    # Slide-7 placeholder aspect is ~11×5.6. python-pptx stretches the inserted
    # image to the placeholder bounds (ignores aspect ratio), so growing fig_h
    # past the placeholder height squashes the text vertically. Cap the growth
    # — the wrap helper already redistributes cell heights proportionally.
    fig_h = min(7.5, max(5.6, 0.16 * sum(line_counts) + 0.3))
    fig, ax = plt.subplots(figsize=(fig_width, fig_h), facecolor="white")
    ax.axis("off")

    table = ax.table(
        cellText=cell_rows,
        cellLoc="left",
        colWidths=col_widths,
        loc="upper left",
        bbox=[0, 0, 1, 1],
    )
    table.auto_set_font_size(False)
    # Dynamic font scaling: shrink if many rows/lines crowd the table
    total_lines = sum(line_counts)
    if total_lines <= 30:
        font_sz = 8.3
    elif total_lines <= 40:
        font_sz = 7.5
    elif total_lines <= 50:
        font_sz = 7.0
    else:
        font_sz = 6.5
    table.set_fontsize(font_sz)
    _apply_row_line_heights(table, line_counts)

    for (r, c), cell in table.get_celld().items():
        cell.set_edgecolor("#D5DCE8")
        cell.set_linewidth(0.6)
        if styles[r] == "header":
            cell.set_facecolor("#1F4690")
            cell.get_text().set_color("white")
            cell.get_text().set_fontweight("bold")
            cell.get_text().set_ha("center")
        else:
            if c == 0:
                cell.set_facecolor("#F7F7F7")
                cell.get_text().set_ha("center")
            elif c == 1:
                cat = rows[r - 1][1]
                cell.set_facecolor(category_colors.get(cat, "#B7C3D0"))
                cell.get_text().set_color("white")
                cell.get_text().set_fontweight("bold")
                cell.get_text().set_ha("center")
            else:
                cell.set_facecolor("white")
                cell.PAD = 0.01

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=220, bbox_inches="tight", pad_inches=0.02, facecolor=fig.get_facecolor())
    plt.close(fig)
    return buf.getvalue()


def _read_peer_compare_sections(excel_path: str) -> dict[str, dict]:
    from openpyxl import load_workbook

    wb = load_workbook(excel_path, data_only=True)
    ws = wb["Peer_Compare"]

    def read_block(title_row: int, header_row: int, data_start: int, data_end: int) -> dict:
        title = str(ws.cell(title_row, 1).value or "").strip()
        headers = [str(ws.cell(header_row, c).value or "").strip() for c in range(2, 7)]
        rows = []
        for r in range(data_start, data_end + 1):
            name = ws.cell(r, 1).value
            if not name:
                continue
            values = [ws.cell(r, c).value for c in range(2, 7)]
            rows.append({"company": str(name), "values": values})
        return {"title": title, "headers": headers, "rows": rows}

    valuation_headers = [str(ws.cell(26, c).value or "").strip() for c in range(2, 7)]
    valuation_rows = []
    for r in range(27, 32):
        name = ws.cell(r, 1).value
        if not name:
            continue
        valuation_rows.append({"company": str(name), "values": [ws.cell(r, c).value for c in range(2, 7)]})

    return {
        "revenue": read_block(4, 5, 6, 9),
        "ebitda_margin": read_block(11, 12, 13, 16),
        "pat": read_block(18, 19, 20, 23),
        "valuation": {"title": str(ws.cell(25, 1).value or "").strip(), "headers": valuation_headers, "rows": valuation_rows},
    }


def _fmt_peer_value(value, *, pct: bool = False) -> str:
    if value in (None, "", "-"):
        return "-"
    try:
        num = float(value)
    except Exception:
        return str(value)
    if pct:
        return f"{num * 100:.1f}%"
    if abs(num) >= 100:
        return f"{num:,.0f}"
    return f"{num:,.1f}"


def _render_peer_table(excel_path: str) -> bytes | None:
    sections = _read_peer_compare_sections(excel_path)
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception as exc:
        logger.warning("matplotlib unavailable for peer table render: %s", exc)
        return None

    rows: list[list[str]] = []
    row_styles: list[str] = []
    for key in ["revenue", "ebitda_margin", "pat", "valuation"]:
        section = sections[key]
        rows.append([section["title"], "", "", "", "", ""])
        row_styles.append("section")
        rows.append(["Company", *section["headers"]])
        row_styles.append("header")
        for item in section["rows"]:
            pct = key == "ebitda_margin" or key == "valuation" and False
            vals = [_fmt_peer_value(v, pct=pct) for v in item["values"]]
            if key == "valuation":
                vals = [
                    _fmt_peer_value(item["values"][0]),
                    _fmt_peer_value(item["values"][1]),
                    str(item["values"][2] if item["values"][2] not in (None, "") else "-"),
                    _fmt_peer_value(item["values"][3], pct=True),
                    _fmt_peer_value(item["values"][4], pct=True),
                ]
            rows.append([item["company"], *vals])
            row_styles.append("data")

    fig_width = 8.2
    col_widths = [0.31, 0.14, 0.14, 0.14, 0.14, 0.13]
    rows, line_counts = _wrap_table_cells(rows, col_widths, fig_width_inches=fig_width, chars_per_inch=13.0, skip_header=False)
    # Cap fig_h growth to avoid python-pptx stretching the image into the
    # fixed placeholder bounds (which would squash all text).
    fig_h = min(9.0, max(7.2, 0.13 * sum(line_counts) + 0.4))
    fig, ax = plt.subplots(figsize=(fig_width, fig_h), facecolor="white")
    ax.axis("off")

    table = ax.table(
        cellText=rows,
        cellLoc="center",
        colWidths=col_widths,
        loc="upper left",
        bbox=[0, 0, 1, 1],
    )
    table.auto_set_font_size(False)
    table.set_fontsize(7.0)
    _apply_row_line_heights(table, line_counts)

    for (r, c), cell in table.get_celld().items():
        cell.set_edgecolor("#D5DCE8")
        cell.set_linewidth(0.5)
        style = row_styles[r]
        if style == "section":
            cell.set_facecolor("#FFA500" if c == 0 else "#FFF4D6")
            if c == 0:
                cell.get_text().set_color("white")
                cell.get_text().set_fontweight("bold")
                _set_cell_left_align(cell)
            else:
                cell.get_text().set_text("")
        elif style == "header":
            cell.set_facecolor("#1F4690")
            cell.get_text().set_color("white")
            cell.get_text().set_fontweight("bold")
        else:
            if c == 0:
                cell.set_facecolor("#F4F7FC")
                _set_cell_left_align(cell)
                if rows[r][0].upper().startswith("GRAVITA"):
                    cell.get_text().set_fontweight("bold")
            else:
                cell.set_facecolor("white")
                if rows[r][0].upper().startswith("GRAVITA"):
                    cell.get_text().set_fontweight("bold")

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=220, bbox_inches="tight", pad_inches=0.02, facecolor=fig.get_facecolor())
    plt.close(fig)
    return buf.getvalue()


def _render_peer_bar_chart(title: str, names: list[str], values: list[float], *, percent: bool = False) -> bytes | None:
    if not names or not values:
        return None
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception as exc:
        logger.warning("matplotlib unavailable for peer bar chart render: %s", exc)
        return None

    colors = ["#1F4690"] + ["#FFA500"] * (len(names) - 1)
    fig, ax = plt.subplots(figsize=(5.1, 2.6), facecolor="white")
    bars = ax.barh(names, values, color=colors)
    ax.set_title(title, fontsize=11, color="#1F4690", fontweight="bold", pad=8)
    ax.grid(axis="x", linestyle="--", alpha=0.25)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_visible(False)
    for bar, val in zip(bars, values):
        label = f"{val:.1f}%" if percent else f"{val:,.0f}"
        ax.text(bar.get_width(), bar.get_y() + bar.get_height() / 2, f" {label}", va="center", fontsize=8.5)
    if percent:
        ax.set_xlim(0, max(values) * 1.25 if values else 1)
    ax.tick_params(axis="y", labelsize=8.5)
    ax.tick_params(axis="x", labelsize=8)
    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=220, bbox_inches="tight", pad_inches=0.04, facecolor=fig.get_facecolor())
    plt.close(fig)
    return buf.getvalue()


def inject_competitive_advantage_slide(pptx_path: str, *, excel_path: str | None = None) -> int:
    if not excel_path:
        return 0
    sections = _read_peer_compare_sections(excel_path)
    revenue_rows = sections["revenue"]["rows"]
    margin_rows = sections["ebitda_margin"]["rows"]
    revenue_chart = _render_peer_bar_chart(
        "Revenue FY26A (₹ Cr)",
        [r["company"].replace("GRAVITA INDIA LTD", "Gravita") for r in revenue_rows],
        [float(r["values"][-1]) for r in revenue_rows],
        percent=False,
    )
    margin_chart = _render_peer_bar_chart(
        "EBITDA Margin FY26A",
        [r["company"].replace("GRAVITA INDIA LTD", "Gravita") for r in margin_rows],
        [float(r["values"][-1]) * 100 for r in margin_rows],
        percent=True,
    )
    if not revenue_chart and not margin_chart:
        return 0

    prs = Presentation(pptx_path)
    if len(prs.slides) < 9:
        return 0
    slide = prs.slides[8]
    target_shapes = []
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            text = shape.text_frame.text.strip()
            if text in {
                "{{competitive_chart_1}}",
                "{{competitive_chart_2}}",
                "Competitive positioning chart — see Excel model for details.",
            }:
                target_shapes.append(shape)
    if not target_shapes:
        return 0
    target_shapes.sort(key=lambda s: (int(s.top), int(s.left)))
    grouped: list[list] = []
    for shape in target_shapes:
        if not grouped or abs(int(shape.top) - int(grouped[-1][0].top)) > 250000:
            grouped.append([shape])
        else:
            grouped[-1].append(shape)

    injected = 0
    chart_imgs = [revenue_chart, margin_chart]
    for grp, img in zip(grouped[:2], chart_imgs):
        if not img:
            continue
        anchor = grp[0]
        for extra in grp[1:]:
            _remove_shape(extra)
        _insert_image_into_shape(slide, anchor, img)
        injected += 1
    if injected:
        prs.save(pptx_path)
    return injected


def inject_peer_comparison_slide(pptx_path: str, *, excel_path: str | None = None) -> int:
    if not excel_path:
        return 0
    img = _render_peer_table(excel_path)
    if not img:
        return 0
    prs = Presentation(pptx_path)
    if len(prs.slides) < 10:
        return 0
    slide = prs.slides[9]
    target = None
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            text = shape.text_frame.text.strip()
            if text in {"{{peer_comparision}}", "Peer comparison — see Excel model for details."}:
                target = shape
                break
    if target is None:
        return 0
    _insert_image_into_shape(slide, target, img)
    prs.save(pptx_path)
    return 1


def inject_company_overview_slide(pptx_path: str, fin_model: dict) -> int:
    revenue_mix, ebit_mix = _build_slide6_pie_data(fin_model)
    pie_images = {
        "{{percentage_revenue_pie_chart}}": _render_pie_chart("Revenue Mix %", revenue_mix),
        "{{percentage_EBIT_pie_chart}}": _render_pie_chart("EBIT Mix %", ebit_mix),
    }
    pie_images = {token: img for token, img in pie_images.items() if img}
    if not pie_images:
        return 0

    prs = Presentation(pptx_path)
    if len(prs.slides) < 6:
        return 0
    slide = prs.slides[5]
    injected = 0
    for token, img_bytes in pie_images.items():
        target_shape = None
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame and shape.text_frame.text.strip() == token:
                target_shape = shape
                break
        if target_shape is None:
            continue
        # preserve_aspect: pie shapes render as ovals if stretched to the
        # template's non-square placeholder. Letterbox-fit keeps them circular.
        _insert_image_into_shape(slide, target_shape, img_bytes, preserve_aspect=True)
        injected += 1
    if injected:
        prs.save(pptx_path)
    return injected


def inject_company_timeline_slide(pptx_path: str, *, excel_path: str | None = None) -> int:
    """Render the company timeline table image into the slide-7 placeholder.

    Idempotent: returns 0 (no-op) if the {{company_timeline}} placeholder is
    not present — which is the expected state after the first successful
    injection. The "Re-inject after cleanup" pass calls this a second time;
    without this guard, the earlier fallback ("largest text shape") was
    stamping a tiny duplicate timeline image onto an unrelated text box at
    the bottom-right of slide 7.
    """
    if not excel_path:
        return 0
    img = _render_timeline_table(excel_path)
    if not img:
        return 0
    prs = Presentation(pptx_path)
    if len(prs.slides) < 7:
        return 0
    slide = prs.slides[6]
    target = None
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            text = shape.text_frame.text.strip()
            if text in {"{{company_timeline}}", "{{COMPANY_TIMELINE}}", "Timeline — see Excel model for details."}:
                target = shape
                break
    if target is None:
        return 0
    _insert_image_into_shape(slide, target, img)
    prs.save(pptx_path)
    return 1


def inject_catalyst_timeline_slide(pptx_path: str, *, excel_path: str | None = None) -> int:
    if not excel_path:
        return 0
    img = _render_timeline_table(excel_path, sheet_name="Catalyst_Timeline")
    if not img:
        return 0
    prs = Presentation(pptx_path)
    if len(prs.slides) < 10:
        return 0
    slide = prs.slides[9]
    target = None
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            text = shape.text_frame.text.strip()
            if text in {"{{catalyst_timeline_chart}}", "Catalyst timeline — see Excel model for details."}:
                target = shape
                break
    if target is None:
        return 0
    _insert_image_into_shape(slide, target, img)
    prs.save(pptx_path)
    return 1


def inject_competitive_advantage_slide(pptx_path: str, *, excel_path: str | None = None) -> int:
    if not excel_path:
        return 0
    sections = _read_peer_compare_sections(excel_path)
    revenue_rows = sections["revenue"]["rows"]
    margin_rows = sections["ebitda_margin"]["rows"]
    revenue_chart = _render_peer_bar_chart(
        "Revenue FY26A (₹ Cr)",
        [r["company"].replace("GRAVITA INDIA LTD", "Gravita") for r in revenue_rows],
        [float(r["values"][-1]) for r in revenue_rows],
        percent=False,
    )
    margin_chart = _render_peer_bar_chart(
        "EBITDA Margin FY26A",
        [r["company"].replace("GRAVITA INDIA LTD", "Gravita") for r in margin_rows],
        [float(r["values"][-1]) * 100 for r in margin_rows],
        percent=True,
    )
    prs = Presentation(pptx_path)
    if len(prs.slides) < 9:
        return 0
    slide = prs.slides[8]
    mapping = {
        "{{peer_comparison_chart_1}}": revenue_chart,
        "{{peer_comparison_chart_2}}": margin_chart,
    }
    injected = 0
    for token, img in mapping.items():
        if not img:
            continue
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame and shape.text_frame.text.strip() == token:
                _insert_image_into_shape(slide, shape, img)
                injected += 1
                break
    if injected:
        prs.save(pptx_path)
    return injected


def _render_named_sheet_table(excel_path: str, sheet_name: str, *, max_rows: int = 60) -> bytes | None:
    from openpyxl import load_workbook

    try:
        wb = load_workbook(excel_path, data_only=True)
    except Exception as exc:
        logger.warning("Failed to open workbook for sheet render %s: %s", sheet_name, exc)
        return None
    if sheet_name not in wb.sheetnames:
        return None
    try:
        return excel_injector.render_sheet_as_image(wb[sheet_name], title="", max_rows=max_rows)
    except Exception as exc:
        logger.warning("Failed to render sheet %s: %s", sheet_name, exc)
        return None


def _render_governance_table(excel_path: str) -> bytes | None:
    from openpyxl import load_workbook
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception as exc:
        logger.warning("matplotlib unavailable for governance render: %s", exc)
        return None

    wb = load_workbook(excel_path, data_only=True)
    ws = wb["Governance"]
    board = [[ws.cell(r, c).value for c in range(1, 6)] for r in range(5, 12)]
    shareholding = [[ws.cell(r, c).value for c in range(1, 6)] for r in range(14, 19)]

    fig = plt.figure(figsize=(6.0, 6.8), facecolor="white")
    gs = fig.add_gridspec(2, 1, height_ratios=[1.65, 1.0], hspace=0.16)
    axes = [fig.add_subplot(gs[0]), fig.add_subplot(gs[1])]
    for ax in axes:
        ax.axis("off")

    sections = [
        ("BOARD OF DIRECTORS", board, [0.28, 0.30, 0.20, 0.12, 0.10]),
        ("SHAREHOLDING PATTERN %", shareholding, [0.32, 0.17, 0.17, 0.17, 0.17]),
    ]
    for ax, (title, rows, widths) in zip(axes, sections):
        ax.text(0, 1.03, title, transform=ax.transAxes, fontsize=10, color="white", fontweight="bold",
                bbox=dict(facecolor="#FFA500", edgecolor="#FFA500", boxstyle="square,pad=0.25"))
        formatted = []
        for ridx, row in enumerate(rows):
            out = []
            for cidx, val in enumerate(row):
                if ridx > 0 and cidx > 0 and isinstance(val, (int, float)) and cidx < len(row):
                    if title.startswith("SHARE"):
                        out.append(f"{float(val) * 100:.1f}%")
                    else:
                        out.append(str(int(val)) if float(val).is_integer() else f"{float(val):.1f}")
                else:
                    out.append("" if val is None else str(val))
            formatted.append(out)
        # fontsize 8.1 renders at ~10 chars/inch (proportional font). The
        # default chars_per_inch=11 was over-counting per-line capacity, so
        # names like "Mahavir Prasad Agarwal" stayed unwrapped and clipped at
        # the column edge. 10 is conservative — wraps anything close to width.
        formatted, line_counts = _wrap_table_cells(formatted, widths, fig_width_inches=6.0, chars_per_inch=10.0, skip_header=False)
        table = ax.table(cellText=formatted, cellLoc="left", colWidths=widths, loc="upper left", bbox=[0, 0, 1, 0.95])
        table.auto_set_font_size(False)
        table.set_fontsize(8.1)
        _apply_row_line_heights(table, line_counts)
        for (r, c), cell in table.get_celld().items():
            cell.set_edgecolor("#D5DCE8")
            cell.set_linewidth(0.5)
            if r == 0:
                cell.set_facecolor("#1F4690")
                cell.get_text().set_color("white")
                cell.get_text().set_fontweight("bold")
                cell.get_text().set_ha("center")
                cell.get_text().set_x(0.5)
            else:
                cell.set_facecolor("white")
                if c == 0:
                    _set_cell_left_align(cell)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=220, bbox_inches="tight", pad_inches=0.02, facecolor=fig.get_facecolor())
    plt.close(fig)
    return buf.getvalue()


def _render_key_risks_table(excel_path: str) -> bytes | None:
    from openpyxl import load_workbook
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception as exc:
        logger.warning("matplotlib unavailable for key-risks render: %s", exc)
        return None

    wb = load_workbook(excel_path, data_only=True)
    ws = wb["Key_Risks"]
    rows = [[ws.cell(r, c).value for c in range(1, 9)] for r in range(4, 13)]
    formatted = [["" if v is None else str(v) for v in row] for row in rows]
    fig_width = 11.2
    # Column widths: '#' is narrow. Risk Category/Factor at 0.12/0.13 hold
    # 2-3 word phrases comfortably. Description/Mitigation get the lion's
    # share. Probability/Impact are single-letter chips. Overall Rating only
    # needs to hold "MEDIUM" (~6 chars) so 0.095 is plenty — we steal the
    # remainder back for Description/Mitigation which carry the long text.
    # Sums to 1.0 so the table fills the figure width exactly.
    col_widths = [0.04, 0.12, 0.13, 0.25, 0.225, 0.07, 0.07, 0.095]
    # Bumped chars_per_inch (12 → 14) so Description/Mitigation pack more
    # chars per line and produce fewer wrapped rows — the previous run was
    # overflowing because there were too many lines for the placeholder.
    formatted, line_counts = _wrap_table_cells(formatted, col_widths, fig_width_inches=fig_width, chars_per_inch=14.0, skip_header=False)
    # Slide-18 placeholder is ~11×5.4. python-pptx stretches the image to the
    # placeholder bounds, so a too-tall figure squashes text vertically. We
    # allow a modest extra ~30% of height for breathing room — combined with
    # the wider columns + slight fontsize trim below, this prevents the
    # cell-overflow we saw previously.
    fig_h = min(7.0, max(5.6, 0.16 * sum(line_counts) + 0.5))
    fig, ax = plt.subplots(figsize=(fig_width, fig_h), facecolor="white")
    ax.axis("off")
    table = ax.table(
        cellText=formatted,
        cellLoc="left",
        colWidths=col_widths,
        loc="upper left",
        bbox=[0, 0, 1, 1],
    )
    table.auto_set_font_size(False)
    # Auto-adjust font size based on total line count to prevent vertical overlapping/overflow
    total_lines = sum(line_counts)
    if total_lines > 50:
        fontsize = 5.2
    elif total_lines > 40:
        fontsize = 6.0
    elif total_lines > 30:
        fontsize = 6.6
    else:
        fontsize = 7.2
    table.set_fontsize(fontsize)
    _apply_row_line_heights(table, line_counts)
    color_map = {"H": "#D00000", "M": "#F5A623", "L": "#0B7D20", "MEDIUM": "#F5A623", "LOW": "#0B7D20", "HIGH": "#D00000"}
    for (r, c), cell in table.get_celld().items():
        cell.set_edgecolor("#D5DCE8")
        cell.set_linewidth(0.5)
        if r == 0:
            cell.set_facecolor("#1F4690")
            cell.get_text().set_color("white")
            cell.get_text().set_fontweight("bold")
            cell.get_text().set_ha("center")
        else:
            txt = formatted[r][c]
            if c in {5, 6, 7}:
                bg = color_map.get(txt, "#FFFFFF")
                cell.set_facecolor(bg)
                cell.get_text().set_color("white")
                cell.get_text().set_fontweight("bold")
                cell.get_text().set_ha("center")
            elif c == 0:
                cell.set_facecolor("white")
                cell.get_text().set_ha("center")
            else:
                cell.set_facecolor("white")
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=220, bbox_inches="tight", pad_inches=0.02, facecolor=fig.get_facecolor())
    plt.close(fig)
    return buf.getvalue()


def _render_formula_sheet_table(
    excel_path: str,
    sheet_name: str,
    *,
    max_row: int,
    max_col: int,
) -> bytes | None:
    from openpyxl import load_workbook
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception as exc:
        logger.warning("matplotlib unavailable for formula sheet render %s: %s", sheet_name, exc)
        return None

    # Earnings_Forecast and Valuations_Table are formula-only sheets — every
    # data cell is `='Financials_Table'!XN` or a derived expression. openpyxl
    # cannot compute formulas itself, so loading with data_only=True returns
    # None for unsaved formula cells (the workbook was never opened in Excel).
    # We must evaluate via our local formula resolver for these sheets.
    use_formula_eval = sheet_name in {"Financials_Table", "Earnings_Forecast", "Valuations_Table"}
    wb = load_workbook(excel_path, data_only=not use_formula_eval)
    ws = wb[sheet_name]
    cache: dict[tuple[str, str], float] = {}

    def _is_pct_label(label: str) -> bool:
        ll = label.lower()
        return ("%" in label) or ("margin" in ll) or ("yield" in ll) or ("rate" in ll) or ("growth" in ll) or ("roe" in ll) or ("roce" in ll)

    def _fmt_pct(val: float) -> str:
        # Excel stores percentages as fractions (0.10 = 10%). But some inputs
        # already arrive in percentage scale (15.0 meaning 15%). If the
        # magnitude is plainly > 1.5, the value is already in % units.
        if abs(val) > 1.5:
            return f"{val:.1f}%"
        return f"{val * 100:.1f}%"

    def display_value(r: int, c: int) -> str:
        cell = ws.cell(r, c)
        raw = cell.value
        if isinstance(raw, str) and raw.strip() == "-":
            return "-"
        label = str(ws.cell(r, 1).value or "")
        if c == 1:
            return label
        if use_formula_eval:
            val = _evaluate_excel_formula_cell(wb, sheet_name, cell.coordinate, cache)
            if val is None:
                return ""
            if _is_pct_label(label):
                return _fmt_pct(float(val))
            if "(x)" in label.lower():
                return f"{float(val):.1f}x"
            if abs(float(val)) >= 100:
                return f"{float(val):,.0f}"
            return f"{float(val):,.2f}".rstrip("0").rstrip(".")
        if raw in (None, ""):
            return ""
        if isinstance(raw, str):
            return raw
        try:
            val = float(raw)
        except (TypeError, ValueError):
            return str(raw)
        if _is_pct_label(label):
            return _fmt_pct(float(val))
        if "(x)" in label.lower():
            return f"{float(val):.1f}x"
        if abs(float(val)) >= 100:
            return f"{float(val):,.0f}"
        return f"{float(val):,.2f}".rstrip("0").rstrip(".")

    rows: list[list[str]] = []
    styles: list[str] = []
    for r in range(4, max_row + 1):
        first = ws.cell(r, 1).value
        if first is None:
            continue
        first_text = str(first).strip()
        if r == 4:
            row = [str(ws.cell(r, c).value or "") for c in range(1, max_col + 1)]
            styles.append("header")
        elif all((ws.cell(r, c).value in (None, "")) for c in range(2, max_col + 1)):
            row = [display_value(r, c) for c in range(1, max_col + 1)]
            styles.append("section")
        else:
            row = [display_value(r, c) for c in range(1, max_col + 1)]
            styles.append("data")
        rows.append(row)

    fig_width = 11.2
    # Slightly wider Particulars column (0.18 → 0.20) to host long section
    # labels like "PE SENSITIVITY (Target Price)" — they were overflowing the
    # cell on slide 15. Remaining columns shrink a hair from 0.082 → 0.080.
    col_widths = [0.20] + [0.080] * (max_col - 1)
    # fontsize 7.6 renders at ~12 chars/inch. Use 10 for a safety margin so
    # long labels in the leftmost column wrap before they hit the cell edge.
    rows, line_counts = _wrap_table_cells(rows, col_widths, fig_width_inches=fig_width, chars_per_inch=10.0)
    # Cap fig_h growth to keep aspect close to the slide placeholder.
    fig_h = min(7.5, max(5.5, 0.12 * sum(line_counts) + 0.4))
    fig, ax = plt.subplots(figsize=(fig_width, fig_h), facecolor="white")
    ax.axis("off")
    table = ax.table(cellText=rows, cellLoc="center", colWidths=col_widths, loc="upper left", bbox=[0, 0, 1, 1])
    table.auto_set_font_size(False)
    table.set_fontsize(7.6)
    _apply_row_line_heights(table, line_counts)

    for (r, c), cell in table.get_celld().items():
        cell.set_edgecolor("#D5DCE8")
        cell.set_linewidth(0.45)
        style = styles[r]
        if style == "header":
            cell.set_facecolor("#1F4690" if c < 6 else "#FFA500")
            cell.get_text().set_color("white")
            cell.get_text().set_fontweight("bold")
        elif style == "section":
            if c == 0:
                cell.set_facecolor("#FFA500")
                cell.get_text().set_color("white")
                cell.get_text().set_fontweight("bold")
                _set_cell_left_align(cell)
            else:
                cell.set_facecolor("#FFFFFF")
                cell.get_text().set_text("")
        else:
            if c == 0:
                cell.set_facecolor("#FFFFFF")
                _set_cell_left_align(cell)
                cell.PAD = 0.01
            else:
                cell.set_facecolor("#F7FAFF" if c < 6 else "#FFF1D9")

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=220, bbox_inches="tight", pad_inches=0.02, facecolor=fig.get_facecolor())
    plt.close(fig)
    return buf.getvalue()


def inject_sheet_visual_slide(
    pptx_path: str,
    *,
    excel_path: str | None,
    slide_number: int,
    sheet_name: str,
    tokens: list[str],
    fallback_texts: list[str],
    max_rows: int = 60,
) -> int:
    if not excel_path:
        return 0
    img = _render_named_sheet_table(excel_path, sheet_name, max_rows=max_rows)
    if not img:
        return 0

    prs = Presentation(pptx_path)
    if len(prs.slides) < slide_number:
        return 0
    slide = prs.slides[slide_number - 1]
    target = None
    match_texts = set(tokens) | set(fallback_texts)
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            text = shape.text_frame.text.strip()
            if text in match_texts:
                target = shape
                break
    if target is None:
        return 0
    _insert_image_into_shape(slide, target, img)
    prs.save(pptx_path)
    return 1


def inject_governance_slide(pptx_path: str, *, excel_path: str | None = None) -> int:
    if not excel_path:
        return 0
    img = _render_governance_table(excel_path)
    if not img:
        return 0
    prs = Presentation(pptx_path)
    if len(prs.slides) < 12:
        return 0
    slide = prs.slides[11]
    target = None
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            text = shape.text_frame.text.strip()
            if text in {"{{governance_table}}", "Governance — see Excel model for details."}:
                target = shape
                break
    if target is None:
        return 0
    _insert_image_into_shape(slide, target, img)
    prs.save(pptx_path)
    return 1


def inject_formula_table_slide(
    pptx_path: str,
    *,
    excel_path: str | None,
    slide_number: int,
    sheet_name: str,
    token: str,
    fallback_text: str,
    max_row: int,
    max_col: int,
) -> int:
    if not excel_path:
        return 0
    img = _render_formula_sheet_table(excel_path, sheet_name, max_row=max_row, max_col=max_col)
    if not img:
        return 0
    prs = Presentation(pptx_path)
    if len(prs.slides) < slide_number:
        return 0
    slide = prs.slides[slide_number - 1]
    target = None
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            text = shape.text_frame.text.strip()
            if text in {token, fallback_text}:
                target = shape
                break
    if target is None:
        return 0
    _insert_image_into_shape(slide, target, img)
    prs.save(pptx_path)
    return 1


def inject_key_risks_slide(pptx_path: str, *, excel_path: str | None = None) -> int:
    if not excel_path:
        return 0
    img = _render_key_risks_table(excel_path)
    if not img:
        return 0
    prs = Presentation(pptx_path)
    if len(prs.slides) < 18:
        return 0
    slide = prs.slides[17]
    target = None
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            text = shape.text_frame.text.strip()
            if text in {"{{key_risks_table}}", "Key risks — see Excel model for details."}:
                target = shape
                break
    if target is None:
        return 0
    _insert_image_into_shape(slide, target, img)
    prs.save(pptx_path)
    return 1


def _render_probability_weight_table(fin_model: dict) -> bytes | None:
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception as exc:
        logger.warning("matplotlib unavailable for probability table render: %s", exc)
        return None

    scenarios = fin_model.get("scenarios") or []
    rows: list[list[str]] = [["Scenario", "Target Price", "Probability", "Weighted TP"]]
    weighted_total = 0.0
    added = 0

    for scenario in scenarios:
        name = str(scenario.get("name", "")).strip().title()
        tp = _parse_number(scenario.get("target_price") or "") or 0.0
        prob = _parse_number(scenario.get("probability_pct") or "") or 0.0
        weighted = round(tp * prob / 100.0, 1)
        weighted_total += weighted
        rows.append([name or "-", f"{tp:.1f}", f"{prob:.0f}%", f"{weighted:.1f}"])
        added += 1

    if not added:
        return None

    rows.append(["Total", "", "100%", f"{weighted_total:.1f}"])

    fig, ax = plt.subplots(figsize=(11.2, 2.1), facecolor="white")
    ax.axis("off")
    table = ax.table(
        cellText=rows,
        cellLoc="center",
        colWidths=[0.30, 0.22, 0.20, 0.28],
        loc="upper left",
        bbox=[0, 0, 1, 1],
    )
    table.auto_set_font_size(False)
    table.set_fontsize(10.5)
    table.scale(1, 1.35)

    for (r, c), cell in table.get_celld().items():
        cell.set_edgecolor("#D5DCE8")
        cell.set_linewidth(0.8)
        if r == 0:
            cell.set_facecolor("#1F4690")
            cell.get_text().set_color("white")
            cell.get_text().set_fontweight("bold")
        elif r == len(rows) - 1:
            cell.set_facecolor("#E8EEF9")
            cell.get_text().set_fontweight("bold")
        else:
            cell.set_facecolor("white")

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=220, bbox_inches="tight", pad_inches=0.02, facecolor=fig.get_facecolor())
    plt.close(fig)
    return buf.getvalue()


def inject_probability_weight_slide(pptx_path: str, fin_model: dict) -> int:
    img = _render_probability_weight_table(fin_model)
    if not img:
        return 0
    prs = Presentation(pptx_path)
    if len(prs.slides) < 17:
        return 0
    slide = prs.slides[16]
    target = None
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            text = shape.text_frame.text.strip()
            if text in {
                "{{probability_weight_table}}",
                "Probability-weighted scenario analysis — see Excel model for details.",
            }:
                target = shape
                break
    if target is None:
        return 0
    _insert_image_into_shape(slide, target, img)
    prs.save(pptx_path)
    return 1


def _insert_image_into_shape(slide, shape, img_bytes: bytes, *, preserve_aspect: bool = False) -> None:
    """Replace a placeholder shape with a generated image at the same bounds.

    By default the picture is stretched to fill the placeholder (matches the
    template designer's intent for charts/tables sized to fit their box).

    When ``preserve_aspect=True`` we measure the actual image dimensions and
    fit it inside the placeholder while keeping its native aspect ratio,
    centring it within the placeholder bounds. Used for pie charts whose
    placeholders aren't square — without this, the circle gets stretched to
    fill the rectangle and renders as an oval.
    """
    import io as _io
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    sp_elem = shape._element
    sp_elem.getparent().remove(sp_elem)

    new_left, new_top, new_width, new_height = left, top, width, height
    if preserve_aspect:
        try:
            from PIL import Image as _Image
            img = _Image.open(_io.BytesIO(img_bytes))
            iw, ih = img.size
            if iw > 0 and ih > 0 and width > 0 and height > 0:
                img_aspect = iw / ih
                box_aspect = width / height
                if img_aspect > box_aspect:
                    # Image is wider than box — fit width, shrink height, centre vertically.
                    new_width = width
                    new_height = int(width / img_aspect)
                    new_top = top + (height - new_height) // 2
                else:
                    # Image is taller — fit height, shrink width, centre horizontally.
                    new_height = height
                    new_width = int(height * img_aspect)
                    new_left = left + (width - new_width) // 2
        except Exception as exc:
            logger.debug("preserve_aspect fallback (PIL unavailable or read failed): %s", exc)

    slide.shapes.add_picture(_io.BytesIO(img_bytes), new_left, new_top, new_width, new_height)


def _shape_is_excel_placeholder(shape) -> bool:
    """Return True if the shape's entire text matches an Excel injection token.

    These shapes must be preserved intact so the COM-based excel_injector
    can locate them after fill_master_template() saves the PPTX.
    """
    if not hasattr(shape, "text_frame") or not shape.text_frame:
        return False
    try:
        full_text = shape.text_frame.text.strip()
    except Exception:
        return False
    return full_text in _EXCEL_INJECTION_TOKENS


def _series_values(fin_model: dict, name: str) -> tuple[list[str], list[float]]:
    wanted = name.lower().strip()
    series_list = fin_model.get("series") or []
    for mode in ("exact", "contains"):
        for series in series_list:
            series_name = str(series.get("name", "")).lower().strip()
            if (mode == "exact" and series_name == wanted) or (mode == "contains" and wanted in series_name):
                periods = [str(p) for p in (series.get("periods") or [])]
                values: list[float] = []
                for raw in (series.get("values") or []):
                    try:
                        values.append(float(str(raw).replace(",", "")))
                    except (TypeError, ValueError):
                        values.append(0.0)
                return periods, values
    return [], []


def _format_actual_year_labels(periods: list[str]) -> list[str]:
    out: list[str] = []
    for period in periods:
        clean = str(period).strip()
        if clean.endswith("A"):
            out.append(clean)
        elif "E" not in clean.upper():
            out.append(f"{clean}A" if clean.startswith("FY") else clean)
        else:
            out.append(clean)
    return out


def _parse_excel_ref(ref: str) -> tuple[str | None, str]:
    if "!" in ref:
        sheet, cell = ref.split("!", 1)
        return sheet.strip().strip("'"), cell.strip()
    return None, ref.strip()


def _evaluate_excel_formula_cell(wb, ws_name: str, cell_ref: str, cache: dict[tuple[str, str], float]) -> float | None:
    key = (ws_name, cell_ref)
    if key in cache:
        return cache[key]

    ws = wb[ws_name]
    value = ws[cell_ref].value
    if isinstance(value, (int, float)):
        cache[key] = float(value)
        return cache[key]
    if value is None:
        cache[key] = 0.0
        return 0.0
    if not isinstance(value, str) or not value.startswith("="):
        try:
            cache[key] = float(str(value).replace(",", ""))
            return cache[key]
        except (TypeError, ValueError):
            cache[key] = 0.0
            return 0.0

    expr = value[1:].strip()
    ref_re = re.compile(r"(?:'[^']+'|[A-Za-z0-9_ ]+)?!\$?[A-Z]{1,3}\$?\d+|\$?[A-Z]{1,3}\$?\d+")

    def _split_excel_args(arg_str: str) -> list[str]:
        args: list[str] = []
        depth = 0
        start = 0
        for idx, ch in enumerate(arg_str):
            if ch == "(":
                depth += 1
            elif ch == ")":
                depth -= 1
            elif ch == "," and depth == 0:
                args.append(arg_str[start:idx].strip())
                start = idx + 1
        args.append(arg_str[start:].strip())
        return [arg for arg in args if arg]

    def _range_sum(range_expr: str) -> float:
        start_ref, end_ref = [part.strip() for part in range_expr.split(":", 1)]
        start_sheet, start_cell = _parse_excel_ref(start_ref.replace("$", ""))
        end_sheet, end_cell = _parse_excel_ref(end_ref.replace("$", ""))
        target_sheet = start_sheet or ws_name
        end_target_sheet = end_sheet or target_sheet
        if end_target_sheet != target_sheet:
            return 0.0
        start_col, start_row = coordinate_from_string(start_cell)
        end_col, end_row = coordinate_from_string(end_cell)
        start_col_idx = column_index_from_string(start_col)
        end_col_idx = column_index_from_string(end_col)
        total = 0.0
        for row_idx in range(start_row, end_row + 1):
            for col_idx in range(start_col_idx, end_col_idx + 1):
                coord = f"{get_column_letter(col_idx)}{row_idx}"
                total += _evaluate_excel_formula_cell(wb, target_sheet, coord, cache) or 0.0
        return total

    def repl(match: re.Match[str]) -> str:
        token = match.group(0).replace("$", "")
        sheet_name, inner_ref = _parse_excel_ref(token)
        target_sheet = sheet_name or ws_name
        inner_val = _evaluate_excel_formula_cell(wb, target_sheet, inner_ref, cache)
        return str(inner_val or 0.0)

    if expr.upper().startswith("SUM(") and expr.endswith(")"):
        inner = expr[4:-1]
        total = 0.0
        for part in _split_excel_args(inner):
            if ":" in part:
                total += _range_sum(part)
            else:
                token = part.replace("$", "")
                sheet_name, inner_ref = _parse_excel_ref(token)
                total += _evaluate_excel_formula_cell(wb, sheet_name or ws_name, inner_ref, cache) or 0.0
        cache[key] = float(total)
        return cache[key]

    # IFERROR(formula, fallback) — return formula's evaluated value, else fallback.
    # The financial model uses IFERROR liberally on division-style formulas so
    # this must be handled or every margin / multiple cell silently returns 0.
    if expr.upper().startswith("IFERROR(") and expr.endswith(")"):
        args = _split_excel_args(expr[8:-1])
        if len(args) == 2:
            inner_expr, fallback_expr = args
            try:
                inner_eval = ref_re.sub(repl, inner_expr)
                cache[key] = float(eval(inner_eval, {"__builtins__": {}}, {}))
                return cache[key]
            except Exception:
                # Fallback: try numeric coerce; if literal like "-" or "n/a", return 0.
                fallback_clean = fallback_expr.strip().strip('"').strip("'")
                try:
                    cache[key] = float(fallback_clean)
                except (TypeError, ValueError):
                    cache[key] = 0.0
                return cache[key]

    if expr.upper().startswith("IF(") and expr.endswith(")"):
        args = _split_excel_args(expr[3:-1])
        if len(args) == 3:
            cond_expr, true_expr, false_expr = args
            cond_eval = ref_re.sub(repl, cond_expr).replace("<>", "!=")
            cond_eval = re.sub(r"(?<![<>=!])=(?!=)", "==", cond_eval)
            try:
                cond_value = bool(eval(cond_eval, {"__builtins__": {}}, {}))
            except Exception:
                cond_value = False
            chosen = true_expr if cond_value else false_expr
            chosen_eval = ref_re.sub(repl, chosen)
            try:
                cache[key] = float(eval(chosen_eval, {"__builtins__": {}}, {}))
            except Exception:
                try:
                    cache[key] = float(chosen_eval)
                except Exception:
                    cache[key] = 0.0
            return cache[key]

    expr = ref_re.sub(repl, expr)
    try:
        result = eval(expr, {"__builtins__": {}}, {})
        cache[key] = float(result)
    except Exception:
        cache[key] = 0.0
    return cache[key]


def _extract_financial_chart_history_from_excel(excel_path: str, cmp_value: float | None) -> dict[str, tuple[list[str], list[float]]]:
    from openpyxl import load_workbook

    wb = load_workbook(excel_path, data_only=False)
    ws = wb["P&L"]
    periods = [str(ws.cell(2, c).value).strip() for c in range(2, ws.max_column + 1)]
    actual_cols = [c for c in range(2, ws.max_column + 1) if str(ws.cell(2, c).value).strip().endswith("A")][-5:]
    actual_periods = _format_actual_year_labels([ws.cell(2, c).value for c in actual_cols])
    cache: dict[tuple[str, str], float] = {}

    def row_values(row_num: int) -> list[float]:
        vals: list[float] = []
        for col in actual_cols:
            vals.append(_evaluate_excel_formula_cell(wb, "P&L", ws.cell(row_num, col).coordinate, cache) or 0.0)
        return vals

    revenue_vals = row_values(3)
    ebitda_vals = row_values(12)
    pat_vals = row_values(20)
    eps_vals = row_values(23)

    # P/E: prefer the pre-computed value from the Valuations_Table sheet — its
    # EPS is derived from the valuation model and stays smooth, whereas P&L row
    # 23 EPS spikes anomalously (e.g. Gravita FY26 = 4.13) and turns CMP/EPS
    # into a noisy 5-digit ratio. Fall back to CMP/EPS only if the sheet or
    # row is missing.
    pe_vals: list[float] = []
    if "Valuations_Table" in wb.sheetnames:
        try:
            vt = wb["Valuations_Table"]
            # Header row (4) labels each column with a period. Map by period.
            vt_headers = {
                str(vt.cell(4, c).value).strip(): c
                for c in range(2, vt.max_column + 1)
                if vt.cell(4, c).value
            }
            pe_row = None
            for r in range(1, vt.max_row + 1):
                label = str(vt.cell(r, 1).value or "").strip().lower()
                if label.startswith("p/e"):
                    pe_row = r
                    break
            if pe_row:
                for period in actual_periods:
                    col = vt_headers.get(period)
                    val = _evaluate_excel_formula_cell(wb, "Valuations_Table", vt.cell(pe_row, col).coordinate, cache) if col else None
                    pe_vals.append(round(float(val), 1) if val else 0.0)
        except Exception as exc:
            logger.debug("Valuations_Table P/E read failed: %s", exc)
            pe_vals = []
    if not pe_vals or not any(pe_vals):
        pe_vals = []
        for eps in eps_vals:
            if cmp_value and eps:
                pe_vals.append(round(cmp_value / eps, 1))
            else:
                pe_vals.append(0.0)

    return {
        "Revenue": (actual_periods, revenue_vals),
        "EBITDA": (actual_periods, ebitda_vals),
        "PAT": (actual_periods, pat_vals),
        "P/E": (actual_periods, pe_vals),
    }


def _extract_summary_dashboard_from_excel(excel_path: str) -> tuple[list[str], list[tuple[str, list[tuple[str, list[float]]]]]]:
    from openpyxl import load_workbook

    wb = load_workbook(excel_path, data_only=False)
    cache: dict[tuple[str, str], float] = {}
    pnl = wb["P&L"]
    bs = wb["Balance Sheet"]
    cf = wb["Cash Flow"]
    ratios = wb["Ratios"]

    all_headers = [str(pnl.cell(2, c).value).strip() for c in range(2, pnl.max_column + 1)]
    actual_cols = [c for c in range(2, pnl.max_column + 1) if str(pnl.cell(2, c).value).strip().endswith("A")]
    proj_cols = [c for c in range(2, pnl.max_column + 1) if str(pnl.cell(2, c).value).strip().endswith("E")]
    selected_cols = actual_cols[-3:] + proj_cols[:2]
    headers = [str(pnl.cell(2, c).value).strip() for c in selected_cols]

    def find_row(ws, label: str) -> int:
        for r in range(1, ws.max_row + 1):
            if str(ws.cell(r, 1).value).strip() == label:
                return r
        raise KeyError(f"Row '{label}' not found in {ws.title}")

    def vals(ws, row_label: str) -> list[float]:
        row = find_row(ws, row_label)
        out: list[float] = []
        for col in selected_cols:
            out.append(_evaluate_excel_formula_cell(wb, ws.title, ws.cell(row, col).coordinate, cache) or 0.0)
        return out

    sections: list[tuple[str, list[tuple[str, list[float]]]]] = [
        ("PROFIT & LOSS", [
            ("Net Revenue", vals(pnl, "Revenue")),
            ("EBITDA", vals(pnl, "EBITDA")),
            ("PAT", vals(pnl, "Profit After Tax (PAT)")),
            ("PAT Margin %", vals(pnl, "PAT Margin %")),
            ("EPS", vals(pnl, "EPS (₹)")),
        ]),
        ("BALANCE SHEET", [
            ("Net Worth", vals(bs, "Total Equity")),
            ("Total Debt", vals(bs, "Borrowings")),
            ("Capital Employed", [a + b for a, b in zip(vals(bs, "Total Equity"), vals(bs, "Borrowings"))]),
            ("Net Fixed Assets", vals(bs, "Net Block")),
            ("Working Capital", [
                rec + inv - liab
                for rec, inv, liab in zip(
                    vals(bs, "Trade Receivables"),
                    vals(bs, "Inventory"),
                    vals(bs, "Other Liabilities"),
                )
            ]),
            ("Debt/Equity", vals(ratios, "Debt/Equity (x)")),
        ]),
        ("CASH FLOW", [
            ("CFO", vals(cf, "CFO")),
            ("Capex", vals(pnl, "Depreciation")),  # temporary replaced below
            ("Free Cash Flow", vals(cf, "Free Cash Flow (FCF)")),
            ("CFO/EBITDA %", vals(ratios, "CFO/EBITDA")),
        ]),
        ("KEY RATIOS", [
            ("ROE %", vals(ratios, "ROE %")),
            ("ROCE %", vals(ratios, "ROCE %")),
            ("Inventory Days", vals(ratios, "Inventory Days")),
        ]),
        ("VALUATIONS (AT CMP)", [
            ("P/E (x)", [
                round((cmp / eps), 1) if eps else 0.0
                for cmp, eps in zip([vals(ratios, "EPS (₹)")[-1]] * len(headers), vals(ratios, "EPS (₹)"))
            ]),
            ("EV/EBITDA (x)", [
                round(((debt + equity) / ebitda), 1) if ebitda else 0.0
                for debt, equity, ebitda in zip(
                    vals(bs, "Borrowings"),
                    vals(bs, "Total Equity"),
                    vals(pnl, "EBITDA"),
                )
            ]),
        ]),
    ]

    # Replace capex with actual projected/historical capex row from Assumptions if available.
    if "Assumptions" in wb.sheetnames:
        asm = wb["Assumptions"]
        try:
            capex_row = find_row(asm, "Capex")
            capex_vals = [(_evaluate_excel_formula_cell(wb, asm.title, asm.cell(capex_row, col).coordinate, cache) or 0.0) for col in selected_cols]
            sections[2][1][1] = ("Capex", capex_vals)
        except Exception:
            pass

    return headers, sections


def _extract_summary_dashboard_from_excel_safe(
    excel_path: str,
) -> tuple[list[str], list[tuple[str, list[tuple[str, list[float]]]]]]:
    from openpyxl import load_workbook

    wb = load_workbook(excel_path, data_only=False)
    cache: dict[tuple[str, str], float] = {}
    pnl = wb["P&L"]
    bs = wb["Balance Sheet"]
    cf = wb["Cash Flow"]
    ratios = wb["Ratios"]

    actual_cols = [c for c in range(2, pnl.max_column + 1) if str(pnl.cell(2, c).value).strip().endswith("A")]
    proj_cols = [c for c in range(2, pnl.max_column + 1) if str(pnl.cell(2, c).value).strip().endswith("E")]
    selected_cols = actual_cols[-3:] + proj_cols[:2]
    headers = [str(pnl.cell(2, c).value).strip() for c in selected_cols]

    def _norm_label(value: object) -> str:
        text = str(value or "").strip()
        text = text.replace("â‚¹", "₹").replace("Rs.", "Rs")
        return " ".join(text.lower().split())

    def find_row(ws, *labels: str) -> int:
        wanted = {_norm_label(label) for label in labels if label}
        for r in range(1, ws.max_row + 1):
            for c in range(1, min(ws.max_column, 3) + 1):
                if _norm_label(ws.cell(r, c).value) in wanted:
                    return r
        raise KeyError(f"Row {labels!r} not found in {ws.title}")

    def vals(ws, *row_labels: str, default: list[float] | None = None) -> list[float]:
        try:
            row = find_row(ws, *row_labels)
        except KeyError:
            if default is not None:
                return list(default)
            raise
        out: list[float] = []
        for col in selected_cols:
            out.append(_evaluate_excel_formula_cell(wb, ws.title, ws.cell(row, col).coordinate, cache) or 0.0)
        return out

    def pct_from_series(numerator: list[float], denominator: list[float]) -> list[float]:
        return [round((num / den) * 100, 1) if den else 0.0 for num, den in zip(numerator, denominator)]

    ebitda_vals = vals(pnl, "EBITDA")
    cfo_vals = vals(cf, "CFO")
    eps_vals = vals(pnl, "EPS (₹)", "EPS (Rs)", "EPS", default=vals(ratios, "EPS (₹)", "EPS (Rs)", "EPS"))

    sections: list[tuple[str, list[tuple[str, list[float]]]]] = [
        ("PROFIT & LOSS", [
            ("Net Revenue", vals(pnl, "Revenue")),
            ("EBITDA", ebitda_vals),
            ("PAT", vals(pnl, "Profit After Tax (PAT)")),
            ("PAT Margin %", vals(pnl, "PAT Margin %")),
            ("EPS", eps_vals),
        ]),
        ("BALANCE SHEET", [
            ("Net Worth", vals(bs, "Total Equity")),
            ("Total Debt", vals(bs, "Borrowings")),
            ("Capital Employed", [a + b for a, b in zip(vals(bs, "Total Equity"), vals(bs, "Borrowings"))]),
            ("Net Fixed Assets", vals(bs, "Net Block")),
            ("Working Capital", [
                rec + inv - liab
                for rec, inv, liab in zip(
                    vals(bs, "Trade Receivables"),
                    vals(bs, "Inventory"),
                    vals(bs, "Other Liabilities"),
                )
            ]),
            ("Debt/Equity", vals(ratios, "Debt/Equity (x)")),
        ]),
        ("CASH FLOW", [
            ("CFO", cfo_vals),
            ("Capex", vals(pnl, "Depreciation")),
            ("Free Cash Flow", vals(cf, "Free Cash Flow (FCF)")),
            (
                "CFO/EBITDA %",
                vals(
                    ratios,
                    "CFO/EBITDA %",
                    "CFO/EBITDA",
                    default=pct_from_series(cfo_vals, ebitda_vals),
                ),
            ),
        ]),
        ("KEY RATIOS", [
            ("ROE %", vals(ratios, "ROE %")),
            ("ROCE %", vals(ratios, "ROCE %")),
            ("Inventory Days", vals(ratios, "Inventory Days")),
        ]),
        ("VALUATIONS (AT CMP)", [
            ("P/E (x)", [
                round((cmp / eps), 1) if eps else 0.0
                for cmp, eps in zip([eps_vals[-1]] * len(headers), eps_vals)
            ]),
            ("EV/EBITDA (x)", [
                round(((debt + equity) / ebitda), 1) if ebitda else 0.0
                for debt, equity, ebitda in zip(
                    vals(bs, "Borrowings"),
                    vals(bs, "Total Equity"),
                    ebitda_vals,
                )
            ]),
        ]),
    ]

    if "Assumptions" in wb.sheetnames:
        asm = wb["Assumptions"]
        try:
            capex_row = find_row(asm, "Capex")
            capex_vals = [
                (_evaluate_excel_formula_cell(wb, asm.title, asm.cell(capex_row, col).coordinate, cache) or 0.0)
                for col in selected_cols
            ]
            sections[2][1][1] = ("Capex", capex_vals)
        except Exception:
            pass

    return headers, sections


def _fmt_summary_cell(label: str, value: float) -> str:
    pct_labels = {"PAT Margin %", "CFO/EBITDA %", "ROE %", "ROCE %"}
    ratio_labels = {"Debt/Equity", "P/E (x)", "EV/EBITDA (x)"}
    if label in pct_labels:
        return f"{value:.1f}%"
    if label in ratio_labels:
        return f"{value:.1f}x"
    if label in {"EPS"}:
        return f"{value:.2f}"
    return f"{value:,.0f}"


def _render_financial_summary_dashboard(excel_path: str, company_name: str) -> bytes | None:
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception as exc:
        logger.warning("matplotlib unavailable for summary dashboard: %s", exc)
        return None

    headers, sections = _extract_summary_dashboard_from_excel_safe(excel_path)
    rows: list[list[str]] = [["Particulars", *headers]]
    row_styles: list[str] = ["header"]
    for section_name, items in sections:
        rows.append([section_name, *[""] * len(headers)])
        row_styles.append("section")
        for label, values in items:
            rows.append([label, *[_fmt_summary_cell(label, v) for v in values]])
            row_styles.append("data")

    fig_width = 6.0
    col_widths = [0.34] + [0.132] * len(headers)
    rows, line_counts = _wrap_table_cells(rows, col_widths, fig_width_inches=fig_width, chars_per_inch=13.0)
    # The slide-1 summary placeholder is tall (~6×9), so we allow more growth
    # here than the wide tables on slides 7/14/15/18 — but still cap it to
    # avoid python-pptx vertical-squashing.
    fig_h = min(10.5, max(6.2, 0.16 * sum(line_counts) + 0.9))
    fig, ax = plt.subplots(figsize=(fig_width, fig_h), facecolor="white")
    ax.axis("off")

    table = ax.table(
        cellText=rows,
        cellLoc="center",
        colWidths=col_widths,
        loc="upper left",
        bbox=[0, 0, 1, 1],
    )
    table.auto_set_font_size(False)
    table.set_fontsize(7.2)
    _apply_row_line_heights(table, line_counts)

    for (r, c), cell in table.get_celld().items():
        cell.set_edgecolor("#D5DCE8")
        cell.set_linewidth(0.5)
        if row_styles[r] == "header":
            cell.set_facecolor("#1F4690")
            cell.get_text().set_color("white")
            cell.get_text().set_fontweight("bold")
        elif row_styles[r] == "section":
            if c == 0:
                cell.set_facecolor("#FFA500")
                cell.get_text().set_color("white")
                cell.get_text().set_fontweight("bold")
            else:
                cell.set_facecolor("#FFF4D6")
                cell.get_text().set_text("")
        else:
            if c == 0:
                cell.set_facecolor("#EEF2FA")
                _set_cell_left_align(cell)
                cell.PAD = 0.02
            else:
                cell.set_facecolor("#F7FAFF" if c <= 3 else "#FFF1D9")

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=240, bbox_inches="tight", pad_inches=0.02, facecolor=fig.get_facecolor())
    plt.close(fig)
    return buf.getvalue()


def _last_five_actual_periods(periods: list[str], values: list[float]) -> tuple[list[str], list[float]]:
    pairs = [
        (str(period).strip(), value)
        for period, value in zip(periods, values)
        if period is not None
    ]
    actual_pairs = [
        (period, value)
        for period, value in pairs
        if "E" not in period.upper()
    ]
    actual_pairs = actual_pairs[-5:]
    return [p for p, _ in actual_pairs], [v for _, v in actual_pairs]


def _last_five_actual_operational(op: dict, key: str) -> tuple[list[str], list[float]]:
    years = [str(y).strip() for y in (op.get("years") or [])]
    values = [float(v) if v is not None else 0.0 for v in (op.get(key) or [])]
    return _last_five_actual_periods(years, values)


def _last_five_actual_segment_series(op: dict) -> tuple[list[str], dict[str, list[float]]]:
    years = [str(y).strip() for y in (op.get("years") or [])]
    actual_idx = [idx for idx, year in enumerate(years) if "E" not in year.upper()][-5:]
    actual_years = [years[idx] for idx in actual_idx]
    segments = op.get("volume_segments") or {}
    out: dict[str, list[float]] = {}
    for name, raw_vals in segments.items():
        vals = [float(v) if v is not None else 0.0 for v in (raw_vals or [])]
        out[name] = [vals[idx] for idx in actual_idx if idx < len(vals)]
    return actual_years, out


def _render_story_chart_collage(
    fin_model: dict,
    company_name: str,
    *,
    operational: bool,
    financial_history: dict[str, tuple[list[str], list[float]]] | None = None,
) -> bytes | None:
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception as exc:
        logger.warning("matplotlib unavailable for story-chart collage: %s", exc)
        return None

    # 3 columns × 2 rows = 6 charts per slide.
    fig, axes = plt.subplots(2, 3, figsize=(15.0, 6.6), facecolor="white")
    fig.suptitle(
        f"{company_name} — {'Operational Charts' if operational else 'Financial Charts'}",
        fontsize=14,
        color="#1F4690",
        fontweight="bold",
        y=0.98,
    )

    for ax in axes.flat:
        ax.set_facecolor("#F8FAFD")
        ax.grid(axis="y", color="#E5EAF3", linewidth=0.8)
        for spine in ax.spines.values():
            spine.set_visible(False)

    def _pad_top(ax, vals):
        """Push the upper y-limit ~10% above the max value so plotted lines
        don't kiss the top edge of the chart area (a common 'line getting cut'
        complaint when the data peaks right at axis_max)."""
        try:
            mx = max(v for v in vals if v is not None)
        except (TypeError, ValueError):
            return
        if mx <= 0:
            return
        ax.set_ylim(bottom=0, top=mx * 1.12)

    def _style_axes(ax, title: str, xlabel: str = "Fiscal Year", ylabel: str = "") -> None:
        ax.set_title(title, fontsize=10, color="#1F4690", fontweight="bold")
        ax.set_xlabel(xlabel, fontsize=8, color="#5C6B82")
        if ylabel:
            ax.set_ylabel(ylabel, fontsize=8, color="#5C6B82")
        ax.tick_params(axis="x", labelsize=8)
        ax.tick_params(axis="y", labelsize=8)

    def _empty(ax, title: str, msg: str = "n/a") -> None:
        _style_axes(ax, title)
        ax.text(0.5, 0.5, msg, ha="center", va="center", fontsize=10,
                color="#9CA3AF", transform=ax.transAxes)
        ax.set_xticks([])
        ax.set_yticks([])

    if operational:
        op = fin_model.get("operational") or {}

        # Shared palette for stacked/grouped charts.
        _PALETTE = ["#1F4690", "#FFA500", "#3A5BA0", "#16A34A", "#9CA3AF", "#7C3AED"]

        def _render_pie(ax, labels, sizes, title):
            """Donut with internal % labels and a legend BELOW the chart, so long
            category names never overlap the wedges. labels=None on ax.pie() prevents
            python-matplotlib from rendering them on the slices."""
            wedges, _txts, _pct = ax.pie(
                sizes,
                labels=None,
                autopct="%1.0f%%",
                pctdistance=0.7,
                colors=_PALETTE[:len(sizes)],
                textprops={"fontsize": 7, "color": "white", "fontweight": "bold"},
                wedgeprops={"width": 0.35, "edgecolor": "white"},
                startangle=90,
            )
            ax.set_title(title, fontsize=10, color="#1F4690", fontweight="bold")
            ax.legend(
                wedges,
                labels,
                loc="lower center",
                bbox_to_anchor=(0.5, -0.18),
                ncol=2,
                fontsize=7,
                frameon=False,
            )

        charts_spec = op.get("charts")
        if not charts_spec:
            # ── 1. Capacity (MTPA / kT) ───────────────────────────────────────
            cap_years, cap_vals = _last_five_actual_operational(op, "capacity_mtpa")
            if not cap_years:
                cap_years, cap_vals = _last_five_actual_operational(op, "capacity_kt")
            if not cap_years:
                seg_years_all, vol_segs_all = _last_five_actual_segment_series(op)
                if seg_years_all and vol_segs_all:
                    totals = [0.0] * len(seg_years_all)
                    for vals in vol_segs_all.values():
                        for i, v in enumerate(vals[:len(totals)]):
                            totals[i] += v
                    cap_years, cap_vals = seg_years_all, [t / 1000.0 for t in totals]

            if cap_years and cap_vals:
                axes[0, 0].bar(cap_years, cap_vals, color="#1F4690")
                _style_axes(axes[0, 0], "Capacity (kT)", ylabel="kT")
                _pad_top(axes[0, 0], cap_vals)
            else:
                _empty(axes[0, 0], "Capacity")

            # ── 2. Volume Sold (MT) — total across segments ────────────────────
            seg_years, volume_segments = _last_five_actual_segment_series(op)
            if seg_years and volume_segments:
                totals = [0.0] * len(seg_years)
                for vals in volume_segments.values():
                    for i, v in enumerate(vals[:len(totals)]):
                        totals[i] += v
                axes[0, 1].bar(seg_years, totals, color="#FFA500")
                _style_axes(axes[0, 1], "Volume Sold (MT)", ylabel="MT")
                _pad_top(axes[0, 1], totals)
            else:
                _empty(axes[0, 1], "Volume Sold")

            # ── 3. Volume by Category ──────────────────────────────────────────
            # Determine dynamic category (e.g. Metal, Product, Segment)
            explicit_category = op.get("volume_category_name") or op.get("volume_dimension") or op.get("volume_by_label")
            if explicit_category and isinstance(explicit_category, str) and explicit_category.strip():
                vol_category = explicit_category.strip()
            else:
                # Detect based on keys
                detect_keys = list(volume_segments.keys()) if volume_segments else []
                if not detect_keys:
                    metal_mix_temp = op.get("volume_by_metal") or op.get("metal_volume_pct") or {}
                    if isinstance(metal_mix_temp, dict):
                        detect_keys = list(metal_mix_temp.keys())
                
                metals = {"lead", "zinc", "copper", "aluminum", "aluminium", "nickel", "iron", "steel", "gold", "silver", "cobalt", "brass", "bronze", "tin", "metal", "alloy", "scrap", "wire", "cable"}
                products = {"polymer", "chemical", "plastic", "paper", "oil", "gas", "coal", "fuel", "material", "product", "battery", "batteries", "grid", "packaging"}
                
                keys_lower = [str(k).lower() for k in detect_keys]
                metal_count = sum(1 for k in keys_lower if any(m in k for m in metals))
                product_count = sum(1 for k in keys_lower if any(p in k for p in products))
                
                if metal_count > 0 and metal_count >= product_count:
                    vol_category = "Metal"
                elif product_count > 0:
                    vol_category = "Product"
                else:
                    vol_category = "Segment"

            # Read volume unit dynamically from model data, defaulting to MT
            vol_unit = op.get("volume_unit") or op.get("volume_units") or op.get("unit") or "MT"

            if seg_years and volume_segments:
                bottoms = [0.0] * len(seg_years)
                for idx, (name, vals) in enumerate(volume_segments.items()):
                    axes[0, 2].bar(seg_years, vals[:len(seg_years)], bottom=bottoms,
                                   color=_PALETTE[idx % len(_PALETTE)], label=name)
                    bottoms = [b + v for b, v in zip(bottoms, vals[:len(bottoms)])]
                _style_axes(axes[0, 2], f"{vol_category} Volume ({vol_unit})", ylabel=vol_unit)
                _pad_top(axes[0, 2], bottoms)
                axes[0, 2].legend(fontsize=7, frameon=False, loc="upper left")
            else:
                # Single-year dict fallback — render as donut.
                metal_mix = op.get("volume_by_metal") or op.get("metal_volume_pct") or {}
                if isinstance(metal_mix, dict) and metal_mix:
                    labels_v = list(metal_mix.keys())
                    sizes_v = []
                    for v in metal_mix.values():
                        try:
                            fv = float(v)
                        except (TypeError, ValueError):
                            fv = 0.0
                        if fv <= 1.0:
                            fv *= 100.0
                        sizes_v.append(fv)

                    if sum(sizes_v) > 0:
                        _render_pie(axes[0, 2], labels_v, sizes_v, f"{vol_category} Volume Mix")
                    else:
                        _empty(axes[0, 2], f"{vol_category} Volume")
                else:
                    _empty(axes[0, 2], f"{vol_category} Volume")

            # ── 4. Geographic Revenue Mix (India vs International) ─────────────
            geo_series = op.get("revenue_by_geography") or op.get("geographic_mix_series") or {}
            rendered_geo = False
            if isinstance(geo_series, dict) and geo_series and all(isinstance(v, list) for v in geo_series.values()):
                geo_year_axis = [str(y).strip() for y in (op.get("years") or [])]
                actual_idx = [i for i, y in enumerate(geo_year_axis) if "E" not in y.upper()][-5:]
                geo_year_axis = [geo_year_axis[i] for i in actual_idx]
                if geo_year_axis:
                    bottoms = [0.0] * len(geo_year_axis)
                    for idx, (name, vals) in enumerate(list(geo_series.items())[:5]):
                        aligned = [float(vals[i]) if i < len(vals) and vals[i] is not None else 0.0 for i in actual_idx]
                        axes[1, 0].bar(geo_year_axis, aligned, bottom=bottoms,
                                       color=_PALETTE[idx % len(_PALETTE)], label=name)
                        bottoms = [b + a for b, a in zip(bottoms, aligned)]
                    _style_axes(axes[1, 0], "India vs International", ylabel="%")
                    _pad_top(axes[1, 0], bottoms)
                    axes[1, 0].legend(fontsize=7, frameon=False, loc="upper left")
                    rendered_geo = True
            if not rendered_geo:
                geo_mix = op.get("geography_mix_pct") or op.get("geographic_mix") or {}
                if isinstance(geo_mix, dict) and geo_mix:
                    india_share = 0.0
                    intl_share = 0.0
                    for k, v in geo_mix.items():
                        try:
                            fv = float(v)
                        except (TypeError, ValueError):
                            continue
                        if fv <= 1.0:
                            fv *= 100.0
                        if "india" in str(k).lower():
                            india_share += fv
                        else:
                            intl_share += fv
                    if india_share + intl_share > 0:
                        _render_pie(axes[1, 0], ["India", "International"], [india_share, intl_share],
                                    "India vs International")
                    else:
                        _empty(axes[1, 0], "India vs International")
                else:
                    _empty(axes[1, 0], "India vs International")

            # ── 5. Plants (India vs Overseas — grouped bars) ───────────────────
            plant_years, india = _last_five_actual_operational(op, "plants_india")
            _, overseas = _last_five_actual_operational(op, "plants_overseas")
            if plant_years and (india or overseas):
                import numpy as _np
                india = india or [0.0] * len(plant_years)
                overseas = overseas or [0.0] * len(plant_years)
                x = _np.arange(len(plant_years))
                w = 0.38
                axes[1, 1].bar(x - w / 2, india,    width=w, color="#1F4690", label="India")
                axes[1, 1].bar(x + w / 2, overseas, width=w, color="#FFA500", label="Overseas")
                axes[1, 1].set_xticks(x)
                axes[1, 1].set_xticklabels(plant_years)
                _style_axes(axes[1, 1], "Plants", ylabel="Count")
                _pad_top(axes[1, 1], [max(a or 0, b or 0) for a, b in zip(india, overseas)])
                axes[1, 1].legend(fontsize=7, frameon=False, loc="upper left")
            else:
                _empty(axes[1, 1], "Plants")

            # ── 6. Capacity Utilization % ──────────────────────────────────────
            util_years, util_vals = _last_five_actual_operational(op, "capacity_utilisation_pct")
            if not util_years:
                util_years, util_vals = _last_five_actual_operational(op, "capacity_utilization_pct")
            if not util_years:
                util_years, util_vals = _last_five_actual_operational(op, "utilization_pct")
            if util_years and util_vals:
                axes[1, 2].plot(util_years, util_vals, color="#1F4690", linewidth=2.5, marker="o")
                _style_axes(axes[1, 2], "Capacity Utilization (%)", ylabel="%")
                _pad_top(axes[1, 2], util_vals)
            else:
                _empty(axes[1, 2], "Capacity Utilization")
        else:
            # ──── DYNAMIC METADATA-DRIVEN RENDERER ────
            # Clear default subplots figure since we will construct a custom grid shape
            plt.close(fig)
            
            n_charts = len(charts_spec)
            if n_charts <= 0:
                fig, axes = plt.subplots(1, 1, figsize=(15.0, 6.6), facecolor="white")
                axes_list = [axes]
            elif n_charts <= 3:
                fig, axes = plt.subplots(1, n_charts, figsize=(15.0, 6.6), facecolor="white")
                axes_list = list(axes) if n_charts > 1 else [axes]
            else:
                rows = 2
                cols = (n_charts + 1) // 2
                fig, axes = plt.subplots(rows, cols, figsize=(15.0, 6.6), facecolor="white")
                axes_list = list(axes.flat)

            fig.suptitle(
                f"{company_name} — Operational Performance",
                fontsize=14,
                color="#1F4690",
                fontweight="bold",
                y=0.98,
            )

            # Re-style all axes in our grid
            for ax in axes_list:
                ax.set_facecolor("#F8FAFD")
                ax.grid(axis="y", color="#E5EAF3", linewidth=0.8)
                for spine in ax.spines.values():
                    spine.set_visible(False)

            # Draw each dynamic chart
            for idx, chart in enumerate(charts_spec):
                if idx >= len(axes_list):
                    break
                ax = axes_list[idx]
                title = chart.get("title") or "Metric"
                chart_type = chart.get("chart_type") or "bar"
                years = [str(y) for y in (chart.get("years") or [])]
                ylabel = chart.get("ylabel") or ""

                stacked_values = chart.get("stacked_values") or {}
                series_names = chart.get("series_names") or list(stacked_values.keys())

                if chart_type == "stacked_bar" and series_names and stacked_values:
                    bottoms = [0.0] * len(years)
                    for s_idx, s_name in enumerate(series_names):
                        s_vals = stacked_values.get(s_name) or []
                        s_vals_num = [float(v) if v is not None else 0.0 for v in s_vals]
                        s_vals_num = s_vals_num[:len(years)] + [0.0] * (len(years) - len(s_vals_num))
                        ax.bar(years, s_vals_num, bottom=bottoms,
                               color=_PALETTE[s_idx % len(_PALETTE)], label=s_name)
                        bottoms = [b + v for b, v in zip(bottoms, s_vals_num)]
                    _style_axes(ax, title, ylabel=ylabel)
                    _pad_top(ax, bottoms)
                    ax.legend(fontsize=7, frameon=False, loc="upper left")

                elif chart_type == "donut":
                    raw_vals = chart.get("values") or []
                    sizes = [float(v) if v is not None else 0.0 for v in raw_vals]
                    labels = [years[i] if i < len(years) else f"Slice {i}" for i in range(len(sizes))]
                    valid_pairs = [(l, s) for l, s in zip(labels, sizes) if s > 0]
                    if valid_pairs:
                        _render_pie(ax, [p[0] for p in valid_pairs], [p[1] for p in valid_pairs], title)
                    else:
                        _empty(ax, title, "n/a")

                elif chart_type == "line":
                    raw_vals = chart.get("values") or []
                    vals = [float(v) if v is not None else 0.0 for v in raw_vals]
                    vals = vals[:len(years)] + [0.0] * (len(years) - len(vals))
                    ax.plot(years, vals, color="#1F4690", linewidth=2.5, marker="o")
                    _style_axes(ax, title, ylabel=ylabel)
                    _pad_top(ax, vals)

                else:  # "bar" chart
                    raw_vals = chart.get("values") or []
                    vals = [float(v) if v is not None else 0.0 for v in raw_vals]
                    vals = vals[:len(years)] + [0.0] * (len(years) - len(vals))
                    ax.bar(years, vals, color="#1F4690")
                    _style_axes(ax, title, ylabel=ylabel)
                    _pad_top(ax, vals)

            # Hide any unused slots in the subplot grid
            for idx in range(n_charts, len(axes_list)):
                axes_list[idx].axis("off")

    else:
        financial_history = financial_history or {}
        rev_years, rev_vals = financial_history.get("Revenue", ([], []))
        ebitda_years, ebitda_vals = financial_history.get("EBITDA", ([], []))
        pat_years, pat_vals = financial_history.get("PAT", ([], []))

        if not rev_years:
            rev_years, rev_vals = _series_values(fin_model, "Revenue")
            rev_years, rev_vals = _last_five_actual_periods(rev_years, rev_vals)
        if not ebitda_years:
            ebitda_years, ebitda_vals = _series_values(fin_model, "EBITDA")
            ebitda_years, ebitda_vals = _last_five_actual_periods(ebitda_years, ebitda_vals)
        if not pat_years:
            pat_years, pat_vals = _series_values(fin_model, "PAT")
            pat_years, pat_vals = _last_five_actual_periods(pat_years, pat_vals)

        # historical_ratios live on the raw JSON, not the mapped series — sniff
        # from a few well-known shapes that _enrich and _build leave behind.
        raw_hist = fin_model.get("historical_ratios") or {}
        hr_years = [str(y).strip() for y in (raw_hist.get("years") or [])]

        def _hr(key: str) -> tuple[list[str], list[float]]:
            vals = raw_hist.get(key)
            if not vals or not hr_years:
                return [], []
            pairs = [(y, v) for y, v in zip(hr_years, vals) if v is not None and "E" not in y.upper()]
            pairs = pairs[-5:]
            return [p for p, _ in pairs], [float(v) for _, v in pairs]

        # 1. Revenue with 5-yr CAGR in title.
        cagr_suffix = ""
        if len(rev_vals) >= 2 and rev_vals[0] and rev_vals[-1]:
            try:
                n = len(rev_vals) - 1
                cagr = ((rev_vals[-1] / rev_vals[0]) ** (1.0 / n) - 1.0) * 100.0
                cagr_suffix = f"  ({n}Y CAGR {cagr:.1f}%)"
            except (ZeroDivisionError, ValueError):
                pass
        if rev_years and rev_vals:
            axes[0, 0].bar(rev_years, rev_vals, color="#1F4690")
            _style_axes(axes[0, 0], f"Revenue (₹ Cr){cagr_suffix}", ylabel="₹ Cr")
            _pad_top(axes[0, 0], rev_vals)
        else:
            _empty(axes[0, 0], "Revenue")

        # 2. EBITDA bars + margin % line overlay.
        ebitda_margin_years, ebitda_margin = _hr("ebitda_margin_pct")
        if ebitda_years and ebitda_vals:
            axes[0, 1].bar(ebitda_years, ebitda_vals, color="#FFA500")
            _style_axes(axes[0, 1], "EBITDA (₹ Cr)", ylabel="₹ Cr")
            _pad_top(axes[0, 1], ebitda_vals)
            if ebitda_margin_years and ebitda_margin:
                ax2 = axes[0, 1].twinx()
                # Align to displayed x-axis
                aligned = {y: m for y, m in zip(ebitda_margin_years, ebitda_margin)}
                m_y = [y for y in ebitda_years if y in aligned]
                m_v = [aligned[y] for y in m_y]
                if m_y:
                    ax2.plot(m_y, m_v, color="#1F4690", linewidth=2.0, marker="o", label="Margin %")
                    ax2.set_ylabel("Margin %", fontsize=8, color="#1F4690")
                    ax2.tick_params(axis="y", labelsize=8, colors="#1F4690")
                    for spine in ax2.spines.values():
                        spine.set_visible(False)
        else:
            _empty(axes[0, 1], "EBITDA")

        # 3. PAT bars + PAT margin % line overlay.
        pat_margin_years, pat_margin = _hr("pat_margin_pct")
        if pat_years and pat_vals:
            axes[0, 2].bar(pat_years, pat_vals, color="#3A5BA0")
            _style_axes(axes[0, 2], "PAT (₹ Cr)", ylabel="₹ Cr")
            _pad_top(axes[0, 2], pat_vals)
            if pat_margin_years and pat_margin:
                ax2 = axes[0, 2].twinx()
                aligned = {y: m for y, m in zip(pat_margin_years, pat_margin)}
                m_y = [y for y in pat_years if y in aligned]
                m_v = [aligned[y] for y in m_y]
                if m_y:
                    ax2.plot(m_y, m_v, color="#FFA500", linewidth=2.0, marker="o", label="Margin %")
                    ax2.set_ylabel("Margin %", fontsize=8, color="#FFA500")
                    ax2.tick_params(axis="y", labelsize=8, colors="#FFA500")
                    for spine in ax2.spines.values():
                        spine.set_visible(False)
        else:
            _empty(axes[0, 2], "PAT")

        # 4. Return Ratios — ROE & ROCE dual line.
        roe_years, roe_vals = _hr("roe_pct")
        roce_years, roce_vals = _hr("roce_pct")
        if (roe_years and roe_vals) or (roce_years and roce_vals):
            all_vals = []
            if roe_years and roe_vals:
                axes[1, 0].plot(roe_years, roe_vals, color="#1F4690", linewidth=2.5, marker="o", label="ROE %")
                all_vals.extend(roe_vals)
            if roce_years and roce_vals:
                axes[1, 0].plot(roce_years, roce_vals, color="#FFA500", linewidth=2.5, marker="s", label="ROCE %")
                all_vals.extend(roce_vals)
            _style_axes(axes[1, 0], "Return Ratios (%)", ylabel="%")
            if all_vals:
                _pad_top(axes[1, 0], all_vals)
            axes[1, 0].legend(fontsize=7, frameon=False, loc="upper left")
        else:
            _empty(axes[1, 0], "Return Ratios")

        # 5. Working Capital Days — receivable + inventory − payable.
        rec_years, rec_vals = _hr("receivable_days")
        if not rec_years:
            rec_years, rec_vals = _hr("debtor_days")
        inv_years, inv_vals = _hr("inventory_days")
        pay_years, pay_vals = _hr("payable_days")
        if not pay_years:
            pay_years, pay_vals = _hr("creditor_days")
        # Choose the longest year list as anchor.
        anchor_years = rec_years or inv_years or pay_years
        if anchor_years:
            def _aligned(yrs, vals):
                m = {y: v for y, v in zip(yrs, vals)}
                return [m.get(y, 0.0) for y in anchor_years]
            r = _aligned(rec_years, rec_vals)
            i = _aligned(inv_years, inv_vals)
            p = _aligned(pay_years, pay_vals)
            wc = [a + b - c for a, b, c in zip(r, i, p)]
            axes[1, 1].bar(anchor_years, wc, color="#16A34A")
            _style_axes(axes[1, 1], "Working Capital Days", ylabel="Days")
            _pad_top(axes[1, 1], wc)
        else:
            _empty(axes[1, 1], "Working Capital Days")

        # 6. P/E & EV/EBITDA — dual line, fall back to single-series with renamed title.
        pe_years, pe_vals = _hr("pe")
        if not pe_years:
            pe_years, pe_vals = _hr("pe_ratio")
        ev_years, ev_vals = _hr("ev_ebitda")
        title6 = "P/E & EV/EBITDA (x)"
        if (pe_years and pe_vals) or (ev_years and ev_vals):
            all_vals = []
            if pe_years and pe_vals:
                axes[1, 2].plot(pe_years, pe_vals, color="#1F4690", linewidth=2.5, marker="o", label="P/E")
                all_vals.extend(pe_vals)
                if not (ev_years and ev_vals):
                    title6 = "P/E (x)"
            if ev_years and ev_vals:
                axes[1, 2].plot(ev_years, ev_vals, color="#FFA500", linewidth=2.5, marker="s", label="EV/EBITDA")
                all_vals.extend(ev_vals)
                if not (pe_years and pe_vals):
                    title6 = "EV/EBITDA (x)"
            _style_axes(axes[1, 2], title6, ylabel="x")
            if all_vals:
                _pad_top(axes[1, 2], all_vals)
            axes[1, 2].legend(fontsize=7, frameon=False, loc="upper left")
        else:
            _empty(axes[1, 2], "P/E & EV/EBITDA")

    fig.tight_layout(rect=[0, 0, 1, 0.95], pad=1.0)
    # Add extra room below for the pie/donut legends and a hair more between
    # the two rows so the operational pies don't get clipped.
    fig.subplots_adjust(bottom=0.10, hspace=0.55)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=200, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)
    return buf.getvalue()


def inject_story_chart_slides(
    pptx_path: str,
    fin_model: dict,
    company_name: str,
    *,
    excel_path: str | None = None,
    cmp_value: float | None = None,
) -> int:
    """Replace slides 2 and 3 placeholders with chart collages built from model data."""
    financial_history: dict[str, tuple[list[str], list[float]]] | None = None
    if excel_path:
        try:
            financial_history = _extract_financial_chart_history_from_excel(excel_path, cmp_value)
        except Exception as exc:
            logger.warning("Excel history extraction failed for story charts: %s", exc)

    prs = Presentation(pptx_path)
    replacements = {
        2: ({"{{financial_charts}}", "{{financial_model_from_excel}}", "Financial model charts — see Excel model for details."}, _render_story_chart_collage(
            fin_model,
            company_name,
            operational=False,
            financial_history=financial_history,
        )),
        3: ({"{{operational_charts}}", "{{financial_model_from_excel_operational_sheet}}", "Operational data — see Excel model for details."}, _render_story_chart_collage(fin_model, company_name, operational=True)),
    }
    injected = 0

    for slide_idx, (tokens, img_bytes) in replacements.items():
        if not img_bytes or slide_idx > len(prs.slides):
            continue
        slide = prs.slides[slide_idx - 1]
        target_shape = None
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame and shape.text_frame.text.strip() in tokens:
                target_shape = shape
                break
        if target_shape is None:
            continue
        _insert_image_into_shape(slide, target_shape, img_bytes)
        injected += 1

    if injected:
        prs.save(pptx_path)
    return injected


def inject_financial_summary_slide(pptx_path: str, company_name: str, *, excel_path: str | None = None) -> int:
    if not excel_path:
        return 0
    try:
        img_bytes = _render_financial_summary_dashboard(excel_path, company_name)
    except Exception as exc:
        logger.warning("Financial summary dashboard render failed: %s", exc)
        return 0
    if not img_bytes:
        return 0

    prs = Presentation(pptx_path)
    if not prs.slides:
        return 0
    slide = prs.slides[0]
    target_shape = None
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame and shape.text_frame.text.strip() == "{{financial_summary_image}}":
            target_shape = shape
            break
    if target_shape is None:
        return 0
    _insert_image_into_shape(slide, target_shape, img_bytes)
    prs.save(pptx_path)
    return 1


def fill_master_template(
    template_path: str,
    output_path: str,
    replacements: dict,
    fin_model: dict,
    company_name: str = "",
    summary_image: bytes | None = None,
) -> None:
    """Fill master_template.pptx per-slide with text replacements, chart data,
    financial tables, and generated chart images.

    Shape naming convention in the template (PowerPoint Selection Pane):
      - Any shape              : plain text with {{placeholder}} tokens → replaced
      - Named 'slide_type'     : text content declares the slide type (e.g. 'earnings_forecast')
      - Named 'chart:<key>'    : replaced with a matplotlib chart image (see chart_generators.py)
      - Named 'table:<key>'    : populated from financial model series data
      - Native PPTX chart      : data updated with per-slide-type series from financial model

    Shapes whose entire text matches an Excel injection placeholder token
    (e.g. {{financial_summary_image}}) are deliberately SKIPPED so the
    downstream COM-based excel_injector can still locate them.
    """
    prs = Presentation(template_path)

    # Stale literal strings baked into the master template that must be rewritten
    # to the current company. These are not {{placeholder}} tokens, so the
    # token-based replacer skips them; we patch the slide text directly.
    literal_disclosure_subs: list[tuple[str, str]] = []
    if company_name:
        literal_disclosure_subs.append(
            ("(Premier Energies Ltd.)", f"({company_name})")
        )

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_type = _detect_slide_type(slide)
        pptx_chart_data = _build_pptx_chart_data(slide_type, fin_model)

        if slide_idx == 20 and literal_disclosure_subs:
            _apply_literal_text_subs(slide, literal_disclosure_subs)

        if slide_idx == 4:
            _replace_slide4_thesis_shapes(slide, replacements)
        if slide_idx == 5:
            _replace_slide6_overview_shapes(slide, replacements)
        if slide_idx == 10:
            _replace_slide11_management_shapes(slide, replacements)
        if slide_idx == 11:
            _replace_slide12_indicator_shapes(slide, replacements)

        # Shapes marked for image insertion must be processed after iteration
        # (removing while iterating breaks the list).  Collect them first.
        image_insertions: list[tuple] = []  # (shape, img_bytes)

        for shape in slide.shapes:
            shape_name = (shape.name or "").lower().strip()

            if (
                summary_image
                and hasattr(shape, "text_frame")
                and shape.text_frame
                and shape.text_frame.text.strip() == "{{financial_summary_image}}"
            ):
                image_insertions.append((shape, summary_image))
                continue

            # ── 0. Preserve Excel injection placeholders for COM injector ─────
            if _shape_is_excel_placeholder(shape):
                logger.debug("Preserving Excel placeholder shape: %s", shape.text_frame.text.strip())
                continue

            # ── 1. chart:<key> shape → generate PNG and queue for insertion ──
            if shape_name.startswith("chart:"):
                chart_key = shape_name[6:]
                img = generate_chart(chart_key, fin_model, company_name)
                if img:
                    image_insertions.append((shape, img))
                continue  # skip further processing on this shape

            # ── 2. Native PPTX chart → update series data per slide type ──────
            if shape.has_chart and pptx_chart_data:
                try:
                    shape.chart.replace_data(pptx_chart_data)
                except Exception as e:
                    logger.warning("chart.replace_data failed on slide '%s': %s", slide_type, e)

            # ── 3. table:<key> or financial slide table → populate from model ─
            if shape.has_table:
                # Determine table key: explicit name wins, else use slide type
                if shape_name.startswith("table:"):
                    table_key = shape_name[6:]
                else:
                    table_key = slide_type
                _fill_table_from_model(shape.table, table_key, fin_model)
                _replace_text_in_table(shape, replacements)

            # ── 4. Text frame → {{placeholder}} replacement ───────────────────
            elif hasattr(shape, "text_frame") and shape.text_frame:
                if slide_idx == 4 and shape.text_frame.text.strip() in {"{{investment_thesis}}", "{{investment_thesis_heading}}"}:
                    continue
                if slide_idx == 7 and shape.text_frame.text.strip() in {"{{COMPANY_TIMELINE}}", "{{company_timeline}}"}:
                    continue
                if slide_idx == 5 and shape.text_frame.text.strip() == "{{COMPANY_OVERVIEW}}":
                    continue
                if slide_idx == 10 and shape.text_frame.text.strip() in {"{{management_content}}", "{{management_commentry_heading}}"}:
                    continue
                if slide_idx == 11 and shape.text_frame.text.strip() == "{{indicators}}":
                    continue
                if slide_idx == 10 and shape.text_frame.text.strip() == "{{catalyst_timeline_chart}}":
                    continue
                _replace_text_in_frame(shape.text_frame, replacements)

        # ── 5. Insert queued chart images (after shape iteration is done) ──────
        for shape, img_bytes in image_insertions:
            try:
                _insert_image_into_shape(slide, shape, img_bytes)
            except Exception as e:
                logger.warning("Image insert failed for '%s': %s", shape.name, e)

    prs.save(output_path)


def _cleanup_excel_placeholders(pptx_path: str, replacements: dict) -> int:
    """Replace any surviving Excel injection placeholder tokens with fallback text.

    After the COM-based excel_injector runs (or is skipped), some placeholder
    shapes like {{financial_summary_image}} may still be present if injection
    failed or was unavailable.  This pass replaces them with text content from
    the replacements dict so raw {{...}} tokens never appear in the final deck.

    Returns the number of placeholders cleaned up.
    """
    # Build a map from the full token text → fallback content
    _FALLBACK_MAP: dict[str, str] = {
        "{{financial_model_from_excel}}": "Financial model charts — see Excel model for details.",
        "{{financial_model_from_excel_operational_sheet}}": "Operational data — see Excel model for details.",
        "{{financial_charts}}": "Financial model charts — see Excel model for details.",
        "{{operational_charts}}": "Operational data — see Excel model for details.",
        "{{financial_summary_image}}": "Financial summary — see Excel model for details.",
        "{{earnings_forecast_table}}": "Earnings forecast — see Excel model for details.",
        "{{financials_table}}": "Financials — see Excel model for details.",
        "{{valuations_table}}": "Valuations — see Excel model for details.",
        "{{key_risks_table}}": "Key risks — see Excel model for details.",
        "{{peer_comparision}}": replacements.get("peer_comparision", "Peer comparison — see Excel model for details."),
        "{{governance_table}}": replacements.get("indicators", "Governance — see Excel model for details."),
        "{{timeline}}": "Timeline — see Excel model for details.",
        "{{company_timeline}}": "Timeline — see Excel model for details.",
        "{{competitive_chart_1}}": "Competitive positioning chart — see Excel model for details.",
        "{{competitive_chart_2}}": "Competitive positioning chart — see Excel model for details.",
        "{{peer_comparison_chart_1}}": "Competitive positioning chart — see Excel model for details.",
        "{{peer_comparison_chart_2}}": "Competitive positioning chart — see Excel model for details.",
        "{{pie_chart_1}}": "Segment breakdown — see Excel model for details.",
        "{{pie_chart_2}}": "Segment breakdown — see Excel model for details.",
        "{{percentage_revenue_pie_chart}}": "Segment breakdown — see Excel model for details.",
        "{{percentage_EBIT_pie_chart}}": "Segment breakdown — see Excel model for details.",
        "{{catalyst_timeline_chart}}": "Catalyst timeline — see Excel model for details.",
        "{{probability_weight_table}}": "Probability-weighted scenario analysis — see Excel model for details.",
    }

    prs = Presentation(pptx_path)
    cleaned = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or not shape.text_frame:
                continue
            try:
                full_text = shape.text_frame.text.strip()
            except Exception:
                continue
            if full_text in _FALLBACK_MAP:
                fallback = _FALLBACK_MAP[full_text]
                # Preserve font from first run
                saved_font: dict = {}
                for para in shape.text_frame.paragraphs:
                    if para.runs:
                        f = para.runs[0].font
                        saved_font["name"] = f.name
                        saved_font["size"] = f.size
                        saved_font["bold"] = f.bold
                        break

                # Clear and replace
                for para in shape.text_frame.paragraphs:
                    para.clear()
                first_para = shape.text_frame.paragraphs[0]
                run = first_para.add_run()
                run.text = fallback
                font_name = _get_fallback_font(saved_font.get("name"))
                if font_name:
                    run.font.name = font_name
                if saved_font.get("size"):
                    run.font.size = saved_font["size"]

                cleaned += 1
                logger.info("Cleaned up unreplaced Excel placeholder: %s", full_text)

    if cleaned > 0:
        prs.save(pptx_path)
        logger.info("Cleanup pass replaced %d surviving Excel placeholders", cleaned)

    return cleaned


def _convert_pptx_to_pdf(pptx_path: Path, output_dir: Path) -> Path | None:
    """Headless LibreOffice conversion from PPTX to PDF."""
    import subprocess
    soffice_bin = shutil.which("soffice") or shutil.which("soffice.exe")
    if not soffice_bin:
        # Check standard Windows paths if not in PATH
        for win_path in [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]:
            if os.path.exists(win_path):
                soffice_bin = win_path
                break

    if not soffice_bin:
        logger.warning("LibreOffice (soffice) not found; PDF conversion skipped.")
        return None

    logger.info("Converting PPTX to PDF: %s -> %s", pptx_path, output_dir)
    try:
        cmd = [
            soffice_bin,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_dir),
            str(pptx_path),
        ]
        res = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        if res.returncode != 0:
            logger.error("LibreOffice conversion failed (code %d): %s", res.returncode, res.stderr)
            return None

        pdf_path = output_dir / pptx_path.with_suffix(".pdf").name
        if pdf_path.exists():
            logger.info("PDF generated successfully: %s", pdf_path)
            return pdf_path
        logger.warning("LibreOffice completed but output PDF not found at %s", pdf_path)
        return None
    except Exception as e:
        logger.error("Error during LibreOffice PDF conversion: %s", e)
        return None


def generate_pptx_for_report(report_id: str, session_id: str, *, use_mock: bool = False) -> dict:
    """Top-level orchestrator. Returns the response payload for /generate-pptx."""
    t0 = time.time()
    warnings: list[str] = []

    client = get_service_client()
    report, session, sections = _fetch_inputs(client, report_id, session_id)

    ticker = (report.get("nse_symbol") or "UNKNOWN").upper()
    logger.info("generate_pptx start report=%s ticker=%s sections=%d", report_id, ticker, len(sections))

    company = _build_company(report, session, sections)
    metadata = _build_metadata(report, sections)
    model_json = _download_model_json(client, ticker, warnings)
    fin_model = _build_financial_model(ticker, model_json, warnings)
    fin_model = _enrich_financial_model_for_house_deck(fin_model, report, sections, metadata)
    md_body = _build_approved_report_md(report, sections)

    with tempfile.TemporaryDirectory(prefix="reportgen_") as tmp:
        tmp_root = Path(tmp)
        company_path = tmp_root / "company.json"
        metadata_path = tmp_root / "metadata.json"
        model_path = tmp_root / "financial_model.json"
        report_path = tmp_root / "approved_report.md"
        bundle_path = tmp_root / "bundle.json"
        output_root = tmp_root / "out"
        output_root.mkdir(parents=True, exist_ok=True)

        company_path.write_text(json.dumps(company, default=str), encoding="utf-8")
        metadata_path.write_text(json.dumps(metadata, default=str), encoding="utf-8")
        model_path.write_text(json.dumps(fin_model, default=str), encoding="utf-8")
        report_path.write_text(md_body, encoding="utf-8")

        bundle = {
            "company_path": str(company_path),
            "metadata_path": str(metadata_path),
            "financial_model_path": str(model_path),
            "approved_report_path": str(report_path),
        }
        bundle_path.write_text(json.dumps(bundle), encoding="utf-8")

        logger.info("Filling master template directly")

        replacements = map_replacements(company, metadata, fin_model, sections)

        # ── PPT copywriting LLM pass output (research_sessions.ppt_content_json) ──
        # This is the dedicated per-placeholder copy produced by runPptCopywriting()
        # in the browser before /generate-pptx is called. It supersedes the
        # heuristic truncate-and-paste values from map_replacements() because it
        # is the only path that gives each card / panel its own box-budgeted
        # text. We apply it AFTER map_replacements (so atomic chips like cmp /
        # target / KPI values still come from metadata) and BEFORE cs_ppt_data
        # so explicit user UI overrides remain the highest-priority source.
        ppt_copy = session.get("ppt_content_json")
        if isinstance(ppt_copy, str):
            try:
                ppt_copy = json.loads(ppt_copy)
            except Exception:
                ppt_copy = None
        if isinstance(ppt_copy, dict) and ppt_copy:
            applied = 0
            for k, v in ppt_copy.items():
                if v is None:
                    continue
                value = str(v).strip()
                if not value:
                    continue
                replacements[k] = value
                applied += 1
            _sync_equivalent_keys(replacements)
            logger.info("Applied %d slide-copy values from ppt_content_json", applied)
            warnings.append(f"Slide copy fields applied: {applied}")
        else:
            logger.info("No ppt_content_json on session; using heuristic copy only")
            warnings.append("Slide copy fields applied: 0 (run 'Generate slide copy' in the UI for better text)")

        # Apply saved PPT placeholder overrides confirmed by the user in the UI
        saved_ppt_raw = report.get("cs_ppt_data") or ""
        if saved_ppt_raw:
            try:
                saved_overrides = json.loads(saved_ppt_raw)
                if isinstance(saved_overrides, dict):
                    replacements.update(saved_overrides)
                    _sync_equivalent_keys(replacements)
                    # Re-normalise bullet placeholders after the override merge
                    # — overrides from the UI / LLM may bypass map_replacements'
                    # post-processing.
                    for _bk in _BULLET_LIST_KEYS:
                        if _bk in replacements and isinstance(replacements[_bk], str):
                            replacements[_bk] = _ensure_bullet_lines(replacements[_bk])
                    logger.info("Applied %d saved PPT overrides from cs_ppt_data", len(saved_overrides))
            except Exception as e:
                logger.warning("Failed to parse cs_ppt_data overrides: %s", e)

        # Check multiple possible paths for the master template (local vs docker)
        possible_paths = [
            Path(__file__).resolve().parent.parent.parent / "master_template.pptx",  # local repo structure
            Path(__file__).resolve().parent / "master_template.pptx",  # If dumped in the same directory
            Path("/app/master_template.pptx"), # Docker container path
            Path("master_template.pptx"), # Current working directory
        ]
        
        template_path = None
        for p in possible_paths:
            if p.exists():
                template_path = p
                break
                
        if not template_path:
            raise RuntimeError(f"Master template not found in any of the expected locations.")

        result_pptx_path = str(output_root / "report.pptx")
        excel_path = _download_model_excel(client, ticker, warnings, tmp_root)
        
        summary_image = None
        if excel_path:
            try:
                summary_image = _render_financial_summary_dashboard(str(excel_path), company.get("name", ""))
            except Exception as exc:
                logger.warning("Financial summary dashboard render failed before template fill: %s", exc)

        if template_path.exists():
            fill_master_template(
                str(template_path), result_pptx_path, replacements, fin_model,
                company_name=company.get("name", ""),
                summary_image=summary_image,
            )
        else:
            raise RuntimeError(f"Master template not found at {template_path}")

        story_injections = inject_story_chart_slides(
            result_pptx_path,
            fin_model,
            company.get("name", ""),
            excel_path=str(excel_path) if excel_path else None,
            cmp_value=_parse_number(metadata.get("cmp")),
        )
        if story_injections:
            logger.info("Injected %d story-chart slide visuals", story_injections)

        overview_injections = inject_company_overview_slide(result_pptx_path, fin_model)
        if overview_injections:
            logger.info("Injected %d company-overview pie visuals", overview_injections)

        timeline_injections = inject_company_timeline_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
        )
        if timeline_injections:
            logger.info("Injected %d company timeline visuals", timeline_injections)

        catalyst_timeline_injections = inject_catalyst_timeline_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
        )
        if catalyst_timeline_injections:
            logger.info("Injected %d catalyst timeline visuals", catalyst_timeline_injections)

        competitive_injections = inject_competitive_advantage_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
        )
        if competitive_injections:
            logger.info("Injected %d competitive-advantage visuals", competitive_injections)

        peer_table_injections = inject_peer_comparison_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
        )
        if peer_table_injections:
            logger.info("Injected %d peer-comparison visuals", peer_table_injections)

        governance_table_injections = inject_governance_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
        )
        if governance_table_injections:
            logger.info("Injected %d governance table visuals", governance_table_injections)

        earnings_table_injections = inject_formula_table_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
            slide_number=13,
            sheet_name="Earnings_Forecast",
            token="{{earnings_forecast_table}}",
            fallback_text="Earnings forecast — see Excel model for details.",
            max_row=34,
            max_col=11,
        )
        if earnings_table_injections:
            logger.info("Injected %d earnings forecast visuals", earnings_table_injections)

        financials_table_injections = inject_formula_table_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
            slide_number=14,
            sheet_name="Financials_Table",
            token="{{financials_table}}",
            fallback_text="Financials — see Excel model for details.",
            max_row=45,
            max_col=12,
        )
        if financials_table_injections:
            logger.info("Injected %d financials table visuals", financials_table_injections)

        valuations_table_injections = inject_formula_table_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
            slide_number=15,
            sheet_name="Valuations_Table",
            token="{{valuations_table}}",
            fallback_text="Valuations — see Excel model for details.",
            max_row=45,
            max_col=11,
        )
        if valuations_table_injections:
            logger.info("Injected %d valuations table visuals", valuations_table_injections)

        key_risks_table_injections = inject_key_risks_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
        )
        if key_risks_table_injections:
            logger.info("Injected %d key-risks visuals", key_risks_table_injections)

        probability_weight_injections = inject_probability_weight_slide(
            result_pptx_path,
            fin_model,
        )
        if probability_weight_injections:
            logger.info("Injected %d probability-weight visuals", probability_weight_injections)

        # Inject Excel visuals: COM on Windows, openpyxl+matplotlib everywhere else
        logger.info("Attempting to inject Excel tables/charts into PPTX...")
        injection_count = 0
        summary_injection_count = 0
        if excel_path and excel_path.exists():
            # inject_excel_visuals_into_ppt tries COM first, falls back to openpyxl renderer
            injection_count = excel_injector.inject_excel_visuals_into_ppt(
                str(excel_path), result_pptx_path
            )
            logger.info("Excel injection completed: %d visuals injected", injection_count)

            # Re-apply the custom slide-1 dashboard after generic Excel
            # injection. This keeps the summary panel deterministic even when
            # placeholder preservation/order differs across code paths.
            summary_injection_count = inject_financial_summary_slide(
                result_pptx_path,
                company.get("name", ""),
                excel_path=str(excel_path),
            )
            if summary_injection_count:
                logger.info(
                    "Financial summary dashboard injected on slide 1 (%d insertion)",
                    summary_injection_count,
                )

            # If main injector returned 0, try direct openpyxl rendering as second chance
            if injection_count == 0:
                logger.info("Trying direct openpyxl sheet rendering as fallback...")
                images = excel_injector.render_all_excel_sheets(str(excel_path))
                if images:
                    injection_count = excel_injector.inject_excel_visuals_into_pptx(
                        result_pptx_path, images
                    )
                    logger.info("Openpyxl fallback injected %d visuals", injection_count)

                    summary_injection_count = inject_financial_summary_slide(
                        result_pptx_path,
                        company.get("name", ""),
                        excel_path=str(excel_path),
                    )
                    if summary_injection_count:
                        logger.info(
                            "Financial summary dashboard re-injected after openpyxl fallback (%d insertion)",
                            summary_injection_count,
                        )
        else:
            logger.warning("No Excel file available — Excel injection skipped entirely")
            warnings.append("Excel model file not found; financial tables/charts use text fallback")

        if injection_count > 0:
            warnings.append(f"Injected {injection_count} Excel tables/charts into the report")
        warnings.append(f"Financial summary slide injection count: {summary_injection_count}")
        warnings.append(f"Company overview pie injection count: {overview_injections}")
        warnings.append(f"Company timeline injection count: {timeline_injections}")
        warnings.append(f"Catalyst timeline injection count: {catalyst_timeline_injections}")
        warnings.append(f"Competitive advantage injection count: {competitive_injections}")
        warnings.append(f"Peer comparison injection count: {peer_table_injections}")
        warnings.append(f"Governance table injection count: {governance_table_injections}")
        warnings.append(f"Earnings forecast injection count: {earnings_table_injections}")
        warnings.append(f"Financials table injection count: {financials_table_injections}")
        warnings.append(f"Valuations table injection count: {valuations_table_injections}")
        warnings.append(f"Key risks table injection count: {key_risks_table_injections}")
        warnings.append(f"Probability weight injection count: {probability_weight_injections}")

        # Cleanup pass: replace any surviving Excel injection placeholder tokens
        # with text-based fallback content so they don't appear as raw {{...}} text.
        _cleanup_excel_placeholders(result_pptx_path, replacements)
        overview_injections = inject_company_overview_slide(result_pptx_path, fin_model)
        if overview_injections:
            logger.info("Re-injected %d company-overview pie visuals after cleanup", overview_injections)
        competitive_injections = inject_competitive_advantage_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
        )
        if competitive_injections:
            logger.info("Re-injected %d competitive-advantage visuals after cleanup", competitive_injections)
        peer_table_injections = inject_peer_comparison_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
        )
        if peer_table_injections:
            logger.info("Re-injected %d peer-comparison visuals after cleanup", peer_table_injections)
        timeline_injections = inject_company_timeline_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
        )
        if timeline_injections:
            logger.info("Re-injected %d company timeline visuals after cleanup", timeline_injections)
        catalyst_timeline_injections = inject_catalyst_timeline_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
        )
        if catalyst_timeline_injections:
            logger.info("Re-injected %d catalyst timeline visuals after cleanup", catalyst_timeline_injections)
        governance_table_injections = inject_governance_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
        )
        if governance_table_injections:
            logger.info("Re-injected %d governance table visuals after cleanup", governance_table_injections)
        earnings_table_injections = inject_formula_table_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
            slide_number=13,
            sheet_name="Earnings_Forecast",
            token="{{earnings_forecast_table}}",
            fallback_text="Earnings forecast — see Excel model for details.",
            max_row=34,
            max_col=11,
        )
        if earnings_table_injections:
            logger.info("Re-injected %d earnings forecast visuals after cleanup", earnings_table_injections)
        financials_table_injections = inject_formula_table_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
            slide_number=14,
            sheet_name="Financials_Table",
            token="{{financials_table}}",
            fallback_text="Financials — see Excel model for details.",
            max_row=45,
            max_col=12,
        )
        if financials_table_injections:
            logger.info("Re-injected %d financials table visuals after cleanup", financials_table_injections)
        valuations_table_injections = inject_formula_table_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
            slide_number=15,
            sheet_name="Valuations_Table",
            token="{{valuations_table}}",
            fallback_text="Valuations — see Excel model for details.",
            max_row=45,
            max_col=11,
        )
        if valuations_table_injections:
            logger.info("Re-injected %d valuations table visuals after cleanup", valuations_table_injections)
        key_risks_table_injections = inject_key_risks_slide(
            result_pptx_path,
            excel_path=str(excel_path) if excel_path else None,
        )
        if key_risks_table_injections:
            logger.info("Re-injected %d key-risks visuals after cleanup", key_risks_table_injections)
        probability_weight_injections = inject_probability_weight_slide(
            result_pptx_path,
            fin_model,
        )
        if probability_weight_injections:
            logger.info("Re-injected %d probability-weight visuals after cleanup", probability_weight_injections)

        # Upload artifacts
        ts = int(time.time())
        pptx_key = f"{ticker}/{report_id}/report_{ts}.pptx"
        pptx_path_out, pptx_url = _upload(client, Path(result_pptx_path), pptx_key, PPTX_CONTENT_TYPE)

        # Upload copy to Google Drive if vault_folder_id is set
        vault_folder_id = session.get("vault_folder_id")
        drive_file_id = None
        drive_file_url = None
        if vault_folder_id:
            gdrive_res = _upload_to_gdrive(Path(result_pptx_path), f"{ticker}_Report_{ts}.pptx", vault_folder_id, "Research Report")
            if gdrive_res:
                drive_file_id = gdrive_res.get("id")
                drive_file_url = gdrive_res.get("url")
                logger.info("Successfully uploaded PPT to Drive: id=%s url=%s", drive_file_id, drive_file_url)
                warnings.append("PowerPoint uploaded to Google Drive vault (Google Slides view ready)")
            else:
                warnings.append("Failed to upload PowerPoint copy to Google Drive vault")

        pdf_path_out: str | None = None
        pdf_url: str | None = None
        
        pdf_path = _convert_pptx_to_pdf(Path(result_pptx_path), output_root)
        if pdf_path and pdf_path.exists():
            pdf_key = f"{ticker}/{report_id}/report_{ts}.pdf"
            pdf_path_out, pdf_url = _upload(client, pdf_path, pdf_key, PDF_CONTENT_TYPE)
        else:
            warnings.append("PDF conversion skipped (LibreOffice soffice failed or not found)")

    # Update DB
    now_iso = datetime.now(timezone.utc).isoformat()
    db_payload = {
        "pptx_file_path": pptx_path_out,
        "pptx_file_url": pptx_url,
        "pptx_pdf_file_path": pdf_path_out,
        "pptx_pdf_file_url": pdf_url,
        "pptx_generated_at": now_iso,
        "pptx_status": "ready",
        "updated_at": now_iso,
    }
    if drive_file_id:
        db_payload["drive_file_id"] = drive_file_id
        db_payload["ppt_file_id"] = drive_file_id
    if drive_file_url:
        db_payload["drive_file_url"] = drive_file_url
        db_payload["ppt_file_url"] = drive_file_url

    client.table("research_reports").update(db_payload).eq("report_id", report_id).execute()

    duration = round(time.time() - t0, 2)
    logger.info("generate_pptx done report=%s duration=%.2fs", report_id, duration)

    return {
        "status": "success",
        "message": f"PPTX generated in {duration}s",
        "pptx_file_url": pptx_url,
        "pptx_file_path": pptx_path_out,
        "pptx_pdf_file_url": pdf_url,
        "pptx_pdf_file_path": pdf_path_out,
        "ppt_file_id": drive_file_id,
        "ppt_file_url": drive_file_url,
        "duration_seconds": duration,
        "warnings": warnings,
    }
